"""Build the C4A Corporate Business Development deck.

This is not a company profile. It is a capability and business development deck:
what we do, how we do it, how you buy it, and what it produces.

Audience: mixed institutional. Hospital owners and boards, hospital groups,
payers and HMOs, investors and DFIs, government and development partners,
and healthcare ventures.

Content sources:
  - consultforafrica.com (services, solutions, about)
  - Healthcare Marketing Agency positioning document (agency team lead)

Outputs:
  docs/cfa-corporate-bd-deck.pptx
  docs/cfa-corporate-bd-deck.pdf

One content definition drives both. The PPTX is the editable master; the PDF is
drawn directly with reportlab so it does not depend on LibreOffice being
installed. PPTX uses Calibri, PDF uses Helvetica, and the PDF backend shrinks
text a step at a time if the wider metrics would overflow a fixed box.

Run:
  python3 scripts/build-cfa-bd-deck.py
"""
from __future__ import annotations

from pathlib import Path

# ── Geometry (inches) ──────────────────────────────────────────────────
SLIDE_W = 13.333
SLIDE_H = 7.5

ROOT = Path(__file__).resolve().parents[1]
DOCS = ROOT / "docs"
# The flat c4a-logo.png / c4a-icon.png have no alpha channel and paint a white
# box on dark artwork. Every placement in this deck sits on a dark ground (the
# navy/ink header band, or a full-bleed dark slide), so we use the transparent
# reversed marks throughout.
LOGO_FULL = str(DOCS / "c4a-logo-reversed.png")
LOGO_ICON = str(DOCS / "c4a-icon-reversed.png")

# ── Brand palette (hex) ────────────────────────────────────────────────
NAVY = "#0B3C5D"
DEEP = "#0F2744"
INK = "#06090F"
PANEL_DK = "#101A2B"
AMBER_DK = "#1A1508"
GOLD = "#D4AF37"
LIGHT = "#F4F7FA"
PANEL = "#E9EFF5"
DARK = "#111827"
WHITE = "#FFFFFF"
GREY = "#6B7280"
MIDGREY = "#9AA3AF"
RULE = "#DDE4EC"
PALE = "#C8D4E2"


# ══════════════════════════════════════════════════════════════════════
# Backends
# ══════════════════════════════════════════════════════════════════════
def _run(t, size=12, bold=False, color=DARK, italic=False, sb=0, ls=1.2):
    """One paragraph inside a text box."""
    return dict(t=t, size=size, bold=bold, color=color, italic=italic, sb=sb, ls=ls)


class PptxBackend:
    ext = "pptx"

    def __init__(self):
        from pptx import Presentation
        from pptx.util import Inches

        self.prs = Presentation()
        self.prs.slide_width = Inches(SLIDE_W)
        self.prs.slide_height = Inches(SLIDE_H)
        self.blank = self.prs.slide_layouts[6]
        self.s = None
        self.font = "Calibri"

    # -- helpers ------------------------------------------------------
    @staticmethod
    def _rgb(hexstr):
        from pptx.dml.color import RGBColor

        h = hexstr.lstrip("#")
        return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))

    @staticmethod
    def _in(v):
        from pptx.util import Inches

        return Inches(v)

    # -- api ----------------------------------------------------------
    def slide(self):
        self.s = self.prs.slides.add_slide(self.blank)

    def rect(self, x, y, w, h, fill, line=None):
        from pptx.enum.shapes import MSO_SHAPE
        from pptx.util import Pt

        shp = self.s.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, self._in(x), self._in(y), self._in(w), self._in(h)
        )
        shp.fill.solid()
        shp.fill.fore_color.rgb = self._rgb(fill)
        if line is None:
            shp.line.fill.background()
        else:
            shp.line.color.rgb = self._rgb(line)
            shp.line.width = Pt(1)
        shp.shadow.inherit = False

    def text(self, x, y, w, h, runs, align="l", anchor="t"):
        from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
        from pptx.util import Pt

        amap = {"l": PP_ALIGN.LEFT, "c": PP_ALIGN.CENTER, "r": PP_ALIGN.RIGHT}
        vmap = {"t": MSO_ANCHOR.TOP, "m": MSO_ANCHOR.MIDDLE, "b": MSO_ANCHOR.BOTTOM}

        box = self.s.shapes.add_textbox(self._in(x), self._in(y), self._in(w), self._in(h))
        tf = box.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = vmap[anchor]
        tf.margin_left = tf.margin_right = self._in(0.04)
        tf.margin_top = tf.margin_bottom = self._in(0.02)
        tf.clear()

        for i, r in enumerate(runs):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.alignment = amap[align]
            p.line_spacing = r["ls"]
            if r["sb"]:
                p.space_before = Pt(r["sb"])
            run = p.add_run()
            run.text = r["t"]
            run.font.size = Pt(r["size"])
            run.font.bold = r["bold"]
            run.font.italic = r["italic"]
            run.font.color.rgb = self._rgb(r["color"])
            run.font.name = self.font

    def image(self, path, x, y, height):
        self.s.shapes.add_picture(path, self._in(x), self._in(y), height=self._in(height))

    def save(self, path):
        self.prs.save(path)


class PdfBackend:
    ext = "pdf"

    def __init__(self):
        from reportlab.pdfgen import canvas
        from reportlab.lib.utils import ImageReader

        self._ImageReader = ImageReader
        self._buf = []
        self.path_tmp = None
        self.c = None
        self._canvas_cls = canvas.Canvas
        self._started = False
        self._img_cache = {}

    def _ensure(self, path):
        if self.c is None:
            self.c = self._canvas_cls(str(path), pagesize=(SLIDE_W * 72, SLIDE_H * 72))
            self.c.setTitle("Consult For Africa | Corporate Capability & Business Development")
            self.c.setAuthor("Consult For Africa")

    # -- helpers ------------------------------------------------------
    @staticmethod
    def _fontname(bold, italic):
        if bold and italic:
            return "Helvetica-BoldOblique"
        if bold:
            return "Helvetica-Bold"
        if italic:
            return "Helvetica-Oblique"
        return "Helvetica"

    @staticmethod
    def _rgb(hexstr):
        from reportlab.lib.colors import HexColor

        return HexColor(hexstr)

    def _wrap(self, text, font, size, maxw):
        from reportlab.pdfbase.pdfmetrics import stringWidth

        words = text.split()
        if not words:
            return [""]
        lines, cur = [], words[0]
        for w in words[1:]:
            trial = cur + " " + w
            if stringWidth(trial, font, size) <= maxw:
                cur = trial
            else:
                lines.append(cur)
                cur = w
        lines.append(cur)
        return lines

    def _layout(self, runs, maxw, scale):
        """Return (lines, total_height_pt) at a given font scale."""
        out, total = [], 0.0
        for r in runs:
            size = r["size"] * scale
            font = self._fontname(r["bold"], r["italic"])
            lines = self._wrap(r["t"], font, size, maxw)
            lh = size * r["ls"]
            total += r["sb"] * scale + lh * len(lines)
            out.append((r, size, font, lines, lh))
        return out, total

    # -- api ----------------------------------------------------------
    def slide(self):
        if self._started:
            self.c.showPage()
        self._started = True

    def rect(self, x, y, w, h, fill, line=None):
        c = self.c
        c.setFillColor(self._rgb(fill))
        if line is None:
            c.setStrokeColor(self._rgb(fill))
            c.rect(x * 72, (SLIDE_H - y - h) * 72, w * 72, h * 72, stroke=0, fill=1)
        else:
            c.setStrokeColor(self._rgb(line))
            c.setLineWidth(0.75)
            c.rect(x * 72, (SLIDE_H - y - h) * 72, w * 72, h * 72, stroke=1, fill=1)

    def text(self, x, y, w, h, runs, align="l", anchor="t"):
        c = self.c
        pad = 0.04 * 72
        maxw = w * 72 - pad * 2
        box_h = h * 72

        # Shrink-to-fit: Helvetica is wider than Calibri, so a box tuned for
        # PPTX can overflow here. Step down until it fits (floor 82%).
        scale = 1.0
        laid, total = self._layout(runs, maxw, scale)
        while total > box_h and scale > 0.82:
            scale -= 0.04
            laid, total = self._layout(runs, maxw, scale)

        if anchor == "m":
            cursor = (SLIDE_H - y) * 72 - (box_h - total) / 2.0
        elif anchor == "b":
            cursor = (SLIDE_H - y - h) * 72 + total
        else:
            cursor = (SLIDE_H - y) * 72 - 0.02 * 72

        for r, size, font, lines, lh in laid:
            cursor -= r["sb"] * scale
            c.setFont(font, size)
            c.setFillColor(self._rgb(r["color"]))
            for ln in lines:
                cursor -= lh
                ty = cursor + lh * 0.22
                if align == "c":
                    c.drawCentredString(x * 72 + w * 72 / 2.0, ty, ln)
                elif align == "r":
                    c.drawRightString(x * 72 + w * 72 - pad, ty, ln)
                else:
                    c.drawString(x * 72 + pad, ty, ln)

    def image(self, path, x, y, height):
        if path not in self._img_cache:
            self._img_cache[path] = self._ImageReader(path)
        img = self._img_cache[path]
        iw, ih = img.getSize()
        width = height * (iw / ih)
        self.c.drawImage(
            img, x * 72, (SLIDE_H - y - height) * 72, width * 72, height * 72,
            mask="auto",
        )

    def save(self, path):
        self._ensure(path)
        self.c.save()

    def begin(self, path):
        self._ensure(path)


# ══════════════════════════════════════════════════════════════════════
# Drawing helpers (backend agnostic)
# ══════════════════════════════════════════════════════════════════════
class Deck:
    def __init__(self, be, total):
        self.be = be
        self.n = 0
        self.total = total

    # -- chrome -------------------------------------------------------
    def new(self, dark=False):
        self.be.slide()
        self.n += 1
        if dark:
            self.be.rect(0, 0, SLIDE_W, SLIDE_H, INK)

    def chrome(self, section, dark=False):
        be = self.be
        be.rect(0, 0, SLIDE_W, 0.52, INK if dark else NAVY)
        be.rect(0, 0.52, SLIDE_W, 0.035, GOLD)
        be.text(0.6, 0.11, 7.0, 0.32, [_run("CONSULT FOR AFRICA", 10, True, WHITE)])
        be.text(7.0, 0.11, 5.0, 0.32, [_run(section.upper(), 10, True, GOLD)], align="r")
        be.image(LOGO_ICON, 12.55, 0.05, 0.42)
        foot = MIDGREY if dark else GREY
        be.text(0.6, 7.08, 8.0, 0.28,
                [_run("consultforafrica.com   |   Lagos  ·  Abuja   |   Maarova™  ·  CadreHealth", 8.5, False, foot)])
        be.text(9.0, 7.08, 3.75, 0.28,
                [_run(f"{self.n:02d} / {self.total:02d}", 8.5, False, foot)], align="r")

    def title(self, title, kicker=None, sub=None, dark=False, y=0.95):
        be = self.be
        tcol = WHITE if dark else NAVY
        scol = MIDGREY if dark else GREY
        if kicker:
            be.text(0.6, y, 12.1, 0.3, [_run(kicker.upper(), 10.5, True, GOLD)])
            y += 0.34
        be.text(0.6, y, 12.1, 0.62, [_run(title, 27, True, tcol, ls=1.1)])
        y += 0.66
        be.rect(0.6, y, 0.7, 0.045, GOLD)
        y += 0.16
        if sub:
            be.text(0.6, y, 11.6, 0.55, [_run(sub, 13, False, scol, ls=1.25)])
            y += 0.62
        return y

    # -- components ---------------------------------------------------
    def stats(self, items, y, x=0.6, total_w=12.1, h=1.45, dark=False, vsize=26):
        n = len(items)
        gap = 0.18
        w = (total_w - gap * (n - 1)) / n
        for i, (value, label) in enumerate(items):
            cx = x + (w + gap) * i
            self.be.rect(cx, y, w, h, PANEL_DK if dark else LIGHT)
            self.be.rect(cx, y, w, 0.075, GOLD)
            self.be.text(cx + 0.22, y + 0.3, w - 0.44, 0.5,
                         [_run(value, vsize, True, GOLD if dark else NAVY)])
            self.be.text(cx + 0.22, y + 0.85, w - 0.44, 0.52,
                         [_run(label, 10, False, MIDGREY if dark else GREY, ls=1.15)])

    def lead_card(self, x, y, w, h, head, body, dark=False, hsize=12.0, bsize=11.0):
        self.be.rect(x, y, w, h, PANEL_DK if dark else WHITE, None if dark else RULE)
        self.be.rect(x, y, 0.045, h, GOLD)
        self.be.text(x + 0.28, y + 0.12, w - 0.5, h - 0.24, [
            _run(head, hsize, True, WHITE if dark else NAVY, ls=1.1),
            _run(body, bsize, False, MIDGREY if dark else GREY, sb=3, ls=1.15),
        ], anchor="m")


# ══════════════════════════════════════════════════════════════════════
# Content
# ══════════════════════════════════════════════════════════════════════
SERVICES = [
    dict(
        num="01", name="Hospital Turnaround & Financial Recovery",
        prop="When cashflow tightens and instability sets in, hospitals need decisive action. We restore financial control, stop revenue leakage, and rebuild operational discipline.",
        points=[
            ("Revenue capture", "Billing integrity, claims recovery, payer reconciliation, and price realisation."),
            ("Cost discipline", "Procurement control, stock and consumable leakage, staffing cost to revenue."),
            ("Productivity", "Theatre utilisation, bed days, clinic throughput, and diagnostics yield."),
            ("Cash and control", "Cashflow stabilisation, financial visibility, and a management reporting cadence."),
        ],
        mandate="Equity-backed turnaround or fixed-term management mandate, 24 to 60 months",
        best="Distressed or underperforming facilities, and lenders or investors holding the asset",
    ),
    dict(
        num="02", name="Strategy, Growth & Commercial Performance",
        prop="Growth comes from aligning clinical strengths with demand, referral flows, and patient access. Not from adding more services and hoping for the best.",
        points=[
            ("Service line strategy", "Where demand actually is, what it is worth, and what you are equipped to win."),
            ("Referral networks", "Physician engagement, partner clinics, and closing referral leakage."),
            ("Payer and pricing", "Payer mix, HMO dependency, tariff negotiation, and self-pay positioning."),
            ("Diversification", "New revenue lines, patient experience, and retention economics."),
        ],
        mandate="Advisory project, 8 to 16 weeks, with an optional execution phase to month 14",
        best="Groups expanding into new markets and facilities with underperforming service lines",
    ),
    dict(
        num="03", name="Clinical Governance & Accreditation",
        prop="Strong quality systems protect patients and build institutional credibility. We strengthen governance structures and prepare institutions for external accreditation.",
        points=[
            ("Governance frameworks", "Committee structure, escalation, clinical accountability, and board reporting."),
            ("Accreditation readiness", "JCI, COHSASA, and SafeCare preparation, gap closure, and mock survey."),
            ("Patient safety", "Incident systems, mortality and morbidity review, and risk registers."),
            ("Clinical audit", "Audit cycles, quality monitoring, and clinical KPI reporting."),
        ],
        mandate="Diagnostic audit then a 6 to 18 month readiness and implementation programme",
        best="Boards facing quality or safety exposure, and facilities pursuing accreditation",
    ),
    dict(
        num="04", name="Digital Health & Technology Leadership",
        prop="We help healthcare organisations and ventures get the technology foundation right, from system selection to digital strategy to the teams that actually run it.",
        points=[
            ("Operating intelligence", "Executive dashboards, KPI trees, and decision-grade reporting."),
            ("Systems", "HIS and EMR selection, implementation oversight, and workflow digitisation."),
            ("CTO as a service", "Technical leadership for healthtech ventures without a full-time hire."),
            ("Commercial flexibility", "Option to convert CTO fees into pre-seed equity."),
        ],
        mandate="Fractional or project mandate, 6 to 18 months, on-site or hybrid",
        best="Hospitals digitising under pressure and healthtech ventures pre-Series A",
    ),
    dict(
        num="05", name="Fractional Leadership & Executive Secondments",
        prop="Not every institution needs a full-time C-suite hire. We embed experienced healthcare executives on fixed-term mandates to fill leadership gaps and drive specific transformations.",
        points=[
            ("Fractional C-suite", "CEO, COO, CMO, and CTO on 2 to 3 days a week, on-site or hybrid."),
            ("Interim leadership", "Hospital Director and Medical Director cover through transitions."),
            ("Embedded operators", "Project-embedded clinical and operational leads inside your team."),
            ("Terms", "Performance-linked mandates with an optional conversion to permanent hire."),
        ],
        mandate="Secondment 3 to 12 months, fractional leadership 6 to 18 months",
        best="Institutions between leaders, in transition, or carrying a capability gap they cannot hire for",
    ),
    dict(
        num="06", name="Health Systems & Public Sector Advisory",
        prop="We support governments, development partners, and NGOs on health system design, hospital network planning, and moving policy into implementation.",
        points=[
            ("System design", "Service configuration, hospital network planning, and referral architecture."),
            ("Policy to delivery", "Turning approved policy into funded, sequenced, and monitored programmes."),
            ("Primary care", "PHC strengthening, facility readiness, and workforce distribution."),
            ("Partner advisory", "Programme design, delivery oversight, and results reporting for funders."),
        ],
        mandate="Programme mandate, typically 6 to 24 months, milestone-based",
        best="Ministries and state governments, DFIs, development partners, and NGOs",
    ),
    dict(
        num="07", name="Healthcare HR Management",
        prop="The workforce crisis is the single biggest threat to African healthcare. We combine consulting with proprietary assessment technology to help institutions hire, develop, and retain clinical leaders.",
        points=[
            ("Retention strategy", "Why your clinicians leave, what it costs, and what actually changes it."),
            ("Executive search", "Senior appointments supported by Maarova™ leadership assessment."),
            ("Leadership development", "Coaching, succession planning, and capability building for clinical leaders."),
            ("Workforce planning", "Establishment review, compensation design, and culture measurement."),
        ],
        mandate="Project or retained, powered by the Maarova™ assessment platform",
        best="Institutions losing senior clinicians and boards appointing into critical roles",
    ),
    dict(
        num="08", name="Healthcare Workforce Recruitment",
        prop="Access a verified pool of healthcare professionals across 16 cadres. Permanent hires, locum cover, and structured pre-employment assessment, run on our own platform.",
        points=[
            ("Permanent recruitment", "Clinical and non-clinical roles sourced from a verified talent pool."),
            ("Locum and surge", "Temporary staffing filled in days, not months."),
            ("Verification", "Credential checks and structured pre-employment assessment."),
            ("Market intelligence", "Salary benchmarking and workforce analytics by cadre and state."),
        ],
        mandate="Ongoing, per-placement or retained, delivered through CadreHealth",
        best="Any facility carrying vacancies, and groups standardising how they hire",
    ),
    dict(
        num="09", name="Healthcare Marketing Agency",
        prop="Empowering health brands through strategic marketing. A specialist agency inside the firm, blending data-driven strategy with medical expertise and creativity.",
        points=[
            ("Brand and positioning", "Brand architecture, value proposition, segmentation, identity, and guidelines."),
            ("Research and insight", "Epidemiological and competitive analysis, journey mapping, health equity audits."),
            ("Campaigns and creative", "Multichannel planning, medically vetted content, video, and influencer programmes."),
            ("Digital and analytics", "Web and app build, SEO, social, CRM, and live dashboards measuring ROI."),
        ],
        mandate="Retainer advisory, fixed-price project, or a hybrid of the two",
        best="Hospitals, HMOs, NGOs, and life sciences companies building brands and moving behaviour",
    ),
]

AGENCY_PHASES = [
    ("01", "Discovery & Insight",
     "Market research and stakeholder interviews with patients, clinicians, and regulators. Channel, competitor, and epidemiological audit."),
    ("02", "Strategy & Planning",
     "Positioning, messaging hierarchy, media strategy, and KPI framework, aligned to business and clinical objectives."),
    ("03", "Creative & Campaign",
     "Branding, storytelling, and digital content that simplify complexity without losing scientific integrity. Medically vetted."),
    ("04", "Execute & Optimise",
     "Multichannel deployment with compliance checks, live dashboards, behavioural journey design, and measured ROI."),
]

AGENCY_OFFERINGS = [
    "Brand strategy & positioning",
    "Market research & insights",
    "Campaign planning & management",
    "Creative services",
    "Digital marketing & technology",
    "Influencer & community engagement",
    "Public relations & advertising",
    "Analytics & CRM",
]


# ══════════════════════════════════════════════════════════════════════
# Slides
# ══════════════════════════════════════════════════════════════════════
def build_deck(d: Deck):
    be = d.be

    # ═══ 01 TITLE ═══════════════════════════════════════════════════
    d.new(dark=True)
    be.rect(0, 0, 0.16, SLIDE_H, GOLD)
    be.image(LOGO_FULL, 0.85, 0.75, 0.95)
    be.text(0.85, 2.55, 9.0, 0.35,
            [_run("CORPORATE CAPABILITY & BUSINESS DEVELOPMENT", 13, True, GOLD)])
    be.text(0.85, 3.05, 11.5, 1.9, [
        _run("We do not write reports", 41, True, WHITE, ls=1.05),
        _run("about African healthcare.", 41, True, WHITE, ls=1.05),
        _run("We run it.", 41, True, GOLD, ls=1.05),
    ])
    be.rect(0.85, 5.35, 1.1, 0.05, GOLD)
    be.text(0.85, 5.6, 10.5, 0.9, [
        _run("Nine service lines. Eight engagement models. Three proprietary platforms.", 15, False, WHITE, ls=1.3),
        _run("Healthcare transformation and management across Africa.", 15, False, MIDGREY, sb=3, ls=1.3),
    ])
    be.text(0.85, 6.85, 11.5, 0.35,
            [_run("consultforafrica.com   |   hello@consultforafrica.com   |   Lagos  ·  Abuja", 11, False, MIDGREY)])

    # ═══ 02 THE FIRM IN ONE PAGE ════════════════════════════════════
    d.new()
    d.chrome("The firm")
    y = d.title("The firm in one page", kicker="Orientation",
                sub="An Africa-focused healthcare management and transformation firm. Built by operators, for operators.")
    d.stats([
        ("135+ yrs", "Combined senior healthcare leadership across the partner network"),
        ("20+", "Senior operators in the C4A network, fielded as integrated teams"),
        ("16", "Professional cadres covered by our workforce platform"),
        ("3", "Proprietary platforms owned and operated by the firm"),
    ], y=y + 0.05)

    rows = [
        ("What we are", "A management and transformation firm, not an advisory house. We embed, implement, and stay long enough to hand over something that works."),
        ("Who we serve", "Private hospitals and groups, payers and HMOs, investors and DFIs, governments and development partners, NGOs, life sciences, and healthcare ventures."),
        ("Where we operate", "Headquartered in Lagos and Abuja, delivering across Nigeria and into Ghana, Kenya, and the wider continent."),
        ("How we are paid", "Fixed-scope projects, monthly retainers, secondment day rates, fractional mandates, and equity-linked turnaround structures."),
    ]
    ry = y + 1.62
    for i, (label, text) in enumerate(rows):
        top = ry + 0.68 * i
        be.rect(0.6, top, 12.1, 0.62, LIGHT if i % 2 == 0 else WHITE)
        be.text(0.85, top + 0.05, 2.5, 0.52, [_run(label, 12.5, True, NAVY)], anchor="m")
        be.text(3.5, top + 0.04, 9.0, 0.54, [_run(text, 11, False, DARK, ls=1.18)], anchor="m")

    # ═══ 03 THE PROBLEM WE EXIST FOR ════════════════════════════════
    d.new(dark=True)
    d.chrome("Why we exist", dark=True)
    d.title("Africa's health systems do not have a knowledge problem",
            kicker="The premise", dark=True)
    be.text(0.6, 2.35, 11.5, 0.6,
            [_run("They have an execution problem. That is the whole thesis of this firm.", 17, False, GOLD, ls=1.25)])

    gaps = [
        ("The strategy exists", "Boards have commissioned the reports. The recommendations are usually right. Almost none of them have been implemented."),
        ("The capability is thin", "Senior leaders were promoted on clinical tenure, not on management capability, and nobody has ever measured the difference."),
        ("The workforce is leaving", "Every vacancy costs months of locum premium and referral leakage while the post sits open."),
        ("Nobody owns the number", "Revenue leaks, theatre lists run half empty, and marketing spend disappears, because no single person is accountable for the outcome."),
    ]
    cw, gap = 2.9, 0.17
    for i, (head, body) in enumerate(gaps):
        cx = 0.6 + (cw + gap) * i
        be.rect(cx, 3.3, cw, 2.5, PANEL_DK)
        be.rect(cx, 3.3, cw, 0.07, GOLD)
        be.text(cx + 0.25, 3.55, cw - 0.5, 0.55, [_run(head, 13.5, True, WHITE, ls=1.1)])
        be.text(cx + 0.25, 4.2, cw - 0.5, 1.45, [_run(body, 10.5, False, MIDGREY, ls=1.25)])
    be.text(0.6, 6.15, 12.1, 0.7,
            [_run("Every service line in this deck exists to close one of those four gaps.", 13, False, WHITE)])

    # ═══ 04 WHAT MAKES US DIFFERENT ═════════════════════════════════
    d.new()
    d.chrome("Positioning")
    d.title("What separates us from a consulting firm", kicker="Positioning",
            sub="The difference is not the analysis. It is who is standing in the building when the plan has to work.")
    hy = 2.6
    be.rect(0.6, hy, 5.85, 0.45, PANEL)
    be.text(0.85, hy + 0.07, 5.3, 0.32, [_run("THE TRADITIONAL MODEL", 10.5, True, GREY)])
    be.rect(6.85, hy, 5.85, 0.45, NAVY)
    be.text(7.1, hy + 0.07, 5.3, 0.32, [_run("THE C4A MODEL", 10.5, True, GOLD)])

    compare = [
        ("Diagnoses, then leaves", "Diagnoses, then stays to deliver"),
        ("Partners sell, juniors deliver", "Former hospital CEOs and clinical leaders on the ground"),
        ("Frameworks imported from other markets", "Models built for African cost, payer, and workforce realities"),
        ("Success measured by the report", "Success measured by revenue, utilisation, and governance"),
        ("Dependency by design", "Capability transferred so performance holds after we leave"),
    ]
    for i, (old, new) in enumerate(compare):
        top = hy + 0.58 + 0.72 * i
        be.rect(0.6, top, 5.85, 0.62, LIGHT)
        be.text(0.85, top + 0.05, 5.3, 0.52, [_run(old, 11.5, False, GREY, ls=1.15)], anchor="m")
        be.rect(6.85, top, 5.85, 0.62, WHITE, RULE)
        be.rect(6.85, top, 0.045, 0.62, GOLD)
        be.text(7.1, top + 0.05, 5.3, 0.52, [_run(new, 11.5, True, NAVY, ls=1.15)], anchor="m")

    # ═══ 05 THE CAPABILITY MAP ══════════════════════════════════════
    d.new(dark=True)
    d.chrome("Capability map", dark=True)
    d.title("Nine service lines, one integrated team", kicker="What we do", dark=True,
            sub="Take one. Take several. Most institutions start with one and widen the mandate once it works.")
    cw, ch, gx, gy = 3.9, 1.24, 0.28, 0.16
    for i, svc in enumerate(SERVICES):
        col, row = i % 3, i // 3
        x = 0.6 + (cw + gx) * col
        y = 2.75 + (ch + gy) * row
        hot = svc["num"] == "09"
        be.rect(x, y, cw, ch, AMBER_DK if hot else PANEL_DK)
        be.rect(x, y, 0.05, ch, GOLD)
        be.text(x + 0.28, y + 0.14, cw - 0.5, 0.24, [_run(svc["num"], 10, True, GOLD)])
        be.text(x + 0.28, y + 0.4, cw - 0.5, 0.5, [_run(svc["name"], 12.5, True, WHITE, ls=1.08)])
        tag = "NEW PRACTICE" if hot else svc["points"][0][0]
        be.text(x + 0.28, y + 0.95, cw - 0.5, 0.3,
                [_run(tag.upper(), 8.5, True, GOLD if hot else MIDGREY)])

    # ═══ 06 to 14: SERVICE LINE DETAIL ══════════════════════════════
    for svc in SERVICES:
        d.new()
        d.chrome(f"Service line {svc['num']}")
        d.title(svc["name"], kicker=f"Service line {svc['num']}")
        be.text(0.6, 2.18, 12.1, 0.7, [_run(svc["prop"], 12.5, False, GREY, ls=1.25)])

        pw, ph = 5.85, 0.95
        for i, (head, body) in enumerate(svc["points"]):
            col, row = i % 2, i // 2
            d.lead_card(0.6 + (pw + 0.4) * col, 3.05 + (ph + 0.2) * row, pw, ph, head, body)

        by = 5.45
        be.rect(0.6, by, 12.1, 1.4, LIGHT)
        be.rect(0.6, by, 12.1, 0.05, NAVY)
        be.text(0.95, by + 0.22, 5.6, 0.25, [_run("TYPICAL MANDATE", 9.5, True, GOLD)])
        be.text(0.95, by + 0.52, 5.6, 0.72, [_run(svc["mandate"], 11.5, True, NAVY, ls=1.2)])
        be.rect(6.75, by + 0.25, 0.012, 0.9, RULE)
        be.text(7.1, by + 0.22, 5.4, 0.25, [_run("BEST FOR", 9.5, True, GOLD)])
        be.text(7.1, by + 0.52, 5.4, 0.72, [_run(svc["best"], 11.5, True, NAVY, ls=1.2)])

    # ═══ 15 INSIDE THE MARKETING AGENCY ═════════════════════════════
    d.new(dark=True)
    d.chrome("Marketing agency", dark=True)
    d.title("Inside the marketing agency", kicker="Spotlight: new practice", dark=True,
            sub="We are more than an agency. We are healthcare partners. Behavioural science, digital analytics, and human-centred design, applied to health brands.")

    pw = (12.1 - 0.18 * 3) / 4
    for i, (num, name, body) in enumerate(AGENCY_PHASES):
        x = 0.6 + (pw + 0.18) * i
        be.rect(x, 2.72, pw, 1.75, PANEL_DK)
        be.rect(x, 2.72, pw, 0.07, GOLD)
        be.text(x + 0.24, 2.95, pw - 0.48, 0.28, [_run(num, 13, True, GOLD)])
        be.text(x + 0.24, 3.26, pw - 0.48, 0.34, [_run(name, 12, True, WHITE, ls=1.05)])
        be.text(x + 0.24, 3.66, pw - 0.48, 0.72, [_run(body, 9.5, False, MIDGREY, ls=1.2)])

    be.text(0.6, 4.66, 12.1, 0.28, [_run("SERVICE OFFERINGS", 10.5, True, GOLD)])
    ow = (12.1 - 0.15 * 3) / 4
    for i, name in enumerate(AGENCY_OFFERINGS):
        col, row = i % 4, i // 4
        x = 0.6 + (ow + 0.15) * col
        y = 5.0 + (0.62 + 0.14) * row
        be.rect(x, y, ow, 0.62, PANEL_DK)
        be.rect(x, y, 0.04, 0.62, GOLD)
        be.text(x + 0.24, y + 0.06, ow - 0.42, 0.5, [_run(name, 10.5, True, WHITE, ls=1.1)], anchor="m")

    be.text(0.6, 6.52, 12.1, 0.45, [
        _run("Engagement: retainer advisory, fixed-price project, or a hybrid of the two. "
             "Serving hospitals, HMOs, NGOs, and life sciences companies.", 11, False, WHITE, ls=1.2)])

    # ═══ 16 PLATFORMS ═══════════════════════════════════════════════
    d.new(dark=True)
    d.chrome("Platforms", dark=True)
    d.title("We own the infrastructure, not just the advice", kicker="Proprietary platforms", dark=True,
            sub="Platforms built by the firm and used inside client engagements, or licensed directly.")
    plats = [
        ("Maarova™", "Leadership assessment & development",
         ["Psychometric assessment built for African healthcare leadership",
          "360 feedback, development reporting, and executive coaching",
          "Used in executive search, succession, and appointment decisions"]),
        ("CadreHealth", "Healthcare workforce platform",
         ["Verified professionals across 16 cadres and 73 hospitals",
          "Permanent and locum recruitment, credential verification",
          "Salary intelligence and workforce benchmarking data"]),
        ("Agent Channel", "Commercial distribution network",
         ["Commission-based sales agents recruited from our clinical network",
          "Territory management, deal tracking, and attribution",
          "Route to market for healthcare products and services"]),
    ]
    pw = 3.9
    for i, (name, tagline, bullets) in enumerate(plats):
        x = 0.6 + (pw + 0.28) * i
        be.rect(x, 2.9, pw, 3.35, PANEL_DK)
        be.rect(x, 2.9, pw, 0.075, GOLD)
        be.text(x + 0.3, 3.15, pw - 0.6, 0.4, [_run(name, 18, True, WHITE)])
        be.text(x + 0.3, 3.6, pw - 0.6, 0.35, [_run(tagline, 10.5, False, GOLD, ls=1.15)])
        be.text(x + 0.3, 4.15, pw - 0.6, 1.9,
                [_run(("—  " if j == 0 else "—  ") + b, 10.5, False, MIDGREY,
                      sb=(0 if j == 0 else 8), ls=1.25) for j, b in enumerate(bullets)])
    be.text(0.6, 6.45, 12.1, 0.5,
            [_run("Clients reach these platforms through an engagement, or licence them directly.", 12, False, WHITE)])

    # ═══ 17 ENGAGEMENT MODELS ═══════════════════════════════════════
    d.new()
    d.chrome("How to engage")
    d.title("Eight ways to work with us", kicker="Engagement models",
            sub="Every institution has different constraints. The commercial structure flexes, the accountability does not.")
    models = [
        ("Advisory project", "8 to 16 weeks", "Fixed scope, defined deliverables, embedded team"),
        ("Retainer advisory", "6 to 12 months rolling", "Monthly hours pool with priority specialist access"),
        ("Embedded secondment", "3 to 12 months", "Full-time consultant reporting into your leadership"),
        ("Fractional leadership", "6 to 18 months", "Part-time C-suite, 2 to 3 days a week"),
        ("Hospital transformation", "24 to 60 months", "Equity-backed turnaround with management control"),
        ("Transaction advisory", "6 to 18 months", "M&A and capital raising for healthcare assets"),
        ("Healthcare recruitment", "Ongoing", "Permanent, locum, and assessment via CadreHealth"),
        ("Commercial distribution", "Ongoing", "Managed commission agent network, Agent Channel"),
    ]
    mw, mh = 5.85, 1.0
    for i, (name, dur, desc) in enumerate(models):
        col, row = i % 2, i // 2
        x = 0.6 + (mw + 0.4) * col
        y = 2.7 + (mh + 0.16) * row
        be.rect(x, y, mw, mh, WHITE, RULE)
        be.rect(x, y, 0.05, mh, GOLD)
        be.text(x + 0.3, y + 0.14, 3.4, 0.3, [_run(name, 13, True, NAVY)])
        be.text(x + 0.3, y + 0.5, mw - 0.6, 0.4, [_run(desc, 10.5, False, GREY, ls=1.15)])
        be.text(x + 3.75, y + 0.14, 1.85, 0.3, [_run(dur, 10, True, GOLD)], align="r")

    # ═══ 18 THE EXECUTION MODEL ═════════════════════════════════════
    d.new(dark=True)
    d.chrome("How we work", dark=True)
    d.title("The C4A execution model", kicker="Delivery method", dark=True,
            sub="Five stages. The last one is the one most firms skip, and it is the reason performance holds.")
    steps = [
        ("01", "Diagnose", "Deep operational and financial assessment. We find the real problem, not the presenting one."),
        ("02", "Design", "A bespoke plan with clear milestones, ownership, and measurable targets."),
        ("03", "Deploy", "Embedded execution. Our team works inside the institution alongside your people."),
        ("04", "Deliver", "Measurable outcomes: revenue recovered, costs reduced, governance strengthened."),
        ("05", "Transfer", "Capability and systems handed to your team so performance sustains after we leave."),
    ]
    sw, sgap = 2.3, 0.19
    for i, (num, name, desc) in enumerate(steps):
        x = 0.6 + (sw + sgap) * i
        be.rect(x, 3.0, sw, 2.85, PANEL_DK)
        be.rect(x, 3.0, sw, 0.075, GOLD)
        be.text(x + 0.25, 3.28, sw - 0.5, 0.45, [_run(num, 22, True, GOLD)])
        be.text(x + 0.25, 3.82, sw - 0.5, 0.35, [_run(name, 15, True, WHITE)])
        be.text(x + 0.25, 4.25, sw - 0.5, 1.45, [_run(desc, 10, False, MIDGREY, ls=1.25)])
        if i < 4:
            be.rect(x + sw + 0.06, 4.35, 0.07, 0.07, GOLD)
    be.text(0.6, 6.15, 12.1, 0.6,
            [_run("We report on a fixed cadence throughout, to management monthly and to the board quarterly.", 12, False, WHITE)])

    # ═══ 19 WHAT IMPROVEMENT LOOKS LIKE ═════════════════════════════
    d.new()
    d.chrome("Impact")
    y = d.title("Where performance improves, value follows", kicker="What we move",
                sub="Four levers carry most of the value in a healthcare institution. We are measured on all four.")
    levers = [
        ("Revenue recovery", "Capture earned income, close billing gaps, and fix payer reconciliation"),
        ("Utilisation gains", "Theatre lists, bed days, clinic slots, and diagnostics running at capacity"),
        ("Cost discipline", "Procurement leakage, consumable waste, and staffing cost to revenue"),
        ("Governance strength", "Board oversight, clinical accountability, and decision-grade reporting"),
    ]
    lw = 2.9
    for i, (head, body) in enumerate(levers):
        x = 0.6 + (lw + 0.17) * i
        be.rect(x, y + 0.05, lw, 1.75, NAVY)
        be.rect(x, y + 0.05, lw, 0.07, GOLD)
        be.text(x + 0.25, y + 0.3, lw - 0.5, 0.4, [_run(head, 14, True, WHITE, ls=1.1)])
        be.text(x + 0.25, y + 0.8, lw - 0.5, 0.9, [_run(body, 10.5, False, PALE, ls=1.25)])

    ty = y + 2.05
    be.text(0.6, ty, 12.1, 0.3, [_run("TRACK RECORD", 10.5, True, GOLD)])
    d.stats([
        ("$1.1M+", "Annual savings delivered in a single engagement"),
        ("135+ yrs", "Combined senior leadership across the partner network"),
        ("20+", "Senior operators fielded as integrated teams"),
        ("16", "Professional cadres in the workforce platform"),
    ], y=ty + 0.35, h=1.35)
    be.text(0.6, 6.5, 12.1, 0.5, [
        _run("“Their embedded execution approach restored operational discipline, strengthened "
             "governance, and rebuilt leadership accountability.”   Hospital board representative",
             11, False, GREY, italic=True)])

    # ═══ 20 WHO WE SERVE ════════════════════════════════════════════
    d.new()
    d.chrome("Clients")
    d.title("Who we serve, and what they come to us for", kicker="Client segments")
    segs = [
        ("Private hospitals and groups", "Turnaround, growth, governance, workforce, and marketing"),
        ("Payers, HMOs and insurers", "Provider network strategy, member enrolment, engagement, and cost control"),
        ("Investors, PE and DFIs", "Commercial and operational due diligence, then post-deal value creation"),
        ("Governments and ministries", "Health system design, network planning, and policy to delivery"),
        ("NGOs and development partners", "Programme design, behaviour change campaigns, and results reporting"),
        ("Life sciences and pharma", "Market entry, brand launch, and medically vetted commercial campaigns"),
        ("Healthcare ventures", "CTO as a service, go to market, and distribution through Agent Channel"),
        ("Diaspora-owned institutions", "Market entry, build and operate mandates, and clinical governance"),
    ]
    sw, sh = 5.85, 1.02
    for i, (name, desc) in enumerate(segs):
        col, row = i % 2, i // 2
        x = 0.6 + (sw + 0.4) * col
        yy = 2.55 + (sh + 0.14) * row
        be.rect(x, yy, sw, sh, LIGHT if row % 2 == 0 else WHITE, RULE)
        be.rect(x + 0.28, yy + 0.44, 0.13, 0.13, GOLD)
        be.text(x + 0.6, yy + 0.16, sw - 0.9, 0.3, [_run(name, 12.5, True, NAVY)])
        be.text(x + 0.6, yy + 0.5, sw - 0.9, 0.4, [_run(desc, 10.5, False, GREY, ls=1.15)])

    # ═══ 21 THE BENCH ═══════════════════════════════════════════════
    d.new(dark=True)
    d.chrome("The bench", dark=True)
    d.title("Who actually shows up", kicker="Our people", dark=True,
            sub="Fielded as an integrated team, not as individual consultants billing hours.")
    bench = [
        ("Hospital CEOs and executive operators",
         "Founders and former CEOs of Nigerian and East African hospital groups, with real P&L accountability behind them."),
        ("Clinical governance authorities",
         "Clinicians who have led accreditation, quality systems, and patient safety programmes at scale."),
        ("Finance and revenue strategists",
         "Specialists in billing integrity, NHIS and HMO strategy, payer negotiation, and cost reduction."),
        ("Marketing, brand and digital specialists",
         "A multidisciplinary agency team spanning strategy, creative, technology, and analytics, vetted by clinicians."),
    ]
    for i, (head, body) in enumerate(bench):
        yy = 2.85 + 0.92 * i
        be.rect(0.6, yy, 12.1, 0.8, PANEL_DK)
        be.rect(0.6, yy, 0.05, 0.8, GOLD)
        be.text(0.95, yy + 0.14, 4.1, 0.55, [_run(head, 12.5, True, WHITE, ls=1.1)], anchor="m")
        be.text(5.3, yy + 0.12, 7.15, 0.6, [_run(body, 10.5, False, MIDGREY, ls=1.2)], anchor="m")
    be.rect(0.6, 6.5, 12.1, 0.5, AMBER_DK)
    be.text(0.95, 6.54, 11.5, 0.42, [
        _run("Led by Dr Debo Odulana, Founding Partner. Former CEO, Cedarcrest Hospitals Abuja. "
             "Former Chief Innovation & Strategy Officer, Evercare Hospital Lekki. Founder of Doctoora.",
             10, False, GOLD, ls=1.1)], anchor="m")

    # ═══ 22 HOW AN ENGAGEMENT STARTS ════════════════════════════════
    d.new()
    d.chrome("Next steps")
    d.title("How an engagement starts", kicker="Commercial pathway",
            sub="Low commitment at the front. The scope widens only once something has been proven.")
    path = [
        ("Step 01", "Discovery call", "30 minutes, confidential, no charge",
         "We establish the presenting problem, the constraints, and whether we are the right firm for it."),
        ("Step 02", "Diagnostic", "2 to 4 weeks, fixed fee",
         "A focused review of the operational, commercial, or brand position, delivered as a board-ready report with a prioritised fix list."),
        ("Step 03", "Scoped mandate", "Structure follows the finding",
         "Project, retainer, secondment, fractional leadership, or a full transformation mandate. Priced against the value at stake."),
        ("Step 04", "Delivery and transfer", "Monthly and quarterly reporting",
         "Embedded execution against agreed targets, then handover of systems and capability to your team."),
    ]
    for i, (step, name, terms, desc) in enumerate(path):
        yy = 2.7 + 1.06 * i
        be.rect(0.6, yy, 12.1, 0.94, LIGHT if i % 2 == 0 else WHITE, RULE)
        be.rect(0.6, yy, 0.05, 0.94, GOLD)
        be.text(0.95, yy + 0.14, 1.1, 0.28, [_run(step.upper(), 9.5, True, GOLD)])
        be.text(0.95, yy + 0.42, 3.0, 0.3, [_run(name, 13.5, True, NAVY)])
        be.text(0.95, yy + 0.7, 3.0, 0.24, [_run(terms, 9.5, False, GREY)])
        be.text(4.3, yy + 0.16, 8.2, 0.65, [_run(desc, 11, False, DARK, ls=1.2)], anchor="m")

    # ═══ 23 CLOSE ═══════════════════════════════════════════════════
    d.new(dark=True)
    be.rect(0, 0, 0.16, SLIDE_H, GOLD)
    be.image(LOGO_FULL, 0.85, 0.7, 0.85)
    be.text(0.85, 2.25, 11.5, 0.35, [_run("START THE CONVERSATION", 12.5, True, GOLD)])
    be.text(0.85, 2.7, 11.2, 1.5, [
        _run("Tell us what is not working.", 34, True, WHITE, ls=1.1),
        _run("We will tell you whether we can fix it.", 34, True, MIDGREY, ls=1.1),
    ])
    be.rect(0.85, 4.5, 1.1, 0.05, GOLD)
    contacts = [
        ("Web", "consultforafrica.com"),
        ("Email", "hello@consultforafrica.com"),
        ("Phone", "+234 913 813 8553"),
        ("Offices", "Lagos  ·  Abuja"),
    ]
    cw = 2.9
    for i, (label, value) in enumerate(contacts):
        x = 0.85 + (cw + 0.15) * i
        be.text(x, 4.95, cw, 0.28, [_run(label.upper(), 9.5, True, GOLD)])
        be.text(x, 5.25, cw, 0.4, [_run(value, 13, True, WHITE, ls=1.1)])
    be.text(0.85, 6.4, 11.5, 0.6, [
        _run("Consult For Africa   |   Healthcare transformation and management across Africa   |   "
             "Maarova™  ·  CadreHealth  ·  Agent Channel", 10.5, False, MIDGREY)])


TOTAL_SLIDES = 23


def main():
    outputs = []

    pptx_be = PptxBackend()
    build_deck(Deck(pptx_be, TOTAL_SLIDES))
    p = DOCS / "cfa-corporate-bd-deck.pptx"
    pptx_be.save(str(p))
    outputs.append(p)

    pdf_be = PdfBackend()
    p = DOCS / "cfa-corporate-bd-deck.pdf"
    pdf_be.begin(p)
    build_deck(Deck(pdf_be, TOTAL_SLIDES))
    pdf_be.save(str(p))
    outputs.append(p)

    for o in outputs:
        print(f"Saved {o}  ({o.stat().st_size // 1024} KB)")


if __name__ == "__main__":
    main()
