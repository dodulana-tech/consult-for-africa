"""Build the complete 17-slide CFA hospital pitch deck.

Single source of truth. Replaces all standalone files. Includes:
- C4A logo on title and CTA slides
- Icon mark in top-right of every content slide header
- Correct page numbering (X of 17)
- Corrected locum rates (Nigerian market reality)
- Foundation Stack price aligned across all references
- Order matches the user's curated merge

Audience: hospital leadership (CEO, MD, HR Director, Board).
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

# Brand palette
NAVY = RGBColor(0x0B, 0x3C, 0x5D)
GOLD = RGBColor(0xD4, 0xAF, 0x37)
TEAL = RGBColor(0x1F, 0x7A, 0x8C)
DEEP_TEAL = RGBColor(0x14, 0x55, 0x63)
LIGHT = RGBColor(0xF4, 0xF7, 0xFA)
DARK = RGBColor(0x11, 0x18, 0x27)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GREY = RGBColor(0x6B, 0x72, 0x80)
SOFT = RGBColor(0xE5, 0xE7, 0xEB)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
TOTAL_SLIDES = 17

LOGO_FULL = "/Users/debo/consult-for-africa/docs/c4a-logo.png"
LOGO_ICON = "/Users/debo/consult-for-africa/docs/c4a-icon.png"


def set_text(frame, text, size=18, bold=False, color=DARK, align=PP_ALIGN.LEFT, font="Calibri"):
    frame.clear()
    p = frame.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font


def add_paragraph(frame, text, size=14, bold=False, color=DARK, align=PP_ALIGN.LEFT, space_before=Pt(6), font="Calibri"):
    p = frame.add_paragraph()
    p.alignment = align
    p.space_before = space_before
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font


def add_rect(slide, x, y, w, h, fill=NAVY, line=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
    shape.shadow.inherit = False
    return shape


def add_textbox(slide, x, y, w, h):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.05)
    tf.margin_bottom = Inches(0.05)
    return tf


def header_band(slide, slide_num):
    """Standard top header band with logo icon and slide number."""
    add_rect(slide, 0, 0, SLIDE_W, Inches(0.6), fill=NAVY)
    tf = add_textbox(slide, Inches(0.6), Inches(0.1), Inches(8), Inches(0.4))
    set_text(tf, "CONSULT FOR AFRICA", size=11, bold=True, color=WHITE)
    # Slide number
    tf2 = add_textbox(slide, Inches(11.2), Inches(0.1), Inches(1.4), Inches(0.4))
    set_text(tf2, f"SLIDE {slide_num:02d}", size=11, bold=True, color=GOLD, align=PP_ALIGN.RIGHT)
    # Logo icon at far right of header
    slide.shapes.add_picture(LOGO_ICON, Inches(12.7), Inches(0.05), height=Inches(0.5))


def footer_band(slide, slide_num):
    """Standard bottom footer band with brand strip and page count."""
    add_rect(slide, 0, Inches(7.2), SLIDE_W, Inches(0.3), fill=LIGHT)
    tf = add_textbox(slide, Inches(0.6), Inches(7.22), Inches(8), Inches(0.25))
    set_text(tf, "consultforafrica.com  |  Maarova  |  CadreHealth", size=9, color=GREY)
    tf2 = add_textbox(slide, Inches(8.733), Inches(7.22), Inches(4.4), Inches(0.25))
    set_text(tf2, f"{slide_num} of {TOTAL_SLIDES}", size=9, color=GREY, align=PP_ALIGN.RIGHT)


def slide_title(slide, title, kicker=None):
    y = Inches(0.85)
    if kicker:
        tf = add_textbox(slide, Inches(0.6), y, Inches(12.1), Inches(0.4))
        set_text(tf, kicker, size=12, bold=True, color=GOLD)
        y = Inches(1.2)
    tf = add_textbox(slide, Inches(0.6), y, Inches(12.1), Inches(0.8))
    set_text(tf, title, size=30, bold=True, color=NAVY)
    add_rect(slide, Inches(0.6), Inches(2.0), Inches(0.8), Inches(0.05), fill=GOLD)


prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
blank = prs.slide_layouts[6]


# ============= SLIDE 1: TITLE =============
s = prs.slides.add_slide(blank)
add_rect(s, 0, 0, SLIDE_W, SLIDE_H, fill=NAVY)
# Logo at top-left (white background area would be ideal but on navy looks fine since logo has navy chevron)
s.shapes.add_picture(LOGO_FULL, Inches(0.6), Inches(0.6), height=Inches(0.9))
# Eyebrow
add_rect(s, Inches(0.6), Inches(2.7), Inches(1.5), Inches(0.08), fill=GOLD)
tf = add_textbox(s, Inches(0.6), Inches(2.2), Inches(8), Inches(0.4))
set_text(tf, "PROPOSAL FOR HOSPITAL LEADERSHIP", size=14, bold=True, color=GOLD)
# Title
tf = add_textbox(s, Inches(0.6), Inches(3.0), Inches(12), Inches(2.0))
set_text(tf, "Healthcare Leadership and Workforce", size=42, bold=True, color=WHITE)
add_paragraph(tf, "A Stack Built for African Hospitals", size=42, bold=True, color=WHITE, space_before=Pt(0))
# Subtitle
tf = add_textbox(s, Inches(0.6), Inches(5.2), Inches(12), Inches(0.6))
set_text(tf, "Maarova for leadership. CadreHealth for the workforce. One operator. One mission.",
         size=18, color=GOLD)
# Footer
tf = add_textbox(s, Inches(0.6), Inches(6.7), Inches(12), Inches(0.4))
set_text(tf, "consultforafrica.com", size=12, color=WHITE)


# ============= SLIDE 2: THE REALITY =============
s = prs.slides.add_slide(blank)
header_band(s, 2)
slide_title(s, "The reality your hospital is operating in", kicker="THE PROBLEM")

stats = [
    ("80%", "of Nigerian healthcare workers intend to emigrate", TEAL),
    ("8 to 12", "weeks lost recruiting each replacement", NAVY),
    ("N3M+", "monthly cost when a senior leader role sits vacant", GOLD),
    ("0", "objective benchmarks of hospital leadership capability", DARK),
]
card_w = Inches(2.85)
card_h = Inches(2.4)
gap = Inches(0.2)
start_x = Inches(0.6)
start_y = Inches(2.6)
for i, (stat, label, color) in enumerate(stats):
    x = start_x + (card_w + gap) * i
    add_rect(s, x, start_y, card_w, card_h, fill=LIGHT)
    add_rect(s, x, start_y, card_w, Inches(0.12), fill=color)
    tf = add_textbox(s, x + Inches(0.2), start_y + Inches(0.4), card_w - Inches(0.4), Inches(1.0))
    set_text(tf, stat, size=44, bold=True, color=color)
    tf = add_textbox(s, x + Inches(0.2), start_y + Inches(1.5), card_w - Inches(0.4), Inches(0.8))
    set_text(tf, label, size=12, color=DARK)

tf = add_textbox(s, Inches(0.6), Inches(5.4), Inches(12.1), Inches(1.3))
set_text(tf, "Two compounding gaps:", size=16, bold=True, color=NAVY)
add_paragraph(tf, "Talent leaves faster than you can replace, and the leaders left behind were promoted on tenure, not on capability.",
              size=15, color=DARK, space_before=Pt(8))
footer_band(s, 2)


# ============= SLIDE 3: COST OF INACTION =============
s = prs.slides.add_slide(blank)
header_band(s, 3)
slide_title(s, "What inaction is already costing you", kicker="COMMERCIAL REALITY")

rows = [
    ("One bad senior hire", "N18M to N40M", "Recruiting, onboarding, severance, lost output. Industry benchmark: 6 to 12 months of compensation."),
    ("One unfilled consultant role for 6 months", "N12M to N20M", "Locum coverage premium plus referral leakage to competing hospitals."),
    ("One leadership failure that triggers staff exodus", "N50M+", "Replacement recruiting, training ramp, patient volume loss, reputational drag."),
    ("One year of unmeasured leadership capability", "Unbounded", "You cannot fix what you have not assessed. Every quarter compounds."),
]
y = Inches(2.4)
for i, (label, cost, detail) in enumerate(rows):
    row_y = y + Inches(0.95) * i
    bg = LIGHT if i % 2 == 0 else WHITE
    add_rect(s, Inches(0.6), row_y, Inches(12.1), Inches(0.85), fill=bg)
    tf = add_textbox(s, Inches(0.85), row_y + Inches(0.1), Inches(4.5), Inches(0.3))
    set_text(tf, label, size=13, bold=True, color=NAVY)
    tf = add_textbox(s, Inches(0.85), row_y + Inches(0.42), Inches(4.5), Inches(0.4))
    set_text(tf, cost, size=18, bold=True, color=GOLD)
    tf = add_textbox(s, Inches(5.6), row_y + Inches(0.18), Inches(7), Inches(0.6))
    set_text(tf, detail, size=11, color=DARK)
footer_band(s, 3)


# ============= SLIDE 4: TWO PRODUCTS =============
s = prs.slides.add_slide(blank)
header_band(s, 4)
slide_title(s, "Two products. One operator. One mission.", kicker="THE STACK")

card_w = Inches(5.95)
card_h = Inches(4.4)
y = Inches(2.5)

# Maarova
add_rect(s, Inches(0.6), y, card_w, card_h, fill=NAVY)
add_rect(s, Inches(0.6), y, card_w, Inches(0.5), fill=GOLD)
tf = add_textbox(s, Inches(0.85), y + Inches(0.05), card_w - Inches(0.5), Inches(0.4))
set_text(tf, "MAAROVA", size=14, bold=True, color=NAVY)
tf = add_textbox(s, Inches(0.85), y + Inches(0.7), card_w - Inches(0.5), Inches(0.5))
set_text(tf, "Leadership Assessment and Coaching", size=20, bold=True, color=WHITE)
tf = add_textbox(s, Inches(0.85), y + Inches(1.4), card_w - Inches(0.5), Inches(2.8))
set_text(tf, "Fixes the top of the org chart.", size=13, color=GOLD, bold=True)
add_paragraph(tf, "Objective baseline of leadership capability across 6 dimensions",
              size=12, color=WHITE, space_before=Pt(10))
add_paragraph(tf, "360-degree feedback from peers, reports, and the next layer down",
              size=12, color=WHITE, space_before=Pt(6))
add_paragraph(tf, "1:1 and group coaching matched to identified gaps",
              size=12, color=WHITE, space_before=Pt(6))
add_paragraph(tf, "Executive retreats for leadership team alignment",
              size=12, color=WHITE, space_before=Pt(6))
add_paragraph(tf, "Built on Ubuntu coaching framework, not retrofitted Western tools",
              size=12, color=WHITE, space_before=Pt(6))

# CadreHealth
x2 = Inches(0.6) + card_w + Inches(0.2)
add_rect(s, x2, y, card_w, card_h, fill=TEAL)
add_rect(s, x2, y, card_w, Inches(0.5), fill=GOLD)
tf = add_textbox(s, x2 + Inches(0.25), y + Inches(0.05), card_w - Inches(0.5), Inches(0.4))
set_text(tf, "CADREHEALTH", size=14, bold=True, color=NAVY)
tf = add_textbox(s, x2 + Inches(0.25), y + Inches(0.7), card_w - Inches(0.5), Inches(0.5))
set_text(tf, "Healthcare Workforce Platform", size=20, bold=True, color=WHITE)
tf = add_textbox(s, x2 + Inches(0.25), y + Inches(1.4), card_w - Inches(0.5), Inches(2.8))
set_text(tf, "Fixes the pipeline of people moving through.", size=13, color=GOLD, bold=True)
add_paragraph(tf, "4,300+ verified healthcare professionals across 16 cadres",
              size=12, color=WHITE, space_before=Pt(10))
add_paragraph(tf, "Hospital reviews from real healthcare workers, anonymous and verified",
              size=12, color=WHITE, space_before=Pt(6))
add_paragraph(tf, "Salary intelligence by cadre, facility, and city",
              size=12, color=WHITE, space_before=Pt(6))
add_paragraph(tf, "Recruitment matching with career readiness scoring",
              size=12, color=WHITE, space_before=Pt(6))
add_paragraph(tf, "Locum marketplace and credential verification",
              size=12, color=WHITE, space_before=Pt(6))
footer_band(s, 4)


# ============= SLIDE 5: WHY BOTH =============
s = prs.slides.add_slide(blank)
header_band(s, 5)
slide_title(s, "Why both, together", kicker="THE INTEGRATED STORY")

col_w = Inches(3.9)
col_h = Inches(3.2)
y = Inches(2.5)
gap = Inches(0.25)

cols = [
    ("Strong Leadership", "Maarova", "Leaders who can build and retain teams. Diagnostic plus coaching plus retreats.", NAVY),
    ("Quality Workforce", "CadreHealth", "Verified talent pool, salary benchmarks, career readiness, hospital reviews.", TEAL),
    ("Better Outcomes", "Combined", "Lower attrition, faster fills, stronger governance, healthier patient outcomes.", GOLD),
]
start_x = Inches(0.6)
for i, (title, source, body, color) in enumerate(cols):
    x = start_x + (col_w + gap) * i
    add_rect(s, x, y, col_w, col_h, fill=LIGHT)
    add_rect(s, x, y, col_w, Inches(0.7), fill=color)
    tf = add_textbox(s, x + Inches(0.2), y + Inches(0.1), col_w - Inches(0.4), Inches(0.5))
    set_text(tf, title, size=16, bold=True, color=WHITE)
    tf = add_textbox(s, x + Inches(0.2), y + Inches(0.85), col_w - Inches(0.4), Inches(0.4))
    set_text(tf, source.upper(), size=10, bold=True, color=color)
    tf = add_textbox(s, x + Inches(0.2), y + Inches(1.3), col_w - Inches(0.4), Inches(1.8))
    set_text(tf, body, size=13, color=DARK)

tf = add_textbox(s, Inches(0.6), Inches(6.0), Inches(12.1), Inches(1.0))
set_text(tf, "Maarova fixes the top. CadreHealth fixes the pipeline.",
         size=18, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
add_paragraph(tf, "Leaders who know how to retain talent, plus a marketplace where they can find it.",
              size=14, color=DARK, align=PP_ALIGN.CENTER, space_before=Pt(8))
footer_band(s, 5)


# ============= SLIDE 6: MAAROVA HOW =============
s = prs.slides.add_slide(blank)
header_band(s, 6)
slide_title(s, "Maarova: how the engagement works", kicker="LEADERSHIP PRODUCT")

steps = [
    ("01", "Assess", "Online assessment plus 360 feedback. 90 minutes total commitment per leader."),
    ("02", "Debrief", "1:1 results review with accredited coach. Top 2 to 3 development goals identified."),
    ("03", "Coach", "12 sessions over 6 months, biweekly. 1:1 or group format. WhatsApp accountability."),
    ("04", "Sustain", "Mini 360 at month 6. Personal development plan. Optional cohort retreat."),
]
y = Inches(2.5)
step_w = Inches(2.95)
gap = Inches(0.15)
start_x = Inches(0.6)
for i, (num, title, body) in enumerate(steps):
    x = start_x + (step_w + gap) * i
    add_rect(s, x, y, step_w, Inches(3.6), fill=LIGHT)
    add_rect(s, x, y, step_w, Inches(0.9), fill=NAVY)
    tf = add_textbox(s, x + Inches(0.25), y + Inches(0.1), step_w - Inches(0.5), Inches(0.4))
    set_text(tf, num, size=14, bold=True, color=GOLD)
    tf = add_textbox(s, x + Inches(0.25), y + Inches(0.4), step_w - Inches(0.5), Inches(0.5))
    set_text(tf, title, size=20, bold=True, color=WHITE)
    tf = add_textbox(s, x + Inches(0.25), y + Inches(1.1), step_w - Inches(0.5), Inches(2.4))
    set_text(tf, body, size=12, color=DARK)
tf = add_textbox(s, Inches(0.6), Inches(6.4), Inches(12.1), Inches(0.5))
set_text(tf, "Outcome: a leader with a clear capability map, a coach, and a development plan tied to your hospital strategy.",
         size=13, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
footer_band(s, 6)


# ============= SLIDE 7: MAAROVA PRICING =============
s = prs.slides.add_slide(blank)
header_band(s, 7)
slide_title(s, "Maarova investment tiers", kicker="COMMERCIALS  /  MAAROVA")

tiers = [
    ("Assessment Only", "N400K to N600K", "Per leader. Full report and 1:1 debrief.", False),
    ("Assessment + Group Coaching", "N800K to N1.2M", "6 monthly group sessions over 6 months.", False),
    ("Assessment + 1:1 Coaching", "N2M to N3.5M", "12 sessions over 6 months. Recommended starting point for senior leaders.", True),
    ("Leadership Intensive", "N4M to N6M", "1:1 plus group plus on-site workshop. For C-suite.", False),
    ("Executive Retreat", "N1.5M to N2.5M", "Per person. 2 to 3 day offsite. Premium venue, peer cohort.", False),
    ("Enterprise Cohort", "N15M to N25M", "10+ leaders. Stack of assessment, coaching, and retreat. Best per-head value.", True),
]
y = Inches(2.4)
row_h = Inches(0.65)
for i, (name, price, body, featured) in enumerate(tiers):
    row_y = y + row_h * i
    bg = GOLD if featured else (LIGHT if i % 2 == 0 else WHITE)
    text_color = NAVY if featured else DARK
    add_rect(s, Inches(0.6), row_y, Inches(12.1), row_h, fill=bg)
    tf = add_textbox(s, Inches(0.85), row_y + Inches(0.1), Inches(4.0), Inches(0.45))
    set_text(tf, name, size=14, bold=True, color=text_color)
    tf = add_textbox(s, Inches(4.9), row_y + Inches(0.1), Inches(2.8), Inches(0.45))
    set_text(tf, price, size=15, bold=True, color=NAVY if featured else GOLD)
    tf = add_textbox(s, Inches(7.8), row_y + Inches(0.13), Inches(4.8), Inches(0.5))
    set_text(tf, body, size=11, color=text_color)
tf = add_textbox(s, Inches(0.6), Inches(6.5), Inches(12.1), Inches(0.5))
set_text(tf, "ANCHOR: LBS Advanced Management Programme is N17 to 18M for ONE leader. Our enterprise cohort sends 10+ leaders for the same outlay.",
         size=11, bold=True, color=NAVY)
footer_band(s, 7)


# ============= SLIDE 8: SECONDMENTS OVERVIEW =============
s = prs.slides.add_slide(blank)
header_band(s, 8)
slide_title(s, "When you need senior capability now, not in 6 months",
            kicker="THIRD MOTION  /  SECONDMENTS AND FRACTIONAL")

card_w = Inches(5.95)
card_h = Inches(4.3)
y = Inches(2.5)

add_rect(s, Inches(0.6), y, card_w, card_h, fill=NAVY)
add_rect(s, Inches(0.6), y, card_w, Inches(0.5), fill=GOLD)
tf = add_textbox(s, Inches(0.85), y + Inches(0.05), card_w - Inches(0.5), Inches(0.4))
set_text(tf, "SECONDMENT", size=14, bold=True, color=NAVY)
tf = add_textbox(s, Inches(0.85), y + Inches(0.7), card_w - Inches(0.5), Inches(0.5))
set_text(tf, "Full-time, embedded", size=20, bold=True, color=WHITE)
tf = add_textbox(s, Inches(0.85), y + Inches(1.35), card_w - Inches(0.5), Inches(2.8))
set_text(tf, "A CFA consultant deployed full-time into your hospital. They sit in your office, attend your stand-ups, own a P&L line.",
         size=12, color=WHITE)
add_paragraph(tf, "TYPICAL DURATION", size=10, bold=True, color=GOLD, space_before=Pt(12))
add_paragraph(tf, "3 to 12 months. Most engagements run 6 months with optional extension.",
              size=12, color=WHITE, space_before=Pt(4))
add_paragraph(tf, "BEST FOR", size=10, bold=True, color=GOLD, space_before=Pt(8))
add_paragraph(tf, "• Interim COO during turnaround", size=11, color=WHITE, space_before=Pt(4))
add_paragraph(tf, "• Embedded transformation lead", size=11, color=WHITE, space_before=Pt(2))
add_paragraph(tf, "• Acting CHRO during HR overhaul", size=11, color=WHITE, space_before=Pt(2))

x2 = Inches(0.6) + card_w + Inches(0.2)
add_rect(s, x2, y, card_w, card_h, fill=TEAL)
add_rect(s, x2, y, card_w, Inches(0.5), fill=GOLD)
tf = add_textbox(s, x2 + Inches(0.25), y + Inches(0.05), card_w - Inches(0.5), Inches(0.4))
set_text(tf, "FRACTIONAL", size=14, bold=True, color=NAVY)
tf = add_textbox(s, x2 + Inches(0.25), y + Inches(0.7), card_w - Inches(0.5), Inches(0.5))
set_text(tf, "Senior, part-time, ongoing", size=20, bold=True, color=WHITE)
tf = add_textbox(s, x2 + Inches(0.25), y + Inches(1.35), card_w - Inches(0.5), Inches(2.8))
set_text(tf, "A senior CFA leader works with you 1 to 3 days a week over a longer horizon. Strategic capability without the cost of a full-time hire.",
         size=12, color=WHITE)
add_paragraph(tf, "TYPICAL DURATION", size=10, bold=True, color=GOLD, space_before=Pt(12))
add_paragraph(tf, "6 to 12 months, often renewed. 1 to 3 days per week.",
              size=12, color=WHITE, space_before=Pt(4))
add_paragraph(tf, "BEST FOR", size=10, bold=True, color=GOLD, space_before=Pt(8))
add_paragraph(tf, "• Fractional CMO during clinical strategy reset", size=11, color=WHITE, space_before=Pt(4))
add_paragraph(tf, "• Fractional CFO for financial discipline", size=11, color=WHITE, space_before=Pt(2))
add_paragraph(tf, "• Fractional CHRO for talent transformation", size=11, color=WHITE, space_before=Pt(2))

tf = add_textbox(s, Inches(0.6), Inches(7.0), Inches(12.1), Inches(0.2))
set_text(tf, "Senior capability you can hire by the month, not the headcount line.",
         size=12, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
footer_band(s, 8)


# ============= SLIDE 9: SENIOR ROLES + BUNDLES =============
s = prs.slides.add_slide(blank)
header_band(s, 9)
slide_title(s, "Senior roles, monthly cost, what you get",
            kicker="COMMERCIALS  /  SECONDMENTS AND FRACTIONAL")

roles = [
    ("Fractional Chief Medical Officer", "N4M to N7M / month",
     "2 to 3 days a week. Clinical strategy, governance, quality oversight, medical leadership.", "FRACTIONAL"),
    ("Fractional CHRO", "N3M to N5M / month",
     "2 to 3 days a week. Workforce strategy, retention, succession planning, HR systems.", "FRACTIONAL"),
    ("Fractional CFO", "N3M to N6M / month",
     "2 days a week. Financial discipline, cash management, payor strategy, board reporting.", "FRACTIONAL"),
    ("Interim COO (Secondment)", "N6M to N10M / month",
     "Full-time, 6 to 12 months. Daily ops, turnaround, capacity planning, vendor renegotiation.", "SECONDMENT"),
    ("Embedded Transformation Lead", "N5M to N8M / month",
     "Full-time, 9 to 12 months. Owns delivery of a named transformation programme end to end.", "SECONDMENT"),
]
y = Inches(2.4)
row_h = Inches(0.6)
for i, (name, price, body, tag) in enumerate(roles):
    row_y = y + row_h * i
    bg = LIGHT if i % 2 == 0 else WHITE
    add_rect(s, Inches(0.6), row_y, Inches(12.1), row_h, fill=bg)
    tag_color = TEAL if tag == "FRACTIONAL" else NAVY
    add_rect(s, Inches(0.85), row_y + Inches(0.18), Inches(0.95), Inches(0.28), fill=tag_color)
    tf = add_textbox(s, Inches(0.85), row_y + Inches(0.18), Inches(0.95), Inches(0.28))
    set_text(tf, tag, size=8, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    tf = add_textbox(s, Inches(1.95), row_y + Inches(0.08), Inches(4.0), Inches(0.5))
    set_text(tf, name, size=13, bold=True, color=NAVY)
    tf = add_textbox(s, Inches(1.95), row_y + Inches(0.32), Inches(3.0), Inches(0.3))
    set_text(tf, price, size=12, bold=True, color=GOLD)
    tf = add_textbox(s, Inches(5.2), row_y + Inches(0.16), Inches(7.4), Inches(0.5))
    set_text(tf, body, size=10, color=DARK)

# Example bundles
y2 = Inches(5.55)
add_rect(s, Inches(0.6), y2, Inches(12.1), Inches(0.4), fill=NAVY)
tf = add_textbox(s, Inches(0.85), y2 + Inches(0.06), Inches(11.6), Inches(0.3))
set_text(tf, "EXAMPLE BUNDLES  /  WHAT YOU GET FOR WHAT", size=11, bold=True, color=GOLD)
bundle_y = y2 + Inches(0.45)
bundles = [
    ("N40M", "6-month Fractional CMO + Maarova for top 5 leaders"),
    ("N75M", "12-month Interim COO + CadreHealth talent search + 10 mid-manager hires"),
    ("N120M+", "Full transformation: 12-month Embedded Lead + Maarova cohort of 10 + CadreHealth"),
]
bw = Inches(3.95)
gap = Inches(0.2)
for i, (price, desc) in enumerate(bundles):
    bx = Inches(0.6) + (bw + gap) * i
    add_rect(s, bx, bundle_y, bw, Inches(0.7), fill=LIGHT)
    tf = add_textbox(s, bx + Inches(0.15), bundle_y + Inches(0.08), Inches(1.5), Inches(0.55))
    set_text(tf, price, size=18, bold=True, color=GOLD)
    tf = add_textbox(s, bx + Inches(1.65), bundle_y + Inches(0.05), bw - Inches(1.75), Inches(0.6))
    set_text(tf, desc, size=10, color=DARK)
footer_band(s, 9)


# ============= SLIDE 10: CADREHEALTH FOR HOSPITALS =============
s = prs.slides.add_slide(blank)
header_band(s, 10)
slide_title(s, "CadreHealth for hospitals", kicker="COMMERCIALS  /  CADREHEALTH")

left_x = Inches(0.6)
right_x = Inches(6.85)
col_w = Inches(6.0)
y = Inches(2.4)

add_rect(s, left_x, y, col_w, Inches(0.6), fill=NAVY)
tf = add_textbox(s, left_x + Inches(0.25), y + Inches(0.12), col_w - Inches(0.5), Inches(0.4))
set_text(tf, "EMPLOYER SUBSCRIPTIONS", size=14, bold=True, color=WHITE)

sub_tiers = [
    ("Verified Profile", "N250K / year", "Verified employer badge, basic facility profile, response to reviews."),
    ("Talent Search", "N1.5M / year", "Profile plus search across 4,300+ verified professionals, 5 mandates per year."),
    ("Workforce Intelligence", "N4M / year", "Talent search plus salary benchmarks, attrition data, anonymous hospital review insights."),
]
for i, (name, price, body) in enumerate(sub_tiers):
    rh = Inches(1.0)
    ry = y + Inches(0.6) + rh * i + Inches(0.05) * i
    bg = LIGHT if i % 2 == 0 else WHITE
    add_rect(s, left_x, ry, col_w, rh, fill=bg)
    tf = add_textbox(s, left_x + Inches(0.25), ry + Inches(0.1), col_w - Inches(0.5), Inches(0.35))
    set_text(tf, name, size=13, bold=True, color=NAVY)
    tf = add_textbox(s, left_x + Inches(0.25), ry + Inches(0.4), col_w - Inches(0.5), Inches(0.3))
    set_text(tf, price, size=14, bold=True, color=GOLD)
    tf = add_textbox(s, left_x + Inches(0.25), ry + Inches(0.7), col_w - Inches(0.5), Inches(0.3))
    set_text(tf, body, size=10, color=DARK)

add_rect(s, right_x, y, col_w, Inches(0.6), fill=TEAL)
tf = add_textbox(s, right_x + Inches(0.25), y + Inches(0.12), col_w - Inches(0.5), Inches(0.4))
set_text(tf, "PERFORMANCE FEES", size=14, bold=True, color=WHITE)

perf = [
    ("Permanent Placement", "20 to 25%", "Of first-year compensation. Pay only on hire. 90-day replacement guarantee."),
    ("Locum Placement", "10 to 15%", "Margin on locum shifts. Instant payment to professionals via Paystack."),
    ("International Recruitment", "Custom", "Diaspora and returning professionals. Negotiated per mandate."),
]
for i, (name, price, body) in enumerate(perf):
    rh = Inches(1.0)
    ry = y + Inches(0.6) + rh * i + Inches(0.05) * i
    bg = LIGHT if i % 2 == 0 else WHITE
    add_rect(s, right_x, ry, col_w, rh, fill=bg)
    tf = add_textbox(s, right_x + Inches(0.25), ry + Inches(0.1), col_w - Inches(0.5), Inches(0.35))
    set_text(tf, name, size=13, bold=True, color=NAVY)
    tf = add_textbox(s, right_x + Inches(0.25), ry + Inches(0.4), col_w - Inches(0.5), Inches(0.3))
    set_text(tf, price, size=14, bold=True, color=GOLD)
    tf = add_textbox(s, right_x + Inches(0.25), ry + Inches(0.7), col_w - Inches(0.5), Inches(0.3))
    set_text(tf, body, size=10, color=DARK)

tf = add_textbox(s, Inches(0.6), Inches(6.5), Inches(12.1), Inches(0.5))
set_text(tf, "ANCHOR: Traditional headhunters charge 25 to 35% with no replacement guarantee. Our subscription plus placement model is structurally cheaper at any hire volume above 2 per year.",
         size=11, bold=True, color=NAVY)
footer_band(s, 10)


# ============= SLIDE 11: BUNDLE PRICING (4 MECE TIERS) =============
s = prs.slides.add_slide(blank)
header_band(s, 11)
slide_title(s, "The integrated stack: best per-naira value", kicker="COMMERCIALS  /  BUNDLE")

profiles = [
    {"label": "UNDER 50 BEDS", "name": "Boutique Stack", "price": "N6M to N10M",
     "duration": "First year",
     "items": ["Maarova for 2 leaders (assessment + 1:1)", "CadreHealth Verified Profile",
               "2 placements credit", "Annual leadership check-in"]},
    {"label": "50 TO 150 BEDS", "name": "Foundation Stack", "price": "N15M to N25M",
     "duration": "First year",
     "items": ["Maarova for 5 leaders (assessment + 1:1)", "CadreHealth Talent Search subscription",
               "5 placements credit", "Quarterly leadership review"]},
    {"label": "150 TO 350 BEDS", "name": "Growth Stack", "price": "N35M to N55M",
     "duration": "First year",
     "items": ["Maarova Enterprise Cohort of 10 leaders", "Executive retreat included",
               "CadreHealth Workforce Intelligence", "12 placements credit", "Half-yearly board reporting"],
     "featured": True},
    {"label": "350+ BEDS / MULTI-SITE", "name": "Enterprise Stack", "price": "From N75M",
     "duration": "First year",
     "items": ["Maarova for 20+ leaders, multi-site", "Two retreats per year",
               "CadreHealth multi-facility deployment", "25+ placements credit",
               "Dedicated CFA engagement manager"]},
]
y = Inches(2.4)
card_w = Inches(2.95)
gap = Inches(0.15)
start_x = Inches(0.6)
for i, p in enumerate(profiles):
    x = start_x + (card_w + gap) * i
    featured = p.get("featured", False)
    fill = NAVY if featured else LIGHT
    text_main = WHITE if featured else DARK
    label_color = GOLD if featured else NAVY
    add_rect(s, x, y, card_w, Inches(4.4), fill=fill)
    if featured:
        add_rect(s, x, y, card_w, Inches(0.3), fill=GOLD)
        tf = add_textbox(s, x + Inches(0.15), y + Inches(0.02), card_w - Inches(0.3), Inches(0.25))
        set_text(tf, "RECOMMENDED", size=9, bold=True, color=NAVY)
        ty = y + Inches(0.4)
    else:
        ty = y + Inches(0.2)
    tf = add_textbox(s, x + Inches(0.15), ty, card_w - Inches(0.3), Inches(0.3))
    set_text(tf, p["label"], size=9, bold=True, color=label_color)
    tf = add_textbox(s, x + Inches(0.15), ty + Inches(0.3), card_w - Inches(0.3), Inches(0.5))
    set_text(tf, p["name"], size=17, bold=True, color=text_main)
    tf = add_textbox(s, x + Inches(0.15), ty + Inches(0.95), card_w - Inches(0.3), Inches(0.5))
    set_text(tf, p["price"], size=18, bold=True, color=GOLD)
    tf = add_textbox(s, x + Inches(0.15), ty + Inches(1.45), card_w - Inches(0.3), Inches(0.3))
    set_text(tf, p["duration"], size=10, color=text_main)
    item_tf = add_textbox(s, x + Inches(0.15), ty + Inches(1.85), card_w - Inches(0.3), Inches(2.3))
    set_text(item_tf, "• " + p["items"][0], size=10, color=text_main)
    for item in p["items"][1:]:
        add_paragraph(item_tf, "• " + item, size=10, color=text_main, space_before=Pt(4))
tf = add_textbox(s, Inches(0.6), Inches(6.95), Inches(12.1), Inches(0.3))
set_text(tf, "Placements above the credit pool charged at standard 20 to 25% of first-year compensation. Bundle savings of 15 to 25% versus à la carte.",
         size=10, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
footer_band(s, 11)


# ============= SLIDE 12: ROI MATH (Foundation Stack reference fixed) =============
s = prs.slides.add_slide(blank)
header_band(s, 12)
slide_title(s, "The ROI math", kicker="COMMERCIALS  /  PAYBACK")

scenarios = [
    {"title": "Avoid one bad senior hire", "save": "N20M+",
     "vs": "Foundation Stack: N15 to 25M",
     "verdict": "Pays for itself with one prevented mis-hire."},
    {"title": "Reduce attrition by 5%", "save": "N40M+",
     "vs": "200-bed hospital, 600 staff at avg N1M/yr churn cost",
     "verdict": "Growth Stack pays back in under 12 months."},
    {"title": "One leader who scales 3 more", "save": "Compounding",
     "vs": "Cost of NOT developing top leadership",
     "verdict": "The compounding return that no benchmark fully captures."},
]
y = Inches(2.4)
card_w = Inches(3.95)
gap = Inches(0.2)
start_x = Inches(0.6)
for i, sc in enumerate(scenarios):
    x = start_x + (card_w + gap) * i
    add_rect(s, x, y, card_w, Inches(3.8), fill=LIGHT)
    add_rect(s, x, y, card_w, Inches(0.08), fill=GOLD)
    tf = add_textbox(s, x + Inches(0.25), y + Inches(0.3), card_w - Inches(0.5), Inches(0.6))
    set_text(tf, sc["title"], size=15, bold=True, color=NAVY)
    tf = add_textbox(s, x + Inches(0.25), y + Inches(1.1), card_w - Inches(0.5), Inches(0.7))
    set_text(tf, "You save", size=11, color=GREY)
    add_paragraph(tf, sc["save"], size=32, bold=True, color=GOLD, space_before=Pt(2))
    tf = add_textbox(s, x + Inches(0.25), y + Inches(2.4), card_w - Inches(0.5), Inches(0.6))
    set_text(tf, sc["vs"], size=10, color=DARK)
    tf = add_textbox(s, x + Inches(0.25), y + Inches(3.05), card_w - Inches(0.5), Inches(0.6))
    set_text(tf, sc["verdict"], size=11, bold=True, color=TEAL)
tf = add_textbox(s, Inches(0.6), Inches(6.5), Inches(12.1), Inches(0.5))
set_text(tf, "Healthcare leadership and workforce decisions are not a cost line. They are the highest-leverage investment a hospital makes.",
         size=13, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
footer_band(s, 12)


# ============= SLIDE 13: HOW WE START =============
s = prs.slides.add_slide(blank)
header_band(s, 13)
slide_title(s, "How we start", kicker="ENGAGEMENT PATH")

steps = [
    ("Week 0", "Discovery Call", "30-minute call. We understand your context. No deck on our end, just questions."),
    ("Week 1 to 2", "Diagnostic", "Light-touch leadership and workforce diagnostic. Combination of interviews and platform data."),
    ("Week 3", "Tailored Proposal", "Stack recommendation with named coaches, scope, timeline, and fixed pricing in NGN."),
    ("Week 4+", "Phased Rollout", "Pilot with one cohort or facility. Scale based on signal. No multi-year lock-in."),
]
y = Inches(2.6)
card_w = Inches(2.95)
gap = Inches(0.15)
start_x = Inches(0.6)
for i, (when, title, body) in enumerate(steps):
    x = start_x + (card_w + gap) * i
    add_rect(s, x, y, card_w, Inches(3.6), fill=WHITE, line=NAVY)
    add_rect(s, x, y, card_w, Inches(0.6), fill=NAVY)
    tf = add_textbox(s, x + Inches(0.25), y + Inches(0.13), card_w - Inches(0.5), Inches(0.4))
    set_text(tf, when.upper(), size=11, bold=True, color=GOLD)
    tf = add_textbox(s, x + Inches(0.25), y + Inches(0.85), card_w - Inches(0.5), Inches(0.6))
    set_text(tf, title, size=18, bold=True, color=NAVY)
    tf = add_textbox(s, x + Inches(0.25), y + Inches(1.6), card_w - Inches(0.5), Inches(1.9))
    set_text(tf, body, size=11, color=DARK)
    badge = s.shapes.add_shape(MSO_SHAPE.OVAL, x + card_w - Inches(0.65), y - Inches(0.25), Inches(0.5), Inches(0.5))
    badge.fill.solid()
    badge.fill.fore_color.rgb = GOLD
    badge.line.fill.background()
    badge_tf = badge.text_frame
    badge_tf.margin_left = Inches(0)
    badge_tf.margin_right = Inches(0)
    set_text(badge_tf, str(i + 1), size=18, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
tf = add_textbox(s, Inches(0.6), Inches(6.5), Inches(12.1), Inches(0.5))
set_text(tf, "No retainer. No multi-year commitment. You buy outcomes, not seat licenses.",
         size=14, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
footer_band(s, 13)


# ============= SLIDE 14: WHY CFA FEE =============
s = prs.slides.add_slide(blank)
header_band(s, 14)

# Custom title (smaller spacing)
tf = add_textbox(s, Inches(0.6), Inches(0.8), Inches(12.1), Inches(0.4))
set_text(tf, "COMMERCIALS  /  WHY THE CFA FEE", size=11, bold=True, color=GOLD)
tf = add_textbox(s, Inches(0.6), Inches(1.15), Inches(12.1), Inches(0.7))
set_text(tf, "What our fee actually pays for", size=28, bold=True, color=NAVY)
add_rect(s, Inches(0.6), Inches(1.85), Inches(0.8), Inches(0.05), fill=GOLD)
tf = add_textbox(s, Inches(0.6), Inches(1.95), Inches(12.1), Inches(0.4))
set_text(tf, "Four layers of value the hospital does not build, hire, or carry alone.",
         size=14, color=DARK)

pillars = [
    {"color": NAVY, "tag": "01", "title": "Vetted Network", "subtitle": "We have already done the sourcing",
     "items": ["4,300+ healthcare professionals, 16 cadres",
               "Senior CFA consultant bench across functions",
               "ICF-credentialed coaching panel",
               "License verified, reference checked, ID confirmed",
               "Skip the 6 to 12 month recruitment cycle"]},
    {"color": TEAL, "tag": "02", "title": "Statutory Machinery", "subtitle": "Payroll, tax, pension, indemnity",
     "items": ["PAYE, pension, NHIS, NSITF, ITF, VAT, WHT",
               "FIRS remittance and audit-ready records",
               "Indemnity cover bundled (Turaco)",
               "NDPR-compliant data handling",
               "Contracts, IP, and exit management"]},
    {"color": GOLD, "tag": "03", "title": "Risk Backstop", "subtitle": "We carry the wrong-hire risk",
     "items": ["Replacement guarantee on every engagement",
               "Performance accountability we own, not pass through",
               "Escalation line to CFA partners",
               "Substitute cover for cancelled locum shifts",
               "If it goes wrong, it is our problem to fix"]},
    {"color": DEEP_TEAL, "tag": "04", "title": "Institutional Capability", "subtitle": "You hire the firm, not one person",
     "items": ["Playbooks and frameworks the consultant brings",
               "Cross-portfolio learning from other engagements",
               "Coaching and development during placement",
               "Advisory access from CFA partners",
               "Knowledge stays even when the person changes"]},
]
cards_y = Inches(2.55)
card_w = Inches(2.95)
card_h = Inches(3.0)
gap = Inches(0.15)
start_x = Inches(0.6)
for i, p in enumerate(pillars):
    x = start_x + (card_w + gap) * i
    add_rect(s, x, cards_y, card_w, card_h, fill=WHITE, line=SOFT)
    add_rect(s, x, cards_y, card_w, Inches(0.7), fill=p["color"])
    tf = add_textbox(s, x + Inches(0.15), cards_y + Inches(0.05), Inches(0.6), Inches(0.4))
    set_text(tf, p["tag"], size=14, bold=True, color=GOLD)
    tf = add_textbox(s, x + Inches(0.6), cards_y + Inches(0.05), card_w - Inches(0.7), Inches(0.4))
    set_text(tf, p["title"], size=15, bold=True, color=WHITE)
    tf = add_textbox(s, x + Inches(0.15), cards_y + Inches(0.42), card_w - Inches(0.3), Inches(0.25))
    set_text(tf, p["subtitle"], size=10, color=GOLD)
    item_tf = add_textbox(s, x + Inches(0.15), cards_y + Inches(0.85), card_w - Inches(0.3), Inches(2.05))
    set_text(item_tf, "• " + p["items"][0], size=10, color=DARK)
    for item in p["items"][1:]:
        add_paragraph(item_tf, "• " + item, size=10, color=DARK, space_before=Pt(4))

# Comparison strip
comp_y = Inches(5.7)
add_rect(s, Inches(0.6), comp_y, Inches(12.1), Inches(0.42), fill=NAVY)
tf = add_textbox(s, Inches(0.85), comp_y + Inches(0.07), Inches(11.6), Inches(0.3))
set_text(tf, "VS. THE ALTERNATIVES", size=11, bold=True, color=GOLD)
table_y = comp_y + Inches(0.42)
col_widths = [Inches(3.0), Inches(3.0), Inches(3.05), Inches(3.05)]
col_x = [Inches(0.6)]
for w in col_widths[:-1]:
    col_x.append(col_x[-1] + w)
hdrs = ["WHAT YOU NEED", "HIRE DIRECT", "TRADITIONAL AGENCY", "WITH CFA"]
hdr_colors = [LIGHT, LIGHT, LIGHT, GOLD]
hdr_text = [NAVY, GREY, GREY, NAVY]
for i, h in enumerate(hdrs):
    add_rect(s, col_x[i], table_y, col_widths[i], Inches(0.32), fill=hdr_colors[i])
    tf = add_textbox(s, col_x[i] + Inches(0.12), table_y + Inches(0.05), col_widths[i] - Inches(0.2), Inches(0.25))
    set_text(tf, h, size=10, bold=True, color=hdr_text[i])
rows = [
    ("Time to fill", "6 to 12 months", "2 to 3 months", "2 to 4 weeks"),
    ("Compliance burden", "All yours", "All yours", "All ours"),
    ("Replacement risk", "All yours", "Yours after 90 days", "Ours, ongoing"),
]
ry = table_y + Inches(0.32)
row_h = Inches(0.27)
for j, row in enumerate(rows):
    row_y = ry + row_h * j
    bg = WHITE if j % 2 == 0 else LIGHT
    for i, cell in enumerate(row):
        if i == 3:
            cell_bg = GOLD; cell_color = NAVY; bold = True
        else:
            cell_bg = bg; cell_color = NAVY if i == 0 else DARK; bold = (i == 0)
        add_rect(s, col_x[i], row_y, col_widths[i], row_h, fill=cell_bg)
        tf = add_textbox(s, col_x[i] + Inches(0.12), row_y + Inches(0.04), col_widths[i] - Inches(0.2), Inches(0.22))
        set_text(tf, cell, size=10, bold=bold, color=cell_color)
footer_band(s, 14)


# ============= SLIDE 15: MID-LEVEL + LOCUM (CORRECTED RATES) =============
s = prs.slides.add_slide(blank)
header_band(s, 15)
tf = add_textbox(s, Inches(0.6), Inches(0.85), Inches(12.1), Inches(0.4))
set_text(tf, "COMMERCIALS  /  MID-LEVEL SECONDMENTS AND LOCUM COVER", size=11, bold=True, color=GOLD)
tf = add_textbox(s, Inches(0.6), Inches(1.2), Inches(12.1), Inches(0.7))
set_text(tf, "Operational managers and clinical cover, on demand", size=26, bold=True, color=NAVY)
add_rect(s, Inches(0.6), Inches(1.95), Inches(0.8), Inches(0.05), fill=GOLD)

top_y = Inches(2.2)
add_rect(s, Inches(0.6), top_y, Inches(12.1), Inches(0.42), fill=NAVY)
tf = add_textbox(s, Inches(0.85), top_y + Inches(0.07), Inches(11.6), Inches(0.3))
set_text(tf, "MID-LEVEL OPERATIONAL SECONDMENTS  /  MONTHLY RATE", size=11, bold=True, color=GOLD)

mid_roles = [
    ("Finance Manager", "N1.8M to N3M", "N700K to N1.2M", "Month-end close, cash flow, payor reconciliation, board pack."),
    ("Supply Chain Manager", "N1.5M to N2.5M", "N600K to N1M", "Procurement systems, vendor terms, stockout reduction, KPI dashboards."),
    ("Quality & Compliance Manager", "N1.5M to N2.5M", "N600K to N1M", "JCI / SafeCare prep, clinical audit cycles, incident review systems."),
    ("HR Manager", "N1.5M to N2.5M", "N600K to N1M", "Payroll discipline, performance reviews, retention, NMA / NARD navigation."),
    ("IT / Health Information Manager", "N1.8M to N3M", "N700K to N1.2M", "EMR rollouts, data integrity, cybersecurity baseline, HMIS reporting."),
]
y = top_y + Inches(0.42)
row_h = Inches(0.45)
header_y = y
add_rect(s, Inches(0.6), header_y, Inches(12.1), Inches(0.32), fill=LIGHT)
hdrs = [(Inches(0.85), Inches(3.5), "ROLE"),
        (Inches(4.4), Inches(2.2), "FULL-TIME / MO"),
        (Inches(6.65), Inches(2.0), "FRACTIONAL / MO"),
        (Inches(8.7), Inches(4.0), "WHAT THEY OWN")]
for x, w, label in hdrs:
    tf = add_textbox(s, x, header_y + Inches(0.06), w, Inches(0.25))
    set_text(tf, label, size=9, bold=True, color=GREY)
y = header_y + Inches(0.32)
for i, (role, ft, frac, body) in enumerate(mid_roles):
    row_y = y + row_h * i
    bg = WHITE if i % 2 == 0 else LIGHT
    add_rect(s, Inches(0.6), row_y, Inches(12.1), row_h, fill=bg)
    tf = add_textbox(s, Inches(0.85), row_y + Inches(0.1), Inches(3.5), Inches(0.3))
    set_text(tf, role, size=11, bold=True, color=NAVY)
    tf = add_textbox(s, Inches(4.4), row_y + Inches(0.1), Inches(2.2), Inches(0.3))
    set_text(tf, ft, size=11, bold=True, color=GOLD)
    tf = add_textbox(s, Inches(6.65), row_y + Inches(0.1), Inches(2.0), Inches(0.3))
    set_text(tf, frac, size=11, bold=True, color=TEAL)
    tf = add_textbox(s, Inches(8.7), row_y + Inches(0.12), Inches(4.0), Inches(0.3))
    set_text(tf, body, size=9, color=DARK)

bot_y = Inches(5.05)
add_rect(s, Inches(0.6), bot_y, Inches(7.7), Inches(0.42), fill=TEAL)
tf = add_textbox(s, Inches(0.85), bot_y + Inches(0.07), Inches(7.2), Inches(0.3))
set_text(tf, "LOCUM CLINICAL COVER  /  PER 8-12 HOUR SHIFT", size=11, bold=True, color=WHITE)

# CORRECTED rates - Nigerian market reality
locum_roles = [
    ("Consultant Specialist", "N80K to N200K"),
    ("Medical Officer", "N20K to N50K"),
    ("ICU / Theatre Nurse", "N15K to N30K"),
    ("General Nurse", "N10K to N20K"),
    ("Pharmacist", "N20K to N40K"),
    ("Lab Scientist", "N15K to N30K"),
    ("Radiographer", "N25K to N50K"),
]
ly = bot_y + Inches(0.42)
lrow_h = Inches(0.27)
for i, (role, rate) in enumerate(locum_roles):
    row_y = ly + lrow_h * i
    bg = WHITE if i % 2 == 0 else LIGHT
    add_rect(s, Inches(0.6), row_y, Inches(7.7), lrow_h, fill=bg)
    tf = add_textbox(s, Inches(0.85), row_y + Inches(0.04), Inches(4.5), Inches(0.22))
    set_text(tf, role, size=10, bold=True, color=NAVY)
    tf = add_textbox(s, Inches(5.4), row_y + Inches(0.04), Inches(2.8), Inches(0.22))
    set_text(tf, rate, size=10, bold=True, color=GOLD)

rx = Inches(8.5)
add_rect(s, rx, bot_y, Inches(4.2), Inches(0.42), fill=NAVY)
tf = add_textbox(s, rx + Inches(0.2), bot_y + Inches(0.07), Inches(4.0), Inches(0.3))
set_text(tf, "HOW LOCUM COVER WORKS", size=11, bold=True, color=GOLD)
how_y = bot_y + Inches(0.42)
add_rect(s, rx, how_y, Inches(4.2), Inches(1.45), fill=LIGHT)
tf = add_textbox(s, rx + Inches(0.2), how_y + Inches(0.1), Inches(4.0), Inches(1.3))
set_text(tf, "• Pay only when you book a shift", size=10, color=DARK)
add_paragraph(tf, "• 10 to 15% platform margin on top of locum rate", size=10, color=DARK, space_before=Pt(5))
add_paragraph(tf, "• 4,300+ verified professionals across 16 cadres", size=10, color=DARK, space_before=Pt(5))
add_paragraph(tf, "• Indemnity insurance bundled (Turaco)", size=10, color=DARK, space_before=Pt(5))
add_paragraph(tf, "• Paystack instant payout to professional", size=10, color=DARK, space_before=Pt(5))

ex_y = how_y + Inches(1.5)
add_rect(s, rx, ex_y, Inches(4.2), Inches(0.32), fill=GOLD)
tf = add_textbox(s, rx + Inches(0.2), ex_y + Inches(0.04), Inches(4.0), Inches(0.25))
set_text(tf, "EXAMPLE MONTHLY SPEND", size=10, bold=True, color=NAVY)
ex_box_y = ex_y + Inches(0.32)
add_rect(s, rx, ex_box_y, Inches(4.2), Inches(0.95), fill=WHITE, line=GREY)
tf = add_textbox(s, rx + Inches(0.2), ex_box_y + Inches(0.08), Inches(4.0), Inches(0.85))
set_text(tf, "20 nurse shifts + 8 MO shifts + 4 specialist days", size=9, bold=True, color=NAVY)
add_paragraph(tf, "≈ N900K to N1.5M / month, fully covered", size=11, bold=True, color=GOLD, space_before=Pt(4))
add_paragraph(tf, "vs. permanent backfill cost: N2.5M+ / month committed", size=9, color=DARK, space_before=Pt(4))
footer_band(s, 15)


# ============= SLIDE 16: ECONOMICS TRANSPARENCY (CORRECTED LOCUM) =============
s = prs.slides.add_slide(blank)
header_band(s, 16)

tf = add_textbox(s, Inches(0.6), Inches(0.85), Inches(12.1), Inches(0.4))
set_text(tf, "COMMERCIALS  /  ECONOMICS AND TRANSPARENCY", size=11, bold=True, color=GOLD)
tf = add_textbox(s, Inches(0.6), Inches(1.2), Inches(12.1), Inches(0.7))
set_text(tf, "Where every naira goes: hospital, taxes, CFA, professional", size=22, bold=True, color=NAVY)
add_rect(s, Inches(0.6), Inches(1.95), Inches(0.8), Inches(0.05), fill=GOLD)


def waterfall_row(slide, x, y, w, label, amount, kind="neutral"):
    color_map = {"gross": NAVY, "cfa": GOLD, "pro": TEAL, "net": NAVY, "deduction": GREY, "neutral": GREY}
    text_color_map = {"gross": WHITE, "cfa": NAVY, "pro": WHITE, "net": WHITE, "deduction": DARK, "neutral": DARK}
    bg = color_map.get(kind, GREY)
    fg = text_color_map.get(kind, DARK)
    h = Inches(0.34)
    add_rect(slide, x, y, w, h, fill=bg)
    tf = add_textbox(slide, x + Inches(0.12), y + Inches(0.05), w - Inches(2.2), Inches(0.25))
    set_text(tf, label, size=10, bold=True, color=fg)
    tf = add_textbox(slide, x + w - Inches(2.05), y + Inches(0.05), Inches(1.9), Inches(0.25))
    set_text(tf, amount, size=10, bold=True, color=fg, align=PP_ALIGN.RIGHT)


def light_row(slide, x, y, w, label, amount):
    h = Inches(0.3)
    add_rect(slide, x, y, w, h, fill=LIGHT)
    tf = add_textbox(slide, x + Inches(0.12), y + Inches(0.04), w - Inches(2.0), Inches(0.22))
    set_text(tf, label, size=9, color=DARK)
    tf = add_textbox(slide, x + w - Inches(1.85), y + Inches(0.04), Inches(1.7), Inches(0.22))
    set_text(tf, amount, size=9, bold=True, color=DARK, align=PP_ALIGN.RIGHT)


card_w = Inches(4.1)
gap = Inches(0.15)
cards_x = Inches(0.6)
cards_y = Inches(2.2)
card_h = Inches(4.65)

# Card 1: Senior Fractional
x = cards_x
add_rect(s, x, cards_y, card_w, card_h, fill=WHITE, line=SOFT)
add_rect(s, x, cards_y, card_w, Inches(0.55), fill=NAVY)
tf = add_textbox(s, x + Inches(0.15), cards_y + Inches(0.05), card_w - Inches(0.3), Inches(0.25))
set_text(tf, "SENIOR FRACTIONAL", size=9, bold=True, color=GOLD)
tf = add_textbox(s, x + Inches(0.15), cards_y + Inches(0.27), card_w - Inches(0.3), Inches(0.3))
set_text(tf, "Fractional CMO at N5M / month", size=12, bold=True, color=WHITE)
ry = cards_y + Inches(0.7)
waterfall_row(s, x + Inches(0.15), ry, card_w - Inches(0.3), "Hospital invoice (ex VAT)", "N5,000,000", "gross"); ry += Inches(0.38)
light_row(s, x + Inches(0.15), ry, card_w - Inches(0.3), "+ VAT 7.5% (input credit if registered)", "N375,000"); ry += Inches(0.32)
light_row(s, x + Inches(0.15), ry, card_w - Inches(0.3), "- WHT 10% (remitted to FIRS)", "N500,000"); ry += Inches(0.4)
waterfall_row(s, x + Inches(0.15), ry, card_w - Inches(0.3), "CFA platform fee (25%)", "N1,250,000", "cfa"); ry += Inches(0.42)
waterfall_row(s, x + Inches(0.15), ry, card_w - Inches(0.3), "Professional gross (contractor)", "N3,750,000", "pro"); ry += Inches(0.38)
light_row(s, x + Inches(0.15), ry, card_w - Inches(0.3), "- Professional own PIT (~25%)", "N940,000"); ry += Inches(0.4)
waterfall_row(s, x + Inches(0.15), ry, card_w - Inches(0.3), "Professional net cash", "N2,810,000", "net")
tf = add_textbox(s, x + Inches(0.15), cards_y + Inches(4.2), card_w - Inches(0.3), Inches(0.4))
set_text(tf, "Senior contractor model. Professional invoices CFA, handles own taxes.", size=8, color=GREY)

# Card 2: Mid-Level Secondment
x = cards_x + (card_w + gap)
add_rect(s, x, cards_y, card_w, card_h, fill=WHITE, line=SOFT)
add_rect(s, x, cards_y, card_w, Inches(0.55), fill=TEAL)
tf = add_textbox(s, x + Inches(0.15), cards_y + Inches(0.05), card_w - Inches(0.3), Inches(0.25))
set_text(tf, "MID-LEVEL SECONDMENT", size=9, bold=True, color=GOLD)
tf = add_textbox(s, x + Inches(0.15), cards_y + Inches(0.27), card_w - Inches(0.3), Inches(0.3))
set_text(tf, "Finance Manager at N2M / month", size=12, bold=True, color=WHITE)
ry = cards_y + Inches(0.7)
waterfall_row(s, x + Inches(0.15), ry, card_w - Inches(0.3), "Hospital invoice (ex VAT)", "N2,000,000", "gross"); ry += Inches(0.38)
light_row(s, x + Inches(0.15), ry, card_w - Inches(0.3), "+ VAT 7.5%", "N150,000"); ry += Inches(0.32)
light_row(s, x + Inches(0.15), ry, card_w - Inches(0.3), "- WHT 10% (FIRS)", "N200,000"); ry += Inches(0.4)
waterfall_row(s, x + Inches(0.15), ry, card_w - Inches(0.3), "CFA platform fee (20%)", "N400,000", "cfa"); ry += Inches(0.38)
light_row(s, x + Inches(0.15), ry, card_w - Inches(0.3), "Employer pension + NSITF + ITF", "N190,000"); ry += Inches(0.4)
waterfall_row(s, x + Inches(0.15), ry, card_w - Inches(0.3), "Professional gross salary", "N1,410,000", "pro"); ry += Inches(0.38)
light_row(s, x + Inches(0.15), ry, card_w - Inches(0.3), "- PAYE + pension + NHIS", "N390,000"); ry += Inches(0.4)
waterfall_row(s, x + Inches(0.15), ry, card_w - Inches(0.3), "Professional take-home", "N1,020,000", "net")
tf = add_textbox(s, x + Inches(0.15), cards_y + Inches(4.2), card_w - Inches(0.3), Inches(0.4))
set_text(tf, "CFA is employer of record. Full statutory compliance handled.", size=8, color=GREY)

# Card 3: Locum (CORRECTED to N35K MO shift)
x = cards_x + (card_w + gap) * 2
add_rect(s, x, cards_y, card_w, card_h, fill=WHITE, line=SOFT)
add_rect(s, x, cards_y, card_w, Inches(0.55), fill=GOLD)
tf = add_textbox(s, x + Inches(0.15), cards_y + Inches(0.05), card_w - Inches(0.3), Inches(0.25))
set_text(tf, "LOCUM CLINICAL COVER", size=9, bold=True, color=NAVY)
tf = add_textbox(s, x + Inches(0.15), cards_y + Inches(0.27), card_w - Inches(0.3), Inches(0.3))
set_text(tf, "Medical Officer shift at N35K", size=12, bold=True, color=NAVY)
ry = cards_y + Inches(0.7)
waterfall_row(s, x + Inches(0.15), ry, card_w - Inches(0.3), "Hospital pays per shift", "N35,000", "gross"); ry += Inches(0.38)
light_row(s, x + Inches(0.15), ry, card_w - Inches(0.3), "- WHT 5% (individual, FIRS)", "N1,750"); ry += Inches(0.4)
waterfall_row(s, x + Inches(0.15), ry, card_w - Inches(0.3), "CFA platform margin (12%)", "N4,200", "cfa"); ry += Inches(0.42)
waterfall_row(s, x + Inches(0.15), ry, card_w - Inches(0.3), "Professional payout (instant)", "N29,050", "pro"); ry += Inches(0.38)
light_row(s, x + Inches(0.15), ry, card_w - Inches(0.3), "Indemnity (Turaco) bundled", "Included"); ry += Inches(0.32)
light_row(s, x + Inches(0.15), ry, card_w - Inches(0.3), "Professional own annual PIT", "Self-managed"); ry += Inches(0.4)
waterfall_row(s, x + Inches(0.15), ry, card_w - Inches(0.3), "Net cash to professional", "N29,050", "net")
tf = add_textbox(s, x + Inches(0.15), cards_y + Inches(4.2), card_w - Inches(0.3), Inches(0.4))
set_text(tf, "Per-shift contractor. Paid via Paystack within 24 hours of shift.", size=8, color=GREY)

# Bottom strip
strip_y = Inches(7.0)
add_rect(s, 0, strip_y, SLIDE_W, Inches(0.2), fill=NAVY)
tf = add_textbox(s, Inches(0.6), strip_y + Inches(0.0), Inches(12.1), Inches(0.2))
set_text(tf, "Statutory taxes (VAT, WHT, PAYE, pension) flow to government and pension funds. CFA fee covers sourcing, vetting, payroll, indemnity, and ongoing support.",
         size=9, color=WHITE, align=PP_ALIGN.CENTER)
footer_band(s, 16)


# ============= SLIDE 17: NEXT STEP / CTA =============
s = prs.slides.add_slide(blank)
add_rect(s, 0, 0, SLIDE_W, SLIDE_H, fill=NAVY)
# Logo top-left
s.shapes.add_picture(LOGO_FULL, Inches(0.6), Inches(0.6), height=Inches(0.7))
# CTA
tf = add_textbox(s, Inches(0.6), Inches(2.3), Inches(12), Inches(0.5))
set_text(tf, "NEXT STEP", size=14, bold=True, color=GOLD)
tf = add_textbox(s, Inches(0.6), Inches(2.8), Inches(12), Inches(2.0))
set_text(tf, "Book a 30-minute discovery call.", size=42, bold=True, color=WHITE)
add_paragraph(tf, "We will map your specific gaps and send a tailored proposal within 7 days.",
              size=22, color=GOLD, space_before=Pt(15))

# Contact block (with logo on right side)
add_rect(s, Inches(0.6), Inches(5.4), Inches(12.1), Inches(1.6), fill=WHITE)
tf = add_textbox(s, Inches(0.85), Inches(5.55), Inches(11.6), Inches(0.4))
set_text(tf, "CONTACT", size=11, bold=True, color=GOLD)
tf = add_textbox(s, Inches(0.85), Inches(5.9), Inches(5.5), Inches(1.0))
set_text(tf, "Consult For Africa", size=18, bold=True, color=NAVY)
add_paragraph(tf, "consultforafrica.com", size=13, color=DARK, space_before=Pt(5))
add_paragraph(tf, "hello@consultforafrica.com", size=13, color=DARK, space_before=Pt(3))
tf = add_textbox(s, Inches(7.0), Inches(5.9), Inches(4.0), Inches(1.0))
set_text(tf, "Maarova", size=14, bold=True, color=NAVY)
add_paragraph(tf, "consultforafrica.com/maarova", size=12, color=DARK, space_before=Pt(3))
add_paragraph(tf, "CadreHealth", size=14, bold=True, color=NAVY, space_before=Pt(8))
add_paragraph(tf, "consultforafrica.com/oncadre", size=12, color=DARK, space_before=Pt(3))
# Logo on right of contact card
s.shapes.add_picture(LOGO_FULL, Inches(11.0), Inches(5.85), height=Inches(0.95))

# Tagline
tf = add_textbox(s, Inches(0.6), Inches(7.15), Inches(12), Inches(0.4))
set_text(tf, "African healthcare that runs better.", size=14, color=GOLD, align=PP_ALIGN.CENTER)


# Save
output_path = "/Users/debo/consult-for-africa/docs/cfa-pitch-deck-final.pptx"
prs.save(output_path)
print(f"Final deck saved: {output_path}")
print(f"Total slides: {len(prs.slides)}")
