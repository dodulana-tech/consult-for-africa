"""Build standalone 2-slide deck on secondments and fractional hires.

Designed to be merged into cfa-pitch-deck.pptx. Same brand palette
and layout language so the slides drop in without restyling.

Slide 1: What secondments and fractional hires are, and when you use each
Slide 2: Commercial examples - what you can get for what price
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

# Match main deck palette
NAVY = RGBColor(0x0B, 0x3C, 0x5D)
GOLD = RGBColor(0xD4, 0xAF, 0x37)
TEAL = RGBColor(0x1F, 0x7A, 0x8C)
LIGHT = RGBColor(0xF4, 0xF7, 0xFA)
DARK = RGBColor(0x11, 0x18, 0x27)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GREY = RGBColor(0x6B, 0x72, 0x80)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def set_text(frame, text, size=18, bold=False, color=DARK, align=PP_ALIGN.LEFT, font="Calibri"):
    frame.clear()
    p = frame.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font


def add_paragraph(frame, text, size=14, bold=False, color=DARK, align=PP_ALIGN.LEFT, space_before=Pt(6), font="Calibri"):
    p = frame.add_paragraph()
    p.alignment = align
    p.space_before = space_before
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font


def add_rect(slide, x, y, w, h, fill=NAVY, line=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
    shape.shadow.inherit = False
    return shape


def add_textbox(slide, x, y, w, h):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.05)
    tf.margin_bottom = Inches(0.05)
    return tf


def header_band(slide, label):
    add_rect(slide, 0, 0, SLIDE_W, Inches(0.6), fill=NAVY)
    tf = add_textbox(slide, Inches(0.6), Inches(0.1), Inches(8), Inches(0.4))
    set_text(tf, "CONSULT FOR AFRICA", size=11, bold=True, color=WHITE)
    tf2 = add_textbox(slide, Inches(8.733), Inches(0.1), Inches(4.4), Inches(0.4))
    set_text(tf2, label.upper(), size=11, bold=True, color=GOLD, align=PP_ALIGN.RIGHT)


def footer_band(slide):
    add_rect(slide, 0, Inches(7.2), SLIDE_W, Inches(0.3), fill=LIGHT)
    tf = add_textbox(slide, Inches(0.6), Inches(7.22), Inches(8), Inches(0.25))
    set_text(tf, "consultforafrica.com  |  Maarova  |  CadreHealth", size=9, color=GREY)


def slide_title(slide, title, kicker=None):
    y = Inches(0.9)
    if kicker:
        tf = add_textbox(slide, Inches(0.6), y, Inches(12.1), Inches(0.4))
        set_text(tf, kicker, size=12, bold=True, color=GOLD)
        y = Inches(1.25)
    tf = add_textbox(slide, Inches(0.6), y, Inches(12.1), Inches(0.8))
    set_text(tf, title, size=32, bold=True, color=NAVY)
    add_rect(slide, Inches(0.6), Inches(2.05), Inches(0.8), Inches(0.05), fill=GOLD)


prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
blank = prs.slide_layouts[6]


# ============= SLIDE A: SECONDMENTS AND FRACTIONAL =============
s = prs.slides.add_slide(blank)
header_band(s, "Insert slide")
slide_title(s, "When you need senior capability now, not in 6 months",
            kicker="THIRD MOTION  /  SECONDMENTS AND FRACTIONAL")

# Two product cards
card_w = Inches(5.95)
card_h = Inches(4.5)
y = Inches(2.5)

# Secondment
add_rect(s, Inches(0.6), y, card_w, card_h, fill=NAVY)
add_rect(s, Inches(0.6), y, card_w, Inches(0.5), fill=GOLD)
tf = add_textbox(s, Inches(0.85), y + Inches(0.05), card_w - Inches(0.5), Inches(0.4))
set_text(tf, "SECONDMENT", size=14, bold=True, color=NAVY)
tf = add_textbox(s, Inches(0.85), y + Inches(0.7), card_w - Inches(0.5), Inches(0.5))
set_text(tf, "Full-time, embedded", size=20, bold=True, color=WHITE)
tf = add_textbox(s, Inches(0.85), y + Inches(1.4), card_w - Inches(0.5), Inches(2.9))
set_text(tf, "A CFA consultant deployed full-time into your hospital. They sit in your office, attend your stand-ups, own a P&L line.",
         size=12, color=WHITE)
add_paragraph(tf, "TYPICAL DURATION", size=10, bold=True, color=GOLD, space_before=Pt(14))
add_paragraph(tf, "3 to 12 months. Most engagements run 6 months with optional extension.",
              size=12, color=WHITE, space_before=Pt(4))
add_paragraph(tf, "BEST FOR", size=10, bold=True, color=GOLD, space_before=Pt(10))
add_paragraph(tf, "• Interim COO during turnaround", size=11, color=WHITE, space_before=Pt(4))
add_paragraph(tf, "• Embedded transformation lead", size=11, color=WHITE, space_before=Pt(2))
add_paragraph(tf, "• Acting CHRO during HR overhaul", size=11, color=WHITE, space_before=Pt(2))

# Fractional
x2 = Inches(0.6) + card_w + Inches(0.2)
add_rect(s, x2, y, card_w, card_h, fill=TEAL)
add_rect(s, x2, y, card_w, Inches(0.5), fill=GOLD)
tf = add_textbox(s, x2 + Inches(0.25), y + Inches(0.05), card_w - Inches(0.5), Inches(0.4))
set_text(tf, "FRACTIONAL", size=14, bold=True, color=NAVY)
tf = add_textbox(s, x2 + Inches(0.25), y + Inches(0.7), card_w - Inches(0.5), Inches(0.5))
set_text(tf, "Senior, part-time, ongoing", size=20, bold=True, color=WHITE)
tf = add_textbox(s, x2 + Inches(0.25), y + Inches(1.4), card_w - Inches(0.5), Inches(2.9))
set_text(tf, "A senior CFA leader works with you 1 to 3 days a week over a longer horizon. Strategic capability without the cost of a full-time hire.",
         size=12, color=WHITE)
add_paragraph(tf, "TYPICAL DURATION", size=10, bold=True, color=GOLD, space_before=Pt(14))
add_paragraph(tf, "6 to 12 months, often renewed. 1 to 3 days per week.",
              size=12, color=WHITE, space_before=Pt(4))
add_paragraph(tf, "BEST FOR", size=10, bold=True, color=GOLD, space_before=Pt(10))
add_paragraph(tf, "• Fractional CMO during clinical strategy reset", size=11, color=WHITE, space_before=Pt(4))
add_paragraph(tf, "• Fractional CFO for financial discipline", size=11, color=WHITE, space_before=Pt(2))
add_paragraph(tf, "• Fractional CHRO for talent transformation", size=11, color=WHITE, space_before=Pt(2))

# Bottom punchline
tf = add_textbox(s, Inches(0.6), Inches(7.0), Inches(12.1), Inches(0.3))
set_text(tf, "Senior capability you can hire by the month, not the headcount line.",
         size=13, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

footer_band(s)


# ============= SLIDE B: COMMERCIALS WITH EXAMPLES =============
s = prs.slides.add_slide(blank)
header_band(s, "Insert slide")
slide_title(s, "Senior roles, monthly cost, what you get",
            kicker="COMMERCIALS  /  SECONDMENTS AND FRACTIONAL")

# Pricing rows
roles = [
    ("Fractional Chief Medical Officer", "N4M to N7M / month",
     "2 to 3 days a week. Clinical strategy, governance, quality oversight, medical leadership.", "FRACTIONAL"),
    ("Fractional CHRO", "N3M to N5M / month",
     "2 to 3 days a week. Workforce strategy, retention, succession planning, HR systems.", "FRACTIONAL"),
    ("Fractional CFO", "N3M to N6M / month",
     "2 days a week. Financial discipline, cash management, payor strategy, board reporting.", "FRACTIONAL"),
    ("Interim COO (Secondment)", "N6M to N10M / month",
     "Full-time, 6 to 12 months. Daily ops, turnaround, capacity planning, vendor renegotiation.", "SECONDMENT"),
    ("Embedded Transformation Lead", "N5M to N8M / month",
     "Full-time, 9 to 12 months. Owns delivery of a named transformation programme end to end.", "SECONDMENT"),
]
y = Inches(2.4)
row_h = Inches(0.7)
for i, (name, price, body, tag) in enumerate(roles):
    row_y = y + row_h * i
    bg = LIGHT if i % 2 == 0 else WHITE
    add_rect(s, Inches(0.6), row_y, Inches(12.1), row_h, fill=bg)
    # Tag pill
    tag_color = TEAL if tag == "FRACTIONAL" else NAVY
    add_rect(s, Inches(0.85), row_y + Inches(0.22), Inches(0.95), Inches(0.28), fill=tag_color)
    tf = add_textbox(s, Inches(0.85), row_y + Inches(0.22), Inches(0.95), Inches(0.28))
    set_text(tf, tag, size=8, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    # Role name
    tf = add_textbox(s, Inches(1.95), row_y + Inches(0.1), Inches(4.0), Inches(0.5))
    set_text(tf, name, size=13, bold=True, color=NAVY)
    # Price
    tf = add_textbox(s, Inches(1.95), row_y + Inches(0.4), Inches(3.0), Inches(0.3))
    set_text(tf, price, size=12, bold=True, color=GOLD)
    # Body
    tf = add_textbox(s, Inches(5.2), row_y + Inches(0.18), Inches(7.4), Inches(0.5))
    set_text(tf, body, size=10, color=DARK)

# Three example bundles
y2 = Inches(6.05)
add_rect(s, Inches(0.6), y2, Inches(12.1), Inches(0.4), fill=NAVY)
tf = add_textbox(s, Inches(0.85), y2 + Inches(0.06), Inches(11.6), Inches(0.3))
set_text(tf, "EXAMPLE BUNDLES  /  WHAT YOU GET FOR WHAT", size=11, bold=True, color=GOLD)

bundle_y = y2 + Inches(0.45)
bundles = [
    ("N40M", "6-month Fractional CMO + Maarova for top 5 leaders"),
    ("N75M", "12-month Interim COO + CadreHealth talent search + 10 mid-manager hires"),
    ("N120M+", "Full transformation: 12-month Embedded Lead + Maarova cohort of 10 + CadreHealth"),
]
bw = Inches(3.95)
gap = Inches(0.2)
for i, (price, desc) in enumerate(bundles):
    bx = Inches(0.6) + (bw + gap) * i
    add_rect(s, bx, bundle_y, bw, Inches(0.6), fill=LIGHT)
    tf = add_textbox(s, bx + Inches(0.15), bundle_y + Inches(0.08), Inches(1.5), Inches(0.45))
    set_text(tf, price, size=16, bold=True, color=GOLD)
    tf = add_textbox(s, bx + Inches(1.6), bundle_y + Inches(0.05), bw - Inches(1.7), Inches(0.55))
    set_text(tf, desc, size=10, color=DARK)

footer_band(s)


# Save
output_path = "/Users/debo/consult-for-africa/docs/cfa-secondments-fractional-slides.pptx"
prs.save(output_path)
print(f"Standalone slides saved: {output_path}")
print(f"Slides: {len(prs.slides)}")
