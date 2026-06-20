"""V2: corrected locum rates based on Nigerian market reality.

Calibration:
- MO base salary CONMESS N200-300K/mo, private N300-500K/mo
- Daily equivalent ~N15-25K, locum premium 1.5-2.5x daily
- Specialist consultant private salary N800K-3M/mo, locum daily ~1.5x

Old rates were 2-3x too high. This V2 corrects them.
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

NAVY = RGBColor(0x0B, 0x3C, 0x5D)
GOLD = RGBColor(0xD4, 0xAF, 0x37)
TEAL = RGBColor(0x1F, 0x7A, 0x8C)
LIGHT = RGBColor(0xF4, 0xF7, 0xFA)
DARK = RGBColor(0x11, 0x18, 0x27)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GREY = RGBColor(0x6B, 0x72, 0x80)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def set_text(frame, text, size=18, bold=False, color=DARK, align=PP_ALIGN.LEFT, font="Calibri"):
    frame.clear()
    p = frame.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font


def add_paragraph(frame, text, size=14, bold=False, color=DARK, align=PP_ALIGN.LEFT, space_before=Pt(6), font="Calibri"):
    p = frame.add_paragraph()
    p.alignment = align
    p.space_before = space_before
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font


def add_rect(slide, x, y, w, h, fill=NAVY, line=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
    shape.shadow.inherit = False
    return shape


def add_textbox(slide, x, y, w, h):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.05)
    tf.margin_bottom = Inches(0.05)
    return tf


prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
blank = prs.slide_layouts[6]
s = prs.slides.add_slide(blank)

add_rect(s, 0, 0, SLIDE_W, Inches(0.6), fill=NAVY)
tf = add_textbox(s, Inches(0.6), Inches(0.1), Inches(8), Inches(0.4))
set_text(tf, "CONSULT FOR AFRICA", size=11, bold=True, color=WHITE)
tf = add_textbox(s, Inches(8.733), Inches(0.1), Inches(4.4), Inches(0.4))
set_text(tf, "INSERT SLIDE", size=11, bold=True, color=GOLD, align=PP_ALIGN.RIGHT)

tf = add_textbox(s, Inches(0.6), Inches(0.85), Inches(12.1), Inches(0.4))
set_text(tf, "COMMERCIALS  /  MID-LEVEL SECONDMENTS AND LOCUM COVER", size=11, bold=True, color=GOLD)
tf = add_textbox(s, Inches(0.6), Inches(1.2), Inches(12.1), Inches(0.7))
set_text(tf, "Operational managers and clinical cover, on demand", size=28, bold=True, color=NAVY)
add_rect(s, Inches(0.6), Inches(1.95), Inches(0.8), Inches(0.05), fill=GOLD)

# ============= TOP: Mid-level secondments =============
top_y = Inches(2.2)
add_rect(s, Inches(0.6), top_y, Inches(12.1), Inches(0.42), fill=NAVY)
tf = add_textbox(s, Inches(0.85), top_y + Inches(0.07), Inches(11.6), Inches(0.3))
set_text(tf, "MID-LEVEL OPERATIONAL SECONDMENTS  /  MONTHLY RATE", size=11, bold=True, color=GOLD)

mid_roles = [
    ("Finance Manager", "N1.8M to N3M", "N700K to N1.2M", "Month-end close, cash flow, payor reconciliation, board pack."),
    ("Supply Chain Manager", "N1.5M to N2.5M", "N600K to N1M", "Procurement systems, vendor terms, stockout reduction, KPI dashboards."),
    ("Quality & Compliance Manager", "N1.5M to N2.5M", "N600K to N1M", "JCI / SafeCare prep, clinical audit cycles, incident review systems."),
    ("HR Manager", "N1.5M to N2.5M", "N600K to N1M", "Payroll discipline, performance reviews, retention, NMA / NARD navigation."),
    ("IT / Health Information Manager", "N1.8M to N3M", "N700K to N1.2M", "EMR rollouts, data integrity, cybersecurity baseline, HMIS reporting."),
]
y = top_y + Inches(0.42)
row_h = Inches(0.45)

header_y = y
add_rect(s, Inches(0.6), header_y, Inches(12.1), Inches(0.32), fill=LIGHT)
hdrs = [(Inches(0.85), Inches(3.5), "ROLE"),
        (Inches(4.4), Inches(2.2), "FULL-TIME / MO"),
        (Inches(6.65), Inches(2.0), "FRACTIONAL / MO"),
        (Inches(8.7), Inches(4.0), "WHAT THEY OWN")]
for x, w, label in hdrs:
    tf = add_textbox(s, x, header_y + Inches(0.06), w, Inches(0.25))
    set_text(tf, label, size=9, bold=True, color=GREY)

y = header_y + Inches(0.32)
for i, (role, ft, frac, body) in enumerate(mid_roles):
    row_y = y + row_h * i
    bg = WHITE if i % 2 == 0 else LIGHT
    add_rect(s, Inches(0.6), row_y, Inches(12.1), row_h, fill=bg)
    tf = add_textbox(s, Inches(0.85), row_y + Inches(0.1), Inches(3.5), Inches(0.3))
    set_text(tf, role, size=11, bold=True, color=NAVY)
    tf = add_textbox(s, Inches(4.4), row_y + Inches(0.1), Inches(2.2), Inches(0.3))
    set_text(tf, ft, size=11, bold=True, color=GOLD)
    tf = add_textbox(s, Inches(6.65), row_y + Inches(0.1), Inches(2.0), Inches(0.3))
    set_text(tf, frac, size=11, bold=True, color=TEAL)
    tf = add_textbox(s, Inches(8.7), row_y + Inches(0.12), Inches(4.0), Inches(0.3))
    set_text(tf, body, size=9, color=DARK)

# ============= BOTTOM: Locum cover with CORRECTED rates =============
bot_y = Inches(5.05)
add_rect(s, Inches(0.6), bot_y, Inches(7.7), Inches(0.42), fill=TEAL)
tf = add_textbox(s, Inches(0.85), bot_y + Inches(0.07), Inches(7.2), Inches(0.3))
set_text(tf, "LOCUM CLINICAL COVER  /  PER 8-12 HOUR SHIFT", size=11, bold=True, color=WHITE)

# Calibrated to Nigerian private market reality
locum_roles = [
    ("Consultant Specialist", "N80K to N200K"),
    ("Medical Officer", "N20K to N50K"),
    ("ICU / Theatre Nurse", "N15K to N30K"),
    ("General Nurse", "N10K to N20K"),
    ("Pharmacist", "N20K to N40K"),
    ("Lab Scientist", "N15K to N30K"),
    ("Radiographer", "N25K to N50K"),
]
ly = bot_y + Inches(0.42)
lrow_h = Inches(0.27)
for i, (role, rate) in enumerate(locum_roles):
    row_y = ly + lrow_h * i
    bg = WHITE if i % 2 == 0 else LIGHT
    add_rect(s, Inches(0.6), row_y, Inches(7.7), lrow_h, fill=bg)
    tf = add_textbox(s, Inches(0.85), row_y + Inches(0.04), Inches(4.5), Inches(0.22))
    set_text(tf, role, size=10, bold=True, color=NAVY)
    tf = add_textbox(s, Inches(5.4), row_y + Inches(0.04), Inches(2.8), Inches(0.22))
    set_text(tf, rate, size=10, bold=True, color=GOLD)

# Right side
rx = Inches(8.5)
add_rect(s, rx, bot_y, Inches(4.2), Inches(0.42), fill=NAVY)
tf = add_textbox(s, rx + Inches(0.2), bot_y + Inches(0.07), Inches(4.0), Inches(0.3))
set_text(tf, "HOW LOCUM COVER WORKS", size=11, bold=True, color=GOLD)

how_y = bot_y + Inches(0.42)
add_rect(s, rx, how_y, Inches(4.2), Inches(1.45), fill=LIGHT)
tf = add_textbox(s, rx + Inches(0.2), how_y + Inches(0.1), Inches(4.0), Inches(1.3))
set_text(tf, "• Pay only when you book a shift", size=10, color=DARK)
add_paragraph(tf, "• 10 to 15% platform margin on top of locum rate", size=10, color=DARK, space_before=Pt(5))
add_paragraph(tf, "• 4,300+ verified professionals across 16 cadres", size=10, color=DARK, space_before=Pt(5))
add_paragraph(tf, "• Indemnity insurance bundled (Turaco)", size=10, color=DARK, space_before=Pt(5))
add_paragraph(tf, "• Paystack instant payout to professional", size=10, color=DARK, space_before=Pt(5))

# Recomputed example monthly spend
ex_y = how_y + Inches(1.5)
add_rect(s, rx, ex_y, Inches(4.2), Inches(0.32), fill=GOLD)
tf = add_textbox(s, rx + Inches(0.2), ex_y + Inches(0.04), Inches(4.0), Inches(0.25))
set_text(tf, "EXAMPLE MONTHLY SPEND", size=10, bold=True, color=NAVY)
ex_box_y = ex_y + Inches(0.32)
add_rect(s, rx, ex_box_y, Inches(4.2), Inches(0.95), fill=WHITE, line=GREY)
tf = add_textbox(s, rx + Inches(0.2), ex_box_y + Inches(0.08), Inches(4.0), Inches(0.85))
# 20 nurse @ N15K + 8 MO @ N35K + 4 specialist @ N120K = N300K + N280K + N480K = ~N1.06M
set_text(tf, "20 nurse shifts + 8 MO shifts + 4 specialist days", size=9, bold=True, color=NAVY)
add_paragraph(tf, "≈ N900K to N1.5M / month, fully covered", size=11, bold=True, color=GOLD, space_before=Pt(4))
add_paragraph(tf, "vs. permanent backfill cost: N2.5M+ / month committed", size=9, color=DARK, space_before=Pt(4))

add_rect(s, 0, Inches(7.2), SLIDE_W, Inches(0.3), fill=LIGHT)
tf = add_textbox(s, Inches(0.6), Inches(7.22), Inches(8), Inches(0.25))
set_text(tf, "consultforafrica.com  |  Maarova  |  CadreHealth", size=9, color=GREY)

output_path = "/Users/debo/consult-for-africa/docs/cfa-midlevel-locum-slide-v2.pptx"
prs.save(output_path)
print(f"Saved: {output_path}")
