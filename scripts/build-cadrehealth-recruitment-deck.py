"""Build CadreHealth Recruitment commercial deck.

Generates three PPTX variants:
  docs/cadrehealth-recruitment-deck.pptx          (generic master)
  docs/cadrehealth-recruitment-deck-pearl.pptx    (Pearl Oncology variant)
  docs/cadrehealth-recruitment-deck-osiris.pptx   (Osiris Health variant)

Run:
  python3 scripts/build-cadrehealth-recruitment-deck.py

Audience: hospital MD/CEO/CMD/HR Director.
Constraints:
- No em dashes anywhere in copy.
- No "AI-powered" or "AI-driven" language.
- NGN primary, USD never needed for these clients.
- Senior medic register: gravitas, specific, no hype.
"""
from __future__ import annotations

from dataclasses import dataclass
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN


ROOT = Path(__file__).resolve().parents[1]
DOCS = ROOT / "docs"


# ---------- palette ---------------------------------------------------------

NAVY = RGBColor(0x0B, 0x3C, 0x5D)
GOLD = RGBColor(0xD4, 0xAF, 0x37)
TEAL = RGBColor(0x1F, 0x7A, 0x8C)
LIGHT = RGBColor(0xF4, 0xF7, 0xFA)
DARK = RGBColor(0x11, 0x18, 0x27)
INK = RGBColor(0x06, 0x09, 0x0F)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GREY = RGBColor(0x6B, 0x72, 0x80)
SOFT_GREY = RGBColor(0xE5, 0xE7, 0xEB)
SUCCESS = RGBColor(0x05, 0x96, 0x69)
WARN = RGBColor(0xD9, 0x73, 0x06)
RED = RGBColor(0xC0, 0x39, 0x2B)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


# ---------- variant config --------------------------------------------------

@dataclass
class Variant:
    slug: str
    cover_title: str
    cover_sub: str
    client_block_eyebrow: str
    client_block_title: str
    client_block_body: list[str]
    client_block_callout: str
    next_steps: list[tuple[str, str]]
    include_partnership_slide: bool
    partnership_lines: list[str]


GENERIC = Variant(
    slug="generic",
    cover_title="Healthcare Recruitment, Done Differently",
    cover_sub="CadreHealth Recruitment, by Consult For Africa",
    client_block_eyebrow="FOR HOSPITAL LEADERSHIP",
    client_block_title="Built for hospitals carrying clinical scarcity",
    client_block_body=[
        "Your hiring funnel is leaking the wrong way. Senior staff abroad, junior CVs that overstate, and four week shortlists that arrive with two interested candidates.",
        "This deck explains what we do, how we price it, and how we begin in the next two weeks if you decide to engage.",
    ],
    client_block_callout="Three engagement tiers. One platform. One partner.",
    next_steps=[
        ("This week", "Confirm intent and select an engagement tier."),
        ("Next week", "Kickoff call. We agree open roles and screening criteria."),
        ("Week 2", "First role intake. JD calibration and salary band confirmed."),
        ("Week 5 to 6", "First shortlist delivered."),
    ],
    include_partnership_slide=False,
    partnership_lines=[],
)


PEARL = Variant(
    slug="pearl",
    cover_title="A Recruitment Proposal for\nPearl Oncology Specialist Hospital",
    cover_sub="Prepared by Consult For Africa, May 2026",
    client_block_eyebrow="FOR PEARL ONCOLOGY",
    client_block_title="A specialist hospital deserves a specialist recruiter",
    client_block_body=[
        "Pearl Oncology has built a specialist cancer facility in Lekki Phase 1 with four cancer centres and over 400 patients treated since 2021. The clinical promise is clear; the hiring challenge is real, and it is the same one facing every specialist hospital in Lagos.",
        "Your open Medical Officer roles at NGN 400,000+ per month sit in a market where Pearl is competing for the same MBBS-holders as five other specialist centres. Speed, screening quality, and a credible network decide who fills the rota first.",
        "This proposal is built around your immediate Medical Officer need, with a pathway to support Pearl as you scale across the four centres.",
    ],
    client_block_callout="Pilot pricing: 15% on your first three Medical Officer placements.",
    next_steps=[
        ("Week 0", "Kemi shares the Medical Officer JD and confirms salary band."),
        ("Week 1", "Kickoff call with Dr Salako or the HR team. SLA signed."),
        ("Week 2 to 3", "Sourcing live. Verified credentials, oncology-context screening."),
        ("Week 4 to 6", "First shortlist of 3 to 5 candidates with assessment notes."),
        ("Week 7 to 10", "Interviews, offer, acceptance. Replacement guarantee from day one."),
    ],
    include_partnership_slide=False,
    partnership_lines=[],
)


OSIRIS = Variant(
    slug="osiris",
    cover_title="A Recruitment & Partnership Proposal for\nOsiris Health",
    cover_sub="Prepared by Consult For Africa, May 2026",
    client_block_eyebrow="FOR OSIRIS HEALTH",
    client_block_title="A multi-site dialysis operator hiring across four cadres",
    client_block_body=[
        "Osiris Health, manager of Life Support Medical Centre, is one of Nigeria's earliest renal care providers. With the recent R-Jolad collaboration and active expansion across Lagos, Ibadan, and Warri, your hiring problem is no longer single-role; it is a multi-cadre, multi-site clinical scaling exercise.",
        "Your current open positions span Operations Manager, Consultant Nephrologist, Medical Officers, and Dialysis Nurses across three cities. That mix needs more than a contingency recruiter. It needs a partner who can hold the retainer-tier search for the Consultant role while delivering volume on the nursing side.",
        "We are proposing a hybrid engagement that prices each role correctly, plus a parallel partnership conversation that reflects what you have signalled to us.",
    ],
    client_block_callout="Hybrid engagement: retainer for the Consultant, volume pricing for nursing, partnership track in parallel.",
    next_steps=[
        ("Week 0", "Ayoola confirms full open headcount by site and salary bands."),
        ("Week 1", "Two kickoff calls. One commercial (SLA), one strategic (partnership)."),
        ("Week 2", "Consultant Nephrologist retainer search begins. Dialysis Nurse batch sourcing live."),
        ("Week 3 to 4", "First nursing shortlists for Ibadan and Warri."),
        ("Week 4 to 6", "Consultant Nephrologist shortlist of 3 names."),
        ("Week 6", "Partnership MoU drafted in parallel."),
    ],
    include_partnership_slide=True,
    partnership_lines=[
        "Founding Hospital Partner status, with first-mover positioning on the CadreHealth platform.",
        "Revenue share on CadreHealth Connect subscriptions sold into the Osiris network (R-Jolad and other affiliates). Indicative: 15% to 20% of subscription revenue, capped for 24 months.",
        "Co-development input on renal-care workforce features: dialysis nurse credentialing, nephrology fellowship register, in-centre staffing analytics.",
        "Equity participation: explicitly scoped as a separate conversation, not embedded in the recruitment SLA. Reserved for a dedicated meeting with both principals.",
        "Logo and case-study rights with attribution. First claim to publish the partnership story.",
    ],
)


# ---------- helpers ---------------------------------------------------------

def add_rect(slide, x, y, w, h, fill=NAVY, line=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
    shape.shadow.inherit = False
    return shape


def add_textbox(slide, x, y, w, h):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.04)
    tf.margin_bottom = Inches(0.04)
    return tf


def set_text(frame, text, size=18, bold=False, color=DARK, align=PP_ALIGN.LEFT, font="Calibri"):
    frame.clear()
    p = frame.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font


def add_paragraph(frame, text, size=14, bold=False, color=DARK, align=PP_ALIGN.LEFT,
                  space_before=Pt(6), font="Calibri"):
    p = frame.add_paragraph()
    p.alignment = align
    p.space_before = space_before
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font


def header_band(slide, label, total_pages):
    add_rect(slide, 0, 0, SLIDE_W, Inches(0.55), fill=NAVY)
    tf = add_textbox(slide, Inches(0.5), Inches(0.1), Inches(8), Inches(0.4))
    set_text(tf, "CADREHEALTH RECRUITMENT  /  by Consult For Africa", size=10, bold=True, color=WHITE)
    tf2 = add_textbox(slide, Inches(8.733), Inches(0.1), Inches(4.4), Inches(0.4))
    set_text(tf2, label.upper(), size=10, bold=True, color=GOLD, align=PP_ALIGN.RIGHT)


def footer_band(slide, page_num, total_pages):
    add_rect(slide, 0, Inches(7.2), SLIDE_W, Inches(0.3), fill=LIGHT)
    tf = add_textbox(slide, Inches(0.5), Inches(7.22), Inches(8), Inches(0.25))
    set_text(tf, "consultforafrica.com  |  oncadre.com  |  hello@consultforafrica.com",
             size=9, color=GREY)
    tf2 = add_textbox(slide, Inches(8.733), Inches(7.22), Inches(4.4), Inches(0.25))
    set_text(tf2, f"{page_num:02d} of {total_pages:02d}", size=9, color=GREY, align=PP_ALIGN.RIGHT)


def slide_title(slide, title, kicker=None, y=Inches(0.85)):
    if kicker:
        tf = add_textbox(slide, Inches(0.55), y, Inches(12.2), Inches(0.4))
        set_text(tf, kicker, size=11, bold=True, color=GOLD)
        y = Inches(1.2)
    tf = add_textbox(slide, Inches(0.55), y, Inches(12.2), Inches(0.9))
    set_text(tf, title, size=30, bold=True, color=NAVY)
    add_rect(slide, Inches(0.55), Inches(2.0), Inches(0.7), Inches(0.05), fill=GOLD)


def new_slide(prs, label, page_num, total_pages, with_chrome=True):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    if with_chrome:
        header_band(s, label, total_pages)
        footer_band(s, page_num, total_pages)
    return s


def bullet_card(slide, x, y, w, h, eyebrow, title, body_lines, accent=GOLD):
    add_rect(slide, x, y, w, h, fill=WHITE, line=SOFT_GREY)
    add_rect(slide, x, y, Inches(0.1), h, fill=accent)
    tf = add_textbox(slide, x + Inches(0.3), y + Inches(0.2), w - Inches(0.4), Inches(0.3))
    set_text(tf, eyebrow.upper(), size=10, bold=True, color=accent)
    tf = add_textbox(slide, x + Inches(0.3), y + Inches(0.55), w - Inches(0.4), Inches(0.5))
    set_text(tf, title, size=16, bold=True, color=NAVY)
    tf = add_textbox(slide, x + Inches(0.3), y + Inches(1.1), w - Inches(0.4), h - Inches(1.2))
    for i, line in enumerate(body_lines):
        if i == 0:
            set_text(tf, "•  " + line, size=11, color=DARK)
        else:
            add_paragraph(tf, "•  " + line, size=11, color=DARK, space_before=Pt(6))


# ---------- slide builders --------------------------------------------------

def slide_cover(prs, variant: Variant):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(s, 0, 0, SLIDE_W, SLIDE_H, fill=NAVY)

    # gold accent bar
    add_rect(s, Inches(0.6), Inches(2.3), Inches(1.5), Inches(0.08), fill=GOLD)

    # eyebrow
    tf = add_textbox(s, Inches(0.6), Inches(1.85), Inches(10), Inches(0.4))
    set_text(tf, "CADREHEALTH RECRUITMENT", size=14, bold=True, color=GOLD)

    # title
    tf = add_textbox(s, Inches(0.6), Inches(2.6), Inches(12.2), Inches(2.3))
    lines = variant.cover_title.split("\n")
    set_text(tf, lines[0], size=42, bold=True, color=WHITE)
    for ln in lines[1:]:
        add_paragraph(tf, ln, size=42, bold=True, color=WHITE, space_before=Pt(0))

    # subtitle
    tf = add_textbox(s, Inches(0.6), Inches(5.2), Inches(12), Inches(0.6))
    set_text(tf, variant.cover_sub, size=18, color=GOLD)

    # footer
    tf = add_textbox(s, Inches(0.6), Inches(6.7), Inches(12), Inches(0.4))
    set_text(tf, "consultforafrica.com  |  oncadre.com", size=12, color=WHITE)


def slide_client_block(prs, variant: Variant, page_num, total):
    s = new_slide(prs, "For You", page_num, total)
    slide_title(s, variant.client_block_title, kicker=variant.client_block_eyebrow)

    tf = add_textbox(s, Inches(0.55), Inches(2.3), Inches(12.2), Inches(3.5))
    set_text(tf, variant.client_block_body[0], size=14, color=DARK)
    for body in variant.client_block_body[1:]:
        add_paragraph(tf, body, size=14, color=DARK, space_before=Pt(12))

    # callout
    add_rect(s, Inches(0.55), Inches(6.0), Inches(12.2), Inches(0.85), fill=NAVY)
    tf = add_textbox(s, Inches(0.85), Inches(6.18), Inches(11.6), Inches(0.55))
    set_text(tf, variant.client_block_callout, size=16, bold=True, color=GOLD)


def slide_the_problem(prs, page_num, total):
    s = new_slide(prs, "The Problem", page_num, total)
    slide_title(s, "Why this conversation, now", kicker="THE HIRING REALITY")

    stats = [
        ("16,156", "Nigerian-trained nurses on the UK NMC register as of September 2025.", GOLD),
        ("57,000", "Nigerian nurses left for the UK alone over the last 5 years.", TEAL),
        ("42,000", "Nigerian nurses departed in 2022 to 2024. Nigeria is now the third largest source of foreign-trained UK nurses.", NAVY),
        ("8 to 12", "weeks the average hospital loses recruiting each clinical replacement, often longer for specialists.", RED),
    ]
    card_w = Inches(2.95)
    card_h = Inches(2.55)
    gap = Inches(0.15)
    start_x = Inches(0.55)
    start_y = Inches(2.4)

    for i, (stat, label, color) in enumerate(stats):
        x = start_x + (card_w + gap) * i
        add_rect(s, x, start_y, card_w, card_h, fill=LIGHT)
        add_rect(s, x, start_y, card_w, Inches(0.1), fill=color)
        tf = add_textbox(s, x + Inches(0.2), start_y + Inches(0.3), card_w - Inches(0.4), Inches(0.9))
        set_text(tf, stat, size=36, bold=True, color=color)
        tf = add_textbox(s, x + Inches(0.2), start_y + Inches(1.35), card_w - Inches(0.4), Inches(1.1))
        set_text(tf, label, size=11, color=DARK)

    tf = add_textbox(s, Inches(0.55), Inches(5.4), Inches(12.2), Inches(1.5))
    set_text(tf, "The bottom of the funnel is leaking.", size=18, bold=True, color=NAVY)
    add_paragraph(
        tf,
        "Every senior nurse who lands at Heathrow is one your competitor will also fail to replace. The hospitals that solve clinical scarcity first will win the next 24 months.",
        size=14, color=DARK, space_before=Pt(10),
    )


def slide_cost_of_vacancy(prs, page_num, total):
    s = new_slide(prs, "Cost of Vacancy", page_num, total)
    slide_title(s, "What an unfilled clinical role costs you each month",
                kicker="THE QUIET DRAIN")

    rows = [
        ("Specialist consultant vacancy", "NGN 5M to 12M per month",
         "Locum cover at 2 to 3x salary equivalent. Referral cases redirected to competing hospitals. Theatre and clinic capacity capped."),
        ("Medical Officer vacancy in a 24/7 service", "NGN 1.2M to 3M per month",
         "Overtime on retained MOs. Patient throughput compressed. Junior doctor burnout accelerates and triggers further attrition."),
        ("Senior Nurse or Dialysis Nurse vacancy", "NGN 800k to 1.8M per month",
         "Agency nurse premium plus higher error rates. Each empty chair in a dialysis centre is direct lost revenue at NGN 35k+ per session."),
        ("Compounding cost: one bad hire", "NGN 18M to 40M",
         "Industry benchmark: 6 to 12 months of compensation. Recruiting, onboarding, severance, lost output, team morale recovery."),
    ]
    y = Inches(2.25)
    for i, (label, cost, detail) in enumerate(rows):
        row_y = y + Inches(1.15) * i
        bg = LIGHT if i % 2 == 0 else WHITE
        add_rect(s, Inches(0.55), row_y, Inches(12.2), Inches(1.05), fill=bg)
        tf = add_textbox(s, Inches(0.8), row_y + Inches(0.1), Inches(4.6), Inches(0.35))
        set_text(tf, label, size=13, bold=True, color=NAVY)
        tf = add_textbox(s, Inches(0.8), row_y + Inches(0.45), Inches(4.6), Inches(0.5))
        set_text(tf, cost, size=18, bold=True, color=GOLD)
        tf = add_textbox(s, Inches(5.6), row_y + Inches(0.2), Inches(7.0), Inches(0.8))
        set_text(tf, detail, size=11, color=DARK)


def slide_whats_broken(prs, page_num, total):
    s = new_slide(prs, "Why Current Hiring Fails", page_num, total)
    slide_title(s, "Why your current hiring is not working", kicker="THE DIAGNOSIS")

    points = [
        ("Inflated CVs", "Generic job boards return CVs polished by an AI tool and a Lagos Whatsapp group. Verifying takes longer than sourcing."),
        ("Ghost candidates", "Two of every five shortlisted candidates stop replying after the first contact. Your HR team chases air."),
        ("Salary opacity", "You discover at offer stage that your band is 30% below market. The role reopens; another 4 weeks lost."),
        ("Credential ambiguity", "MDCN and NMCN status is verified, if at all, after the interview. Forgeries surface only in incident review."),
        ("Diaspora ignored", "Nigerian consultants in the UK and US are open to remote retainer or fractional return. Almost no Nigerian agency reaches them well."),
        ("No retention signal", "You hire and hope. Three months later the candidate is gone and you start again. No data on why."),
    ]
    cols = 3
    rows = 2
    card_w = Inches(4.05)
    card_h = Inches(2.1)
    gap_x = Inches(0.15)
    gap_y = Inches(0.2)
    start_x = Inches(0.55)
    start_y = Inches(2.3)

    for i, (title, body) in enumerate(points):
        col = i % cols
        row = i // cols
        x = start_x + (card_w + gap_x) * col
        y = start_y + (card_h + gap_y) * row
        add_rect(s, x, y, card_w, card_h, fill=WHITE, line=SOFT_GREY)
        add_rect(s, x, y, card_w, Inches(0.08), fill=GOLD)
        tf = add_textbox(s, x + Inches(0.25), y + Inches(0.2), card_w - Inches(0.4), Inches(0.5))
        set_text(tf, title, size=15, bold=True, color=NAVY)
        tf = add_textbox(s, x + Inches(0.25), y + Inches(0.75), card_w - Inches(0.4), card_h - Inches(0.9))
        set_text(tf, body, size=11, color=DARK)


def slide_intro_cadrehealth(prs, page_num, total):
    s = new_slide(prs, "Our Service", page_num, total)
    slide_title(s, "CadreHealth Recruitment, by Consult For Africa",
                kicker="WHAT WE ARE")

    tf = add_textbox(s, Inches(0.55), Inches(2.3), Inches(12.2), Inches(1.8))
    set_text(
        tf,
        "We are a healthcare-only recruitment service, built on top of the CadreHealth platform: a career platform for Nigerian healthcare professionals across 16 cadres.",
        size=16, color=DARK,
    )
    add_paragraph(
        tf,
        "When a hospital engages us, we are not searching the open internet on your behalf. We are matching your role against a registered, credentialed, salary-anchored network of doctors, nurses, pharmacists, and allied health professionals who have already declared their cadre, location, experience, and career intent.",
        size=13, color=DARK, space_before=Pt(12),
    )

    # two columns: platform vs service
    col_w = Inches(5.9)
    col_h = Inches(2.3)
    col_y = Inches(4.6)

    add_rect(s, Inches(0.55), col_y, col_w, col_h, fill=LIGHT)
    add_rect(s, Inches(0.55), col_y, col_w, Inches(0.08), fill=GOLD)
    tf = add_textbox(s, Inches(0.75), col_y + Inches(0.2), col_w - Inches(0.4), Inches(0.4))
    set_text(tf, "THE PLATFORM", size=11, bold=True, color=GOLD)
    tf = add_textbox(s, Inches(0.75), col_y + Inches(0.6), col_w - Inches(0.4), Inches(0.4))
    set_text(tf, "CadreHealth (oncadre.com)", size=16, bold=True, color=NAVY)
    tf = add_textbox(s, Inches(0.75), col_y + Inches(1.05), col_w - Inches(0.4), col_h - Inches(1.1))
    set_text(tf, "Active members across 16 cadres. Real salary data, credential verification, hospital reviews, diaspora mentorship layer.",
             size=12, color=DARK)

    add_rect(s, Inches(6.85), col_y, col_w, col_h, fill=NAVY)
    add_rect(s, Inches(6.85), col_y, col_w, Inches(0.08), fill=GOLD)
    tf = add_textbox(s, Inches(7.05), col_y + Inches(0.2), col_w - Inches(0.4), Inches(0.4))
    set_text(tf, "THE SERVICE", size=11, bold=True, color=GOLD)
    tf = add_textbox(s, Inches(7.05), col_y + Inches(0.6), col_w - Inches(0.4), Inches(0.4))
    set_text(tf, "CadreHealth Recruitment", size=16, bold=True, color=WHITE)
    tf = add_textbox(s, Inches(7.05), col_y + Inches(1.05), col_w - Inches(0.4), col_h - Inches(1.1))
    set_text(tf, "Consult For Africa's healthcare recruitment offering. Three engagement tiers, one platform, one accountable partner.",
             size=12, color=WHITE)


def slide_platform_advantage(prs, page_num, total):
    s = new_slide(prs, "Platform Advantage", page_num, total)
    slide_title(s, "What CadreHealth gives our recruitment team",
                kicker="THE PLATFORM ADVANTAGE")

    points = [
        ("16 cadres in network", "Doctors, Nurses, Midwives, Pharmacists, Lab Scientists, Radiographers, Physiotherapists, Dentists, Optometrists, CHOs/CHEWs, and more."),
        ("Real salary data", "Cadre-by-cadre, city-by-city, facility-tier-by-tier. We anchor your offer correctly before you start interviewing."),
        ("Credential verification", "MDCN, NMCN, PCN, MLSCN status checked on-platform. Renewals and CPD tracked."),
        ("Hospital review intel", "We know which facilities haemorrhage staff. We screen candidates accordingly and brief you honestly."),
        ("Career Readiness Score", "Cadre-specific score with separate UK, US, Canada, Gulf, and domestic dimensions. We can tell you who is staying versus leaving in 18 months."),
        ("Diaspora layer", "Direct lines into MANSAG, ANPA, and Diaspora Friends Council. Real access to consultants in the UK, US, Canada open to fractional and full return."),
    ]
    cols = 2
    card_w = Inches(6.05)
    card_h = Inches(1.55)
    gap_x = Inches(0.15)
    gap_y = Inches(0.2)
    start_x = Inches(0.55)
    start_y = Inches(2.3)

    for i, (title, body) in enumerate(points):
        col = i % cols
        row = i // cols
        x = start_x + (card_w + gap_x) * col
        y = start_y + (card_h + gap_y) * row
        add_rect(s, x, y, card_w, card_h, fill=LIGHT)
        add_rect(s, x, y, Inches(0.08), card_h, fill=GOLD)
        tf = add_textbox(s, x + Inches(0.25), y + Inches(0.15), card_w - Inches(0.4), Inches(0.4))
        set_text(tf, title, size=14, bold=True, color=NAVY)
        tf = add_textbox(s, x + Inches(0.25), y + Inches(0.6), card_w - Inches(0.4), card_h - Inches(0.7))
        set_text(tf, body, size=11, color=DARK)


def slide_three_tiers_overview(prs, page_num, total):
    s = new_slide(prs, "Three Ways to Engage", page_num, total)
    slide_title(s, "Three ways to engage us", kicker="ENGAGEMENT MODELS")

    tiers = [
        ("TIER 01", "Pay-Per-Hire", "Standard contingency. You only pay on placement.",
         "Best for: ad-hoc clinical and operational hires. MOs, nurses, allied health, mid-level admin.",
         "15% to 20% of first-year gross", NAVY),
        ("TIER 02", "Retainer (Engaged Search)", "We commit a dedicated partner. You commit a small upfront.",
         "Best for: Consultants, Heads of Department, C-suite, scarce specialties.",
         "25% of first-year gross + NGN 500k retainer", TEAL),
        ("TIER 03", "CadreHealth Connect", "Annual subscription. Unlimited reach on the platform.",
         "Best for: hospitals that prefer to self-recruit but want our network and infrastructure.",
         "NGN 100,000 per year", GOLD),
    ]
    card_w = Inches(4.05)
    card_h = Inches(4.4)
    gap = Inches(0.15)
    start_x = Inches(0.55)
    start_y = Inches(2.3)

    for i, (tier_label, name, headline, best_for, price, color) in enumerate(tiers):
        x = start_x + (card_w + gap) * i
        add_rect(s, x, start_y, card_w, card_h, fill=WHITE, line=SOFT_GREY)
        add_rect(s, x, start_y, card_w, Inches(0.65), fill=color)
        tf = add_textbox(s, x + Inches(0.25), start_y + Inches(0.13), card_w - Inches(0.4), Inches(0.4))
        set_text(tf, tier_label, size=11, bold=True, color=WHITE)
        tf = add_textbox(s, x + Inches(0.25), start_y + Inches(0.85), card_w - Inches(0.4), Inches(0.5))
        set_text(tf, name, size=18, bold=True, color=NAVY)
        tf = add_textbox(s, x + Inches(0.25), start_y + Inches(1.45), card_w - Inches(0.4), Inches(0.8))
        set_text(tf, headline, size=12, color=DARK)
        tf = add_textbox(s, x + Inches(0.25), start_y + Inches(2.4), card_w - Inches(0.4), Inches(1.2))
        set_text(tf, best_for, size=11, color=GREY)
        # price band
        add_rect(s, x, start_y + card_h - Inches(0.85), card_w, Inches(0.85), fill=LIGHT)
        tf = add_textbox(s, x + Inches(0.25), start_y + card_h - Inches(0.7), card_w - Inches(0.4), Inches(0.5))
        set_text(tf, price, size=13, bold=True, color=NAVY)


def slide_tier1_detail(prs, page_num, total):
    s = new_slide(prs, "Tier 01 / Pay-Per-Hire", page_num, total)
    slide_title(s, "Tier 1: Pay-Per-Hire", kicker="STANDARD CONTINGENCY")

    # left: structure
    add_rect(s, Inches(0.55), Inches(2.3), Inches(6.0), Inches(4.5), fill=LIGHT)
    tf = add_textbox(s, Inches(0.85), Inches(2.5), Inches(5.4), Inches(0.4))
    set_text(tf, "WHAT IT COSTS", size=11, bold=True, color=GOLD)
    tf = add_textbox(s, Inches(0.85), Inches(2.9), Inches(5.4), Inches(3.7))
    set_text(tf, "20% of first-year gross salary, before VAT.", size=14, bold=True, color=NAVY)
    add_paragraph(tf, "Discounts:", size=12, bold=True, color=DARK, space_before=Pt(14))
    add_paragraph(tf, "•  18% if 3 or more same-role hires in one engagement (volume).", size=11, color=DARK, space_before=Pt(6))
    add_paragraph(tf, "•  15% on a pilot engagement (first 3 hires for a new client).", size=11, color=DARK, space_before=Pt(4))
    add_paragraph(tf, "•  15% if client also subscribes to CadreHealth Connect.", size=11, color=DARK, space_before=Pt(4))
    add_paragraph(tf, "•  Floor: 15%. Never lower.", size=11, color=DARK, space_before=Pt(4))
    add_paragraph(tf, "Out-of-pocket expenses (advertising, travel, paid verification) are billed at cost with pre-approval over NGN 50,000.",
                  size=11, color=GREY, space_before=Pt(14))

    # right: deliverables
    add_rect(s, Inches(6.75), Inches(2.3), Inches(6.0), Inches(4.5), fill=NAVY)
    tf = add_textbox(s, Inches(7.05), Inches(2.5), Inches(5.4), Inches(0.4))
    set_text(tf, "WHAT YOU GET", size=11, bold=True, color=GOLD)
    deliverables = [
        "Verified shortlist of 3 to 5 candidates in 3 weeks.",
        "Credential check (MDCN/NMCN/PCN/MLSCN as applicable).",
        "Career Readiness profile and structured assessment notes.",
        "Salary band confirmed before sourcing begins.",
        "Reference checks on offer-stage candidates.",
        "3-month replacement guarantee, voluntary resignations only.",
        "Offer within 8 to 10 weeks of intake.",
        "Bill on offer acceptance, NET 14 working days.",
    ]
    tf = add_textbox(s, Inches(7.05), Inches(2.95), Inches(5.4), Inches(3.7))
    for i, d in enumerate(deliverables):
        if i == 0:
            set_text(tf, "•  " + d, size=12, color=WHITE)
        else:
            add_paragraph(tf, "•  " + d, size=12, color=WHITE, space_before=Pt(8))


def slide_tier2_detail(prs, page_num, total):
    s = new_slide(prs, "Tier 02 / Retainer", page_num, total)
    slide_title(s, "Tier 2: Retainer (Engaged Search)", kicker="FOR SCARCE OR STRATEGIC ROLES")

    add_rect(s, Inches(0.55), Inches(2.3), Inches(6.0), Inches(4.5), fill=LIGHT)
    tf = add_textbox(s, Inches(0.85), Inches(2.5), Inches(5.4), Inches(0.4))
    set_text(tf, "WHAT IT COSTS", size=11, bold=True, color=TEAL)
    tf = add_textbox(s, Inches(0.85), Inches(2.9), Inches(5.4), Inches(3.7))
    set_text(tf, "25% of first-year gross, before VAT.", size=14, bold=True, color=NAVY)
    add_paragraph(tf, "Upfront retainer: NGN 500,000.", size=13, bold=True, color=NAVY, space_before=Pt(10))
    add_paragraph(tf, "Credited 100% against the placement fee. Not refundable if the client withdraws the role.",
                  size=11, color=GREY, space_before=Pt(4))
    add_paragraph(tf, "Use this tier for:", size=12, bold=True, color=DARK, space_before=Pt(14))
    add_paragraph(tf, "•  Consultants and Heads of Department.", size=11, color=DARK, space_before=Pt(6))
    add_paragraph(tf, "•  Group CMO, COO, MD-level appointments.", size=11, color=DARK, space_before=Pt(4))
    add_paragraph(tf, "•  Niche specialties: Nephrology, Oncology, Cardiothoracic, Neonatology, Anaesthesia subspecialties.",
                  size=11, color=DARK, space_before=Pt(4))

    add_rect(s, Inches(6.75), Inches(2.3), Inches(6.0), Inches(4.5), fill=NAVY)
    tf = add_textbox(s, Inches(7.05), Inches(2.5), Inches(5.4), Inches(0.4))
    set_text(tf, "WHAT YOU GET", size=11, bold=True, color=GOLD)
    deliverables = [
        "Dedicated CFA partner accountable for the search.",
        "First shortlist of 3 named candidates in 2 weeks.",
        "Offer within 4 to 6 weeks of intake.",
        "Diaspora outreach activated where domestic supply is thin.",
        "Discreet approach to passive candidates currently employed.",
        "6-month replacement guarantee, voluntary resignations only.",
        "Quarterly market intelligence brief for your specialty.",
        "Retainer billed on kickoff, balance NET 14 from offer acceptance.",
    ]
    tf = add_textbox(s, Inches(7.05), Inches(2.95), Inches(5.4), Inches(3.7))
    for i, d in enumerate(deliverables):
        if i == 0:
            set_text(tf, "•  " + d, size=12, color=WHITE)
        else:
            add_paragraph(tf, "•  " + d, size=12, color=WHITE, space_before=Pt(8))


def slide_tier3_detail(prs, page_num, total):
    s = new_slide(prs, "Tier 03 / Connect Subscription", page_num, total)
    slide_title(s, "Tier 3: CadreHealth Connect", kicker="ANNUAL SUBSCRIPTION")

    add_rect(s, Inches(0.55), Inches(2.3), Inches(6.0), Inches(4.5), fill=LIGHT)
    tf = add_textbox(s, Inches(0.85), Inches(2.5), Inches(5.4), Inches(0.4))
    set_text(tf, "WHAT IT COSTS", size=11, bold=True, color=GOLD)
    tf = add_textbox(s, Inches(0.85), Inches(2.9), Inches(5.4), Inches(3.7))
    set_text(tf, "NGN 100,000 per year, before VAT.", size=14, bold=True, color=NAVY)
    add_paragraph(tf, "Compare:", size=12, bold=True, color=DARK, space_before=Pt(14))
    add_paragraph(tf, "•  A single Jobberman Pro-Recruit listing in 2021 was NGN 129,000, before inflation. We give you unlimited postings for less than one of theirs.",
                  size=11, color=DARK, space_before=Pt(6))
    add_paragraph(tf, "•  Plus a curated shortlist from the CadreHealth network on every posting, which Jobberman does not include.",
                  size=11, color=DARK, space_before=Pt(6))
    add_paragraph(tf, "Loss-leader pricing on purpose. The subscription seeds the relationship; per-hire engagements follow.",
                  size=11, color=GREY, space_before=Pt(14))

    add_rect(s, Inches(6.75), Inches(2.3), Inches(6.0), Inches(4.5), fill=NAVY)
    tf = add_textbox(s, Inches(7.05), Inches(2.5), Inches(5.4), Inches(0.4))
    set_text(tf, "WHAT YOU GET", size=11, bold=True, color=GOLD)
    deliverables = [
        "Unlimited job postings on the CadreHealth job board for 12 months.",
        "Verified employer profile with your hospital review surface.",
        "Curated shortlist of the top 10 candidates from our network on every posting.",
        "Quarterly market salary report for your cadres of interest.",
        "5 percentage-point discount on any subsequent Tier 1 or Tier 2 engagement.",
        "Branding placement in CadreHealth's monthly cadre newsletter.",
        "Renewable; no auto-charge.",
    ]
    tf = add_textbox(s, Inches(7.05), Inches(2.95), Inches(5.4), Inches(3.7))
    for i, d in enumerate(deliverables):
        if i == 0:
            set_text(tf, "•  " + d, size=12, color=WHITE)
        else:
            add_paragraph(tf, "•  " + d, size=12, color=WHITE, space_before=Pt(8))


def slide_how_it_works(prs, page_num, total):
    s = new_slide(prs, "How It Works", page_num, total)
    slide_title(s, "How a CadreHealth recruitment runs", kicker="THE FIVE-STEP PROCESS")

    steps = [
        ("01", "Intake", "Week 1", "JD calibration, salary band confirmation, screening rubric agreed."),
        ("02", "Source", "Week 2 to 3", "Platform-first sourcing across registered cadre members; diaspora layer activated for retainer roles."),
        ("03", "Shortlist", "Week 3 to 4", "Credential checks, structured screening, Career Readiness assessment. We present 3 to 5 named candidates."),
        ("04", "Assess", "Week 4 to 6", "Structured competency interviews. Reference and credential re-verification on offer-stage candidates."),
        ("05", "Offer & Close", "Week 6 to 10", "Offer letter support, salary negotiation backed by market data, acceptance, replacement guarantee live."),
    ]
    card_w = Inches(2.45)
    card_h = Inches(4.0)
    gap = Inches(0.1)
    start_x = Inches(0.55)
    start_y = Inches(2.3)

    for i, (num, label, when, body) in enumerate(steps):
        x = start_x + (card_w + gap) * i
        add_rect(s, x, start_y, card_w, card_h, fill=WHITE, line=SOFT_GREY)
        # number band
        add_rect(s, x, start_y, card_w, Inches(0.7), fill=NAVY)
        tf = add_textbox(s, x + Inches(0.2), start_y + Inches(0.15), card_w - Inches(0.4), Inches(0.45))
        set_text(tf, num, size=22, bold=True, color=GOLD)
        # label
        tf = add_textbox(s, x + Inches(0.2), start_y + Inches(0.85), card_w - Inches(0.4), Inches(0.4))
        set_text(tf, label, size=16, bold=True, color=NAVY)
        # when
        tf = add_textbox(s, x + Inches(0.2), start_y + Inches(1.3), card_w - Inches(0.4), Inches(0.4))
        set_text(tf, when, size=10, bold=True, color=GOLD)
        # body
        tf = add_textbox(s, x + Inches(0.2), start_y + Inches(1.75), card_w - Inches(0.4), card_h - Inches(1.85))
        set_text(tf, body, size=11, color=DARK)


def slide_quality_controls(prs, page_num, total):
    s = new_slide(prs, "Quality Controls", page_num, total)
    slide_title(s, "What makes our shortlists different", kicker="QUALITY CONTROLS")

    points = [
        ("Credential checked, not claimed",
         "Every shortlisted candidate has their MDCN, NMCN, PCN, or MLSCN status verified against the council register before they reach you. No exceptions."),
        ("Career Readiness scored, not guessed",
         "We use a cadre-specific scoring framework with separate domestic, UK, US, Canada and Gulf dimensions. You see which candidates are likely to stay 18 months versus the ones halfway to PLAB."),
        ("Hospital intel applied to matching",
         "CadreHealth carries verified hospital reviews. We know which employers your candidates have worked for, how long, and why they left. That informs cultural-fit screening before you interview."),
        ("Two-stage technical screen",
         "Specialty-specific screening competencies are signed off by a CFA clinical lead. You do not waste interview slots on candidates who cannot defend their CV."),
        ("Diaspora reach where it matters",
         "Through MANSAG, ANPA, and DFC, we surface Nigerian consultants abroad who are open to fractional, retainer, or full-return roles. Mostly for retainer-tier search."),
        ("Honest declines",
         "When a search is genuinely hard, we tell you in week two, not week ten. Hard searches get a scope reset, a salary band revision, or a recommendation to move tier."),
    ]
    cols = 2
    card_w = Inches(6.05)
    card_h = Inches(1.55)
    gap_x = Inches(0.15)
    gap_y = Inches(0.2)
    start_x = Inches(0.55)
    start_y = Inches(2.3)

    for i, (title, body) in enumerate(points):
        col = i % cols
        row = i // cols
        x = start_x + (card_w + gap_x) * col
        y = start_y + (card_h + gap_y) * row
        add_rect(s, x, y, card_w, card_h, fill=LIGHT)
        add_rect(s, x, y, Inches(0.08), card_h, fill=TEAL)
        tf = add_textbox(s, x + Inches(0.25), y + Inches(0.15), card_w - Inches(0.4), Inches(0.4))
        set_text(tf, title, size=14, bold=True, color=NAVY)
        tf = add_textbox(s, x + Inches(0.25), y + Inches(0.6), card_w - Inches(0.4), card_h - Inches(0.7))
        set_text(tf, body, size=11, color=DARK)


def slide_replacement_guarantee(prs, page_num, total):
    s = new_slide(prs, "Replacement Guarantee", page_num, total)
    slide_title(s, "Our replacement guarantee, in plain language",
                kicker="WE STAND BEHIND THE HIRE")

    tf = add_textbox(s, Inches(0.55), Inches(2.3), Inches(12.2), Inches(0.5))
    set_text(tf, "If a candidate we placed resigns voluntarily within the guarantee window, we re-recruit the role at no additional professional fee.",
             size=14, color=DARK)

    # two big cards
    card_w = Inches(6.05)
    card_h = Inches(3.7)
    start_y = Inches(3.2)

    add_rect(s, Inches(0.55), start_y, card_w, card_h, fill=LIGHT)
    add_rect(s, Inches(0.55), start_y, card_w, Inches(0.6), fill=NAVY)
    tf = add_textbox(s, Inches(0.85), start_y + Inches(0.13), card_w - Inches(0.4), Inches(0.4))
    set_text(tf, "TIER 1 / PAY-PER-HIRE", size=11, bold=True, color=GOLD)
    tf = add_textbox(s, Inches(0.85), start_y + Inches(0.85), card_w - Inches(0.4), Inches(0.5))
    set_text(tf, "3-month guarantee", size=20, bold=True, color=NAVY)
    tf = add_textbox(s, Inches(0.85), start_y + Inches(1.45), card_w - Inches(0.4), card_h - Inches(1.5))
    set_text(tf, "Counted from the candidate's start date. Free re-recruit for the same role. We commit the same SLA on the replacement search. Out-of-pocket expenses on the replacement are billed at cost.",
             size=12, color=DARK)

    add_rect(s, Inches(6.75), start_y, card_w, card_h, fill=NAVY)
    add_rect(s, Inches(6.75), start_y, card_w, Inches(0.6), fill=GOLD)
    tf = add_textbox(s, Inches(7.05), start_y + Inches(0.13), card_w - Inches(0.4), Inches(0.4))
    set_text(tf, "TIER 2 / RETAINER", size=11, bold=True, color=NAVY)
    tf = add_textbox(s, Inches(7.05), start_y + Inches(0.85), card_w - Inches(0.4), Inches(0.5))
    set_text(tf, "6-month guarantee", size=20, bold=True, color=WHITE)
    tf = add_textbox(s, Inches(7.05), start_y + Inches(1.45), card_w - Inches(0.4), card_h - Inches(1.5))
    set_text(tf, "Counted from the candidate's start date. Free re-recruit for the same role on retained terms. Dedicated partner re-engaged. Reflects industry norms for senior, scarce searches.",
             size=12, color=WHITE)


def slide_comparison(prs, page_num, total):
    s = new_slide(prs, "How We Compare", page_num, total)
    slide_title(s, "Why CadreHealth Recruitment, versus the alternatives",
                kicker="THE COMPARISON")

    # 5 columns: Dimension, CadreHealth, Traditional recruiter, Job board, LinkedIn DM
    headers = ["Dimension", "CadreHealth Recruitment", "Traditional recruiter", "Job board (Jobberman et al.)", "LinkedIn DM blast"]
    rows = [
        ["Healthcare-specific", "Yes, the whole platform", "Sometimes, as one of many", "No, generalist", "No, generalist"],
        ["Credential verified pre-shortlist", "Yes, every candidate", "Sometimes, varies", "No", "No"],
        ["Real salary data anchored", "Yes, by cadre and city", "Anecdotal", "No", "No"],
        ["Diaspora reach (UK / US / Canada)", "Yes, MANSAG / ANPA / DFC", "Rarely", "No", "Random"],
        ["Hospital review intelligence", "Yes, candidates and clients", "No", "No", "No"],
        ["Replacement guarantee", "3 to 6 months", "Usually 3 months", "None", "None"],
        ["Indicative cost per senior MO hire", "NGN 720k to 960k", "NGN 960k+", "NGN 15k to 129k per post", "Free, plus your time"],
    ]

    # table layout
    col_widths = [Inches(2.5), Inches(3.0), Inches(2.55), Inches(2.55), Inches(2.0)]
    col_xs = [Inches(0.55)]
    for w in col_widths[:-1]:
        col_xs.append(col_xs[-1] + w)

    row_h = Inches(0.5)
    y = Inches(2.3)

    # header row
    for i, h in enumerate(headers):
        fill = NAVY if i == 1 else LIGHT
        color = WHITE if i == 1 else NAVY
        add_rect(s, col_xs[i], y, col_widths[i], row_h, fill=fill)
        tf = add_textbox(s, col_xs[i] + Inches(0.1), y + Inches(0.1), col_widths[i] - Inches(0.2), row_h - Inches(0.1))
        set_text(tf, h, size=11, bold=True, color=color)

    # data rows
    for r, row in enumerate(rows):
        ry = y + row_h * (r + 1)
        for c, cell in enumerate(row):
            fill = WHITE if r % 2 == 0 else LIGHT
            if c == 1:
                fill = RGBColor(0xFC, 0xF6, 0xE3)  # very light gold tint to call out our column
            add_rect(s, col_xs[c], ry, col_widths[c], row_h, fill=fill)
            color = NAVY if c == 1 else DARK
            bold = c == 1
            tf = add_textbox(s, col_xs[c] + Inches(0.1), ry + Inches(0.08), col_widths[c] - Inches(0.2), row_h - Inches(0.1))
            set_text(tf, cell, size=10, bold=bold, color=color)


def slide_pilot_offer(prs, page_num, total):
    s = new_slide(prs, "Pilot Offer", page_num, total)
    slide_title(s, "A first-engagement pilot for new clients", kicker="LOW-RISK ENTRY")

    tf = add_textbox(s, Inches(0.55), Inches(2.3), Inches(12.2), Inches(0.6))
    set_text(tf, "We discount the first engagement so you can test the model on real roles, with real candidates, before committing to a full SLA.",
             size=14, color=DARK)

    benefits = [
        ("15% on your first 3 hires", "Down from 20% standard. Applies to Tier 1 Pay-Per-Hire only."),
        ("Free CadreHealth Connect for year 1", "NGN 100,000 subscription value. Unlimited job board postings included."),
        ("NET 30 payment terms", "Stretched from NET 14 to give your finance team comfort on a new vendor."),
        ("Quarterly partnership review", "We sit down at month 3 and month 6 to recalibrate roles, salary bands, and SLA."),
        ("Locked rate", "Pilot pricing is held for any engagement signed in 2026, even if standard rates move."),
    ]
    y = Inches(3.1)
    for i, (title, body) in enumerate(benefits):
        row_y = y + Inches(0.7) * i
        add_rect(s, Inches(0.55), row_y, Inches(12.2), Inches(0.6), fill=LIGHT if i % 2 == 0 else WHITE)
        add_rect(s, Inches(0.55), row_y, Inches(0.08), Inches(0.6), fill=GOLD)
        tf = add_textbox(s, Inches(0.8), row_y + Inches(0.13), Inches(4.5), Inches(0.4))
        set_text(tf, title, size=13, bold=True, color=NAVY)
        tf = add_textbox(s, Inches(5.4), row_y + Inches(0.13), Inches(7.2), Inches(0.4))
        set_text(tf, body, size=12, color=DARK)


def slide_commercial_terms(prs, page_num, total):
    s = new_slide(prs, "Commercial Terms", page_num, total)
    slide_title(s, "Commercial terms, on one slide", kicker="THE FINE PRINT, EARLY")

    terms = [
        ("Fees", "All quoted fees are exclusive of VAT (7.5%)."),
        ("Out-of-pocket expenses", "Verified credential pulls, paid advertising, candidate travel are billed at cost. Anything over NGN 50,000 requires pre-approval."),
        ("Payment", "NET 14 working days from offer acceptance, or NET 30 under the pilot. Late payment carries 1.5% monthly interest, not daily."),
        ("Replacement guarantee", "Free re-recruit if the candidate resigns voluntarily within 3 months (Tier 1) or 6 months (Tier 2). Termination by the client and constructive dismissal are out of scope."),
        ("Withdrawal", "If you withdraw a role mid-search for reasons other than our non-delivery, a compensation fee of 50% of the would-be placement fee applies. Industry standard, mirrors what your peers pay."),
        ("Exclusivity", "Per role, for the duration of the engagement. You agree not to run a parallel agency on the same vacancy."),
        ("Data protection", "We hold candidate data under the Nigeria Data Protection Regulation. Data is residency-aware, retention is per NDPR plus your preference."),
        ("Governing law", "Federal Republic of Nigeria."),
    ]
    y = Inches(2.3)
    row_h = Inches(0.55)
    for i, (label, body) in enumerate(terms):
        ry = y + row_h * i
        add_rect(s, Inches(0.55), ry, Inches(12.2), row_h, fill=LIGHT if i % 2 == 0 else WHITE)
        tf = add_textbox(s, Inches(0.8), ry + Inches(0.1), Inches(3.0), Inches(0.4))
        set_text(tf, label, size=12, bold=True, color=NAVY)
        tf = add_textbox(s, Inches(4.0), ry + Inches(0.1), Inches(8.5), row_h - Inches(0.1))
        set_text(tf, body, size=11, color=DARK)


def slide_onboarding_timeline(prs, page_num, total):
    s = new_slide(prs, "Onboarding Timeline", page_num, total)
    slide_title(s, "From signed SLA to first shortlist in 4 weeks",
                kicker="ONBOARDING TIMELINE")

    weeks = [
        ("WEEK 0", "Decision", "SLA signed. Roles and tier selected.", NAVY),
        ("WEEK 1", "Kickoff", "Calibration call with HR and clinical lead. Screening rubric agreed. Salary band confirmed.", TEAL),
        ("WEEK 2", "Sourcing live", "Platform-first sourcing across registered cadre members. Diaspora outreach for retainer roles.", GOLD),
        ("WEEK 3", "First shortlist", "3 to 5 named candidates with verified credentials, structured assessment notes, and salary expectations.", SUCCESS),
        ("WEEK 4 to 10", "Interviews and offer", "Structured competency interviews. Reference checks. Offer drafted and supported. Acceptance.", NAVY),
    ]
    card_w = Inches(2.45)
    card_h = Inches(4.0)
    gap = Inches(0.1)
    start_x = Inches(0.55)
    start_y = Inches(2.3)

    for i, (when, label, body, color) in enumerate(weeks):
        x = start_x + (card_w + gap) * i
        add_rect(s, x, start_y, card_w, card_h, fill=WHITE, line=SOFT_GREY)
        add_rect(s, x, start_y, card_w, Inches(0.55), fill=color)
        tf = add_textbox(s, x + Inches(0.2), start_y + Inches(0.13), card_w - Inches(0.4), Inches(0.4))
        set_text(tf, when, size=11, bold=True, color=WHITE)
        tf = add_textbox(s, x + Inches(0.2), start_y + Inches(0.75), card_w - Inches(0.4), Inches(0.5))
        set_text(tf, label, size=16, bold=True, color=NAVY)
        tf = add_textbox(s, x + Inches(0.2), start_y + Inches(1.3), card_w - Inches(0.4), card_h - Inches(1.4))
        set_text(tf, body, size=11, color=DARK)


def slide_about(prs, page_num, total):
    s = new_slide(prs, "About Us", page_num, total)
    slide_title(s, "About Consult For Africa and CadreHealth",
                kicker="CREDIBILITY")

    # left: about CFA
    add_rect(s, Inches(0.55), Inches(2.3), Inches(6.0), Inches(4.5), fill=LIGHT)
    tf = add_textbox(s, Inches(0.85), Inches(2.5), Inches(5.4), Inches(0.4))
    set_text(tf, "CONSULT FOR AFRICA", size=11, bold=True, color=GOLD)
    tf = add_textbox(s, Inches(0.85), Inches(2.9), Inches(5.4), Inches(3.7))
    set_text(tf, "Healthcare management consulting firm operating across Africa.", size=14, bold=True, color=NAVY)
    add_paragraph(tf, "Founded by Dr Debo Odulana. Engagements span hospital turnaround, clinical governance, fractional leadership, and health systems strengthening.",
                  size=12, color=DARK, space_before=Pt(12))
    add_paragraph(tf, "Active client portfolio: Cooked Indoors, Covally, Connexxum, Cureva, with discovery underway at Priscilla Specialist Hospital.",
                  size=12, color=DARK, space_before=Pt(10))
    add_paragraph(tf, "consultforafrica.com", size=11, color=GREY, space_before=Pt(14))

    # right: about CadreHealth
    add_rect(s, Inches(6.75), Inches(2.3), Inches(6.0), Inches(4.5), fill=NAVY)
    tf = add_textbox(s, Inches(7.05), Inches(2.5), Inches(5.4), Inches(0.4))
    set_text(tf, "CADREHEALTH", size=11, bold=True, color=GOLD)
    tf = add_textbox(s, Inches(7.05), Inches(2.9), Inches(5.4), Inches(3.7))
    set_text(tf, "Career platform for Nigerian healthcare professionals.", size=14, bold=True, color=WHITE)
    add_paragraph(tf, "Sixteen cadres covered. Real salary data, verified credentials, hospital reviews, diaspora layer through MANSAG, ANPA, and DFC.",
                  size=12, color=WHITE, space_before=Pt(12))
    add_paragraph(tf, "Maintained by Consult For Africa. The platform that powers our recruitment service is the same platform doctors, nurses, pharmacists, and allied health professionals use to manage their careers.",
                  size=12, color=WHITE, space_before=Pt(10))
    add_paragraph(tf, "oncadre.com", size=11, color=GOLD, space_before=Pt(14))


def slide_partnership(prs, page_num, total, lines):
    s = new_slide(prs, "Partnership Track", page_num, total)
    slide_title(s, "A separate partnership track, in parallel",
                kicker="FOR THE STRATEGIC CONVERSATION")

    tf = add_textbox(s, Inches(0.55), Inches(2.3), Inches(12.2), Inches(1.0))
    set_text(tf,
             "Recruitment is the immediate need. The partnership conversation is the bigger one. We are proposing two parallel tracks so neither slows the other down.",
             size=14, color=DARK)

    y = Inches(3.5)
    for i, line in enumerate(lines):
        row_y = y + Inches(0.6) * i
        add_rect(s, Inches(0.55), row_y, Inches(12.2), Inches(0.5), fill=LIGHT if i % 2 == 0 else WHITE)
        add_rect(s, Inches(0.55), row_y, Inches(0.08), Inches(0.5), fill=GOLD)
        tf = add_textbox(s, Inches(0.85), row_y + Inches(0.1), Inches(11.7), Inches(0.4))
        set_text(tf, line, size=12, color=DARK)


def slide_next_steps(prs, variant: Variant, page_num, total):
    s = new_slide(prs, "Next Steps", page_num, total)
    slide_title(s, "Next steps", kicker="WHAT HAPPENS NOW")

    y = Inches(2.4)
    for i, (when, what) in enumerate(variant.next_steps):
        row_y = y + Inches(0.75) * i
        add_rect(s, Inches(0.55), row_y, Inches(2.2), Inches(0.6), fill=NAVY)
        tf = add_textbox(s, Inches(0.75), row_y + Inches(0.13), Inches(2.0), Inches(0.4))
        set_text(tf, when, size=13, bold=True, color=GOLD)
        add_rect(s, Inches(2.85), row_y, Inches(9.9), Inches(0.6), fill=LIGHT)
        tf = add_textbox(s, Inches(3.1), row_y + Inches(0.13), Inches(9.5), Inches(0.4))
        set_text(tf, what, size=12, color=DARK)

    # contact
    contact_y = Inches(6.3)
    add_rect(s, Inches(0.55), contact_y, Inches(12.2), Inches(0.7), fill=NAVY)
    tf = add_textbox(s, Inches(0.85), contact_y + Inches(0.13), Inches(11.6), Inches(0.5))
    set_text(tf, "Dr Debo Odulana, Founding Partner  /  debo.odulana@consultforafrica.com  /  +234 913 813 8553",
             size=13, bold=True, color=GOLD)


def slide_appendix_sample(prs, page_num, total):
    s = new_slide(prs, "Appendix / Sample Shortlist", page_num, total)
    slide_title(s, "What a delivered shortlist looks like", kicker="APPENDIX")

    # table of 4 candidates
    headers = ["Candidate", "Cadre / Specialty", "Years", "Credential", "Readiness", "Note"]
    rows = [
        ["Candidate A", "Medical Officer", "4 yrs post-NYSC", "MDCN verified", "78 / 100", "Oncology rotation at LASUTH, EMR proficient."],
        ["Candidate B", "Medical Officer", "3 yrs post-NYSC", "MDCN verified", "71 / 100", "Internal medicine background, oncology exposure via fellowship."],
        ["Candidate C", "Dialysis Nurse", "6 yrs", "NMCN verified", "82 / 100", "Renal experience at private dialysis centre, BLS current."],
        ["Candidate D", "Consultant Nephrologist", "12 yrs", "MDCN consultant register", "85 / 100", "Available for fractional retainer; based in Lagos."],
    ]
    col_widths = [Inches(1.7), Inches(2.2), Inches(1.3), Inches(2.0), Inches(1.3), Inches(3.7)]
    col_xs = [Inches(0.55)]
    for w in col_widths[:-1]:
        col_xs.append(col_xs[-1] + w)
    row_h = Inches(0.55)
    y = Inches(2.4)

    for i, h in enumerate(headers):
        add_rect(s, col_xs[i], y, col_widths[i], row_h, fill=NAVY)
        tf = add_textbox(s, col_xs[i] + Inches(0.1), y + Inches(0.13), col_widths[i] - Inches(0.2), row_h - Inches(0.1))
        set_text(tf, h, size=10, bold=True, color=GOLD)

    for r, row in enumerate(rows):
        ry = y + row_h * (r + 1)
        for c, cell in enumerate(row):
            fill = WHITE if r % 2 == 0 else LIGHT
            add_rect(s, col_xs[c], ry, col_widths[c], row_h, fill=fill)
            tf = add_textbox(s, col_xs[c] + Inches(0.1), ry + Inches(0.13), col_widths[c] - Inches(0.2), row_h - Inches(0.1))
            set_text(tf, cell, size=10, color=DARK)

    tf = add_textbox(s, Inches(0.55), Inches(6.0), Inches(12.2), Inches(1.0))
    set_text(tf, "Each shortlisted candidate ships with: verified credential, structured assessment, Career Readiness score, salary expectations, references on offer-stage candidates.",
             size=11, color=GREY)


def slide_appendix_data_protection(prs, page_num, total):
    s = new_slide(prs, "Appendix / NDPR", page_num, total)
    slide_title(s, "Data protection, in plain language", kicker="APPENDIX")

    points = [
        ("Lawful basis", "Candidates register with CadreHealth and consent to their data being processed for recruitment matching. Hospital clients are introduced under a separate processor agreement."),
        ("Residency", "Candidate data is hosted in a Nigeria-region cloud where available, with backups in compliant jurisdictions."),
        ("Retention", "Active candidate data is retained for the duration of platform membership. Hospital-client data is retained per the signed SLA and NDPR defaults."),
        ("Disclosure", "We share candidate data with hospital clients only after explicit candidate consent for that specific role."),
        ("Subject rights", "Candidates retain full access, correction, and erasure rights under NDPR."),
        ("Breach", "Any incident is reported to the affected parties and the Nigeria Data Protection Commission within 72 hours, per regulation."),
    ]
    y = Inches(2.4)
    row_h = Inches(0.65)
    for i, (label, body) in enumerate(points):
        ry = y + row_h * i
        add_rect(s, Inches(0.55), ry, Inches(12.2), row_h, fill=LIGHT if i % 2 == 0 else WHITE)
        tf = add_textbox(s, Inches(0.8), ry + Inches(0.1), Inches(2.5), Inches(0.4))
        set_text(tf, label, size=12, bold=True, color=NAVY)
        tf = add_textbox(s, Inches(3.5), ry + Inches(0.1), Inches(9.0), row_h - Inches(0.1))
        set_text(tf, body, size=11, color=DARK)


# ---------- assemble --------------------------------------------------------

def build_deck(variant: Variant) -> Path:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # Count total pages for the chrome
    # Order: cover, client_block, problem, cost, broken, intro, platform,
    #        tiers_overview, tier1, tier2, tier3, how_it_works, quality,
    #        replacement, comparison, pilot, terms, onboarding, about,
    #        [partnership], next_steps, appendix_sample, appendix_ndpr
    total = 22 + (1 if variant.include_partnership_slide else 0)

    page = 1
    slide_cover(prs, variant)
    page += 1
    slide_client_block(prs, variant, page, total)
    page += 1
    slide_the_problem(prs, page, total)
    page += 1
    slide_cost_of_vacancy(prs, page, total)
    page += 1
    slide_whats_broken(prs, page, total)
    page += 1
    slide_intro_cadrehealth(prs, page, total)
    page += 1
    slide_platform_advantage(prs, page, total)
    page += 1
    slide_three_tiers_overview(prs, page, total)
    page += 1
    slide_tier1_detail(prs, page, total)
    page += 1
    slide_tier2_detail(prs, page, total)
    page += 1
    slide_tier3_detail(prs, page, total)
    page += 1
    slide_how_it_works(prs, page, total)
    page += 1
    slide_quality_controls(prs, page, total)
    page += 1
    slide_replacement_guarantee(prs, page, total)
    page += 1
    slide_comparison(prs, page, total)
    page += 1
    slide_pilot_offer(prs, page, total)
    page += 1
    slide_commercial_terms(prs, page, total)
    page += 1
    slide_onboarding_timeline(prs, page, total)
    page += 1
    slide_about(prs, page, total)
    page += 1
    if variant.include_partnership_slide:
        slide_partnership(prs, page, total, variant.partnership_lines)
        page += 1
    slide_next_steps(prs, variant, page, total)
    page += 1
    slide_appendix_sample(prs, page, total)
    page += 1
    slide_appendix_data_protection(prs, page, total)

    out = DOCS / f"cadrehealth-recruitment-deck-{variant.slug}.pptx" if variant.slug != "generic" \
        else DOCS / "cadrehealth-recruitment-deck.pptx"
    DOCS.mkdir(exist_ok=True)
    prs.save(str(out))
    return out


def main():
    outs = [build_deck(v) for v in (GENERIC, PEARL, OSIRIS)]
    for p in outs:
        print(f"wrote {p.relative_to(ROOT)}  ({p.stat().st_size // 1024} KB)")


if __name__ == "__main__":
    main()
