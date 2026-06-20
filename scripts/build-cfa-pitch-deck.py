"""Build CFA pitch deck introducing CadreHealth and Maarova.

Audience: Hospital leadership (CEO, MD, HR Director, Board).
Focus: Commercials - pricing tiers, anchors, ROI logic.

Constraints:
- No em dashes anywhere in copy
- No 'AI-powered' or 'AI suggestions' language
- All NGN pricing primary, USD secondary
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# Brand palette
NAVY = RGBColor(0x0B, 0x3C, 0x5D)
GOLD = RGBColor(0xD4, 0xAF, 0x37)
TEAL = RGBColor(0x1F, 0x7A, 0x8C)
LIGHT = RGBColor(0xF4, 0xF7, 0xFA)
DARK = RGBColor(0x11, 0x18, 0x27)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GREY = RGBColor(0x6B, 0x72, 0x80)
SOFT_GREY = RGBColor(0xE5, 0xE7, 0xEB)
SUCCESS = RGBColor(0x05, 0x96, 0x69)

# 16:9 widescreen
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def set_text(frame, text, size=18, bold=False, color=DARK, align=PP_ALIGN.LEFT, font="Calibri"):
    frame.clear()
    p = frame.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font


def add_paragraph(frame, text, size=14, bold=False, color=DARK, align=PP_ALIGN.LEFT, space_before=Pt(6), font="Calibri"):
    p = frame.add_paragraph()
    p.alignment = align
    p.space_before = space_before
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font


def add_rect(slide, x, y, w, h, fill=NAVY, line=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
    shape.shadow.inherit = False
    return shape


def add_textbox(slide, x, y, w, h):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.05)
    tf.margin_bottom = Inches(0.05)
    return tf


def header_band(slide, label):
    add_rect(slide, 0, 0, SLIDE_W, Inches(0.6), fill=NAVY)
    tf = add_textbox(slide, Inches(0.6), Inches(0.1), Inches(8), Inches(0.4))
    set_text(tf, "CONSULT FOR AFRICA", size=11, bold=True, color=WHITE)
    tf2 = add_textbox(slide, Inches(8.733), Inches(0.1), Inches(4.4), Inches(0.4))
    set_text(tf2, label.upper(), size=11, bold=True, color=GOLD, align=PP_ALIGN.RIGHT)


def footer_band(slide, page_num):
    add_rect(slide, 0, Inches(7.2), SLIDE_W, Inches(0.3), fill=LIGHT)
    tf = add_textbox(slide, Inches(0.6), Inches(7.22), Inches(8), Inches(0.25))
    set_text(tf, "consultforafrica.com  |  Maarova  |  CadreHealth", size=9, color=GREY)
    tf2 = add_textbox(slide, Inches(8.733), Inches(7.22), Inches(4.4), Inches(0.25))
    set_text(tf2, f"{page_num} of 12", size=9, color=GREY, align=PP_ALIGN.RIGHT)


def slide_title(slide, title, kicker=None, y=Inches(0.9)):
    if kicker:
        tf = add_textbox(slide, Inches(0.6), y, Inches(12.1), Inches(0.4))
        set_text(tf, kicker, size=12, bold=True, color=GOLD)
        y = Inches(1.25)
    tf = add_textbox(slide, Inches(0.6), y, Inches(12.1), Inches(0.8))
    set_text(tf, title, size=32, bold=True, color=NAVY, font="Calibri")
    line = add_rect(slide, Inches(0.6), Inches(2.05), Inches(0.8), Inches(0.05), fill=GOLD)


# Build the deck
prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
blank = prs.slide_layouts[6]


# ============= SLIDE 1: TITLE =============
s = prs.slides.add_slide(blank)
add_rect(s, 0, 0, SLIDE_W, SLIDE_H, fill=NAVY)
# Gold accent bar
add_rect(s, Inches(0.6), Inches(2.5), Inches(1.5), Inches(0.08), fill=GOLD)
# Eyebrow
tf = add_textbox(s, Inches(0.6), Inches(2.0), Inches(8), Inches(0.4))
set_text(tf, "CONSULT FOR AFRICA", size=14, bold=True, color=GOLD)
# Title
tf = add_textbox(s, Inches(0.6), Inches(2.8), Inches(12), Inches(2.0))
set_text(tf, "Healthcare Leadership and Workforce", size=44, bold=True, color=WHITE)
add_paragraph(tf, "A Stack Built for African Hospitals", size=44, bold=True, color=WHITE, space_before=Pt(0))
# Subtitle
tf = add_textbox(s, Inches(0.6), Inches(5.0), Inches(12), Inches(0.6))
set_text(tf, "Maarova for leadership. CadreHealth for the workforce. One operator. One mission.",
         size=18, color=GOLD)
# Footer
tf = add_textbox(s, Inches(0.6), Inches(6.7), Inches(12), Inches(0.4))
set_text(tf, "consultforafrica.com", size=12, color=WHITE)


# ============= SLIDE 2: THE REALITY =============
s = prs.slides.add_slide(blank)
header_band(s, "Slide 02")
slide_title(s, "The reality your hospital is operating in", kicker="THE PROBLEM")

# 4 stat cards
stats = [
    ("80%", "of Nigerian healthcare workers intend to emigrate", TEAL),
    ("8 to 12", "weeks lost recruiting each replacement", NAVY),
    ("N3M+", "monthly cost when a senior leader role sits vacant", GOLD),
    ("0", "objective benchmarks of hospital leadership capability", DARK),
]
card_w = Inches(2.85)
card_h = Inches(2.4)
gap = Inches(0.2)
start_x = Inches(0.6)
start_y = Inches(2.6)

for i, (stat, label, color) in enumerate(stats):
    x = start_x + (card_w + gap) * i
    add_rect(s, x, start_y, card_w, card_h, fill=LIGHT)
    add_rect(s, x, start_y, card_w, Inches(0.12), fill=color)
    tf = add_textbox(s, x + Inches(0.2), start_y + Inches(0.4), card_w - Inches(0.4), Inches(1.0))
    set_text(tf, stat, size=44, bold=True, color=color)
    tf = add_textbox(s, x + Inches(0.2), start_y + Inches(1.5), card_w - Inches(0.4), Inches(0.8))
    set_text(tf, label, size=12, color=DARK)

# Bottom takeaway
tf = add_textbox(s, Inches(0.6), Inches(5.4), Inches(12.1), Inches(1.3))
set_text(tf, "Two compounding gaps:", size=16, bold=True, color=NAVY)
add_paragraph(tf, "Talent leaves faster than you can replace, and the leaders left behind were promoted on tenure, not on capability.",
              size=15, color=DARK, space_before=Pt(8))

footer_band(s, 2)


# ============= SLIDE 3: THE COST OF DOING NOTHING =============
s = prs.slides.add_slide(blank)
header_band(s, "Slide 03")
slide_title(s, "What inaction is already costing you", kicker="COMMERCIAL REALITY")

# Cost rows
rows = [
    ("One bad senior hire", "N18M to N40M", "Recruiting, onboarding, severance, lost output. Industry benchmark: 6 to 12 months of compensation."),
    ("One unfilled consultant role for 6 months", "N12M to N20M", "Locum coverage premium plus referral leakage to competing hospitals."),
    ("One leadership failure that triggers staff exodus", "N50M+", "Replacement recruiting, training ramp, patient volume loss, reputational drag."),
    ("One year of unmeasured leadership capability", "Unbounded", "You cannot fix what you have not assessed. Every quarter compounds."),
]
y = Inches(2.4)
for i, (label, cost, detail) in enumerate(rows):
    row_y = y + Inches(0.95) * i
    bg = LIGHT if i % 2 == 0 else WHITE
    add_rect(s, Inches(0.6), row_y, Inches(12.1), Inches(0.85), fill=bg)
    # label
    tf = add_textbox(s, Inches(0.85), row_y + Inches(0.1), Inches(4.5), Inches(0.3))
    set_text(tf, label, size=13, bold=True, color=NAVY)
    # cost
    tf = add_textbox(s, Inches(0.85), row_y + Inches(0.42), Inches(4.5), Inches(0.4))
    set_text(tf, cost, size=18, bold=True, color=GOLD)
    # detail
    tf = add_textbox(s, Inches(5.6), row_y + Inches(0.18), Inches(7), Inches(0.6))
    set_text(tf, detail, size=11, color=DARK)

footer_band(s, 3)


# ============= SLIDE 4: TWO PRODUCTS =============
s = prs.slides.add_slide(blank)
header_band(s, "Slide 04")
slide_title(s, "Two products. One operator. One mission.", kicker="THE STACK")

# Two large cards
card_w = Inches(5.95)
card_h = Inches(4.4)
y = Inches(2.5)

# Maarova
add_rect(s, Inches(0.6), y, card_w, card_h, fill=NAVY)
add_rect(s, Inches(0.6), y, card_w, Inches(0.5), fill=GOLD)
tf = add_textbox(s, Inches(0.85), y + Inches(0.05), card_w - Inches(0.5), Inches(0.4))
set_text(tf, "MAAROVA", size=14, bold=True, color=NAVY)
tf = add_textbox(s, Inches(0.85), y + Inches(0.7), card_w - Inches(0.5), Inches(0.5))
set_text(tf, "Leadership Assessment and Coaching", size=20, bold=True, color=WHITE)
tf = add_textbox(s, Inches(0.85), y + Inches(1.4), card_w - Inches(0.5), Inches(2.8))
set_text(tf, "Fixes the top of the org chart.", size=13, color=GOLD, bold=True)
add_paragraph(tf, "Objective baseline of leadership capability across 6 dimensions",
              size=12, color=WHITE, space_before=Pt(10))
add_paragraph(tf, "360-degree feedback from peers, reports, and the next layer down",
              size=12, color=WHITE, space_before=Pt(6))
add_paragraph(tf, "1:1 and group coaching matched to identified gaps",
              size=12, color=WHITE, space_before=Pt(6))
add_paragraph(tf, "Executive retreats for leadership team alignment",
              size=12, color=WHITE, space_before=Pt(6))
add_paragraph(tf, "Built on Ubuntu coaching framework, not retrofitted Western tools",
              size=12, color=WHITE, space_before=Pt(6))

# CadreHealth
x2 = Inches(0.6) + card_w + Inches(0.2)
add_rect(s, x2, y, card_w, card_h, fill=TEAL)
add_rect(s, x2, y, card_w, Inches(0.5), fill=GOLD)
tf = add_textbox(s, x2 + Inches(0.25), y + Inches(0.05), card_w - Inches(0.5), Inches(0.4))
set_text(tf, "CADREHEALTH", size=14, bold=True, color=NAVY)
tf = add_textbox(s, x2 + Inches(0.25), y + Inches(0.7), card_w - Inches(0.5), Inches(0.5))
set_text(tf, "Healthcare Workforce Platform", size=20, bold=True, color=WHITE)
tf = add_textbox(s, x2 + Inches(0.25), y + Inches(1.4), card_w - Inches(0.5), Inches(2.8))
set_text(tf, "Fixes the pipeline of people moving through.", size=13, color=GOLD, bold=True)
add_paragraph(tf, "4,300+ verified healthcare professionals across 16 cadres",
              size=12, color=WHITE, space_before=Pt(10))
add_paragraph(tf, "Hospital reviews from real healthcare workers, anonymous and verified",
              size=12, color=WHITE, space_before=Pt(6))
add_paragraph(tf, "Salary intelligence by cadre, facility, and city",
              size=12, color=WHITE, space_before=Pt(6))
add_paragraph(tf, "Recruitment matching with career readiness scoring",
              size=12, color=WHITE, space_before=Pt(6))
add_paragraph(tf, "Locum marketplace and credential verification",
              size=12, color=WHITE, space_before=Pt(6))

footer_band(s, 4)


# ============= SLIDE 5: WHY BOTH TOGETHER =============
s = prs.slides.add_slide(blank)
header_band(s, "Slide 05")
slide_title(s, "Why both, together", kicker="THE INTEGRATED STORY")

# Diagram: Three columns showing the loop
col_w = Inches(3.9)
col_h = Inches(3.2)
y = Inches(2.5)
gap = Inches(0.25)

cols = [
    ("Strong Leadership", "Maarova", "Leaders who can build and retain teams. Diagnostic plus coaching plus retreats.", NAVY),
    ("Quality Workforce", "CadreHealth", "Verified talent pool, salary benchmarks, career readiness, hospital reviews.", TEAL),
    ("Better Outcomes", "Combined", "Lower attrition, faster fills, stronger governance, healthier patient outcomes.", GOLD),
]
start_x = Inches(0.6)
for i, (title, source, body, color) in enumerate(cols):
    x = start_x + (col_w + gap) * i
    add_rect(s, x, y, col_w, col_h, fill=LIGHT)
    add_rect(s, x, y, col_w, Inches(0.7), fill=color)
    tf = add_textbox(s, x + Inches(0.2), y + Inches(0.1), col_w - Inches(0.4), Inches(0.5))
    set_text(tf, title, size=16, bold=True, color=WHITE)
    tf = add_textbox(s, x + Inches(0.2), y + Inches(0.85), col_w - Inches(0.4), Inches(0.4))
    set_text(tf, source.upper(), size=10, bold=True, color=color)
    tf = add_textbox(s, x + Inches(0.2), y + Inches(1.3), col_w - Inches(0.4), Inches(1.8))
    set_text(tf, body, size=13, color=DARK)

# Bottom punchline
tf = add_textbox(s, Inches(0.6), Inches(6.0), Inches(12.1), Inches(1.0))
set_text(tf, "Maarova fixes the top. CadreHealth fixes the pipeline.",
         size=18, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
add_paragraph(tf, "Leaders who know how to retain talent, plus a marketplace where they can find it.",
              size=14, color=DARK, align=PP_ALIGN.CENTER, space_before=Pt(8))

footer_band(s, 5)


# ============= SLIDE 6: MAAROVA HOW IT WORKS =============
s = prs.slides.add_slide(blank)
header_band(s, "Slide 06")
slide_title(s, "Maarova: how the engagement works", kicker="LEADERSHIP PRODUCT")

# 4-step process
steps = [
    ("01", "Assess", "Online assessment plus 360 feedback. 90 minutes total commitment per leader."),
    ("02", "Debrief", "1:1 results review with accredited coach. Top 2 to 3 development goals identified."),
    ("03", "Coach", "12 sessions over 6 months, biweekly. 1:1 or group format. WhatsApp accountability."),
    ("04", "Sustain", "Mini 360 at month 6. Personal development plan. Optional cohort retreat."),
]
y = Inches(2.5)
step_w = Inches(2.95)
gap = Inches(0.15)
start_x = Inches(0.6)

for i, (num, title, body) in enumerate(steps):
    x = start_x + (step_w + gap) * i
    add_rect(s, x, y, step_w, Inches(3.6), fill=LIGHT)
    add_rect(s, x, y, step_w, Inches(0.9), fill=NAVY)
    tf = add_textbox(s, x + Inches(0.25), y + Inches(0.1), step_w - Inches(0.5), Inches(0.4))
    set_text(tf, num, size=14, bold=True, color=GOLD)
    tf = add_textbox(s, x + Inches(0.25), y + Inches(0.4), step_w - Inches(0.5), Inches(0.5))
    set_text(tf, title, size=20, bold=True, color=WHITE)
    tf = add_textbox(s, x + Inches(0.25), y + Inches(1.1), step_w - Inches(0.5), Inches(2.4))
    set_text(tf, body, size=12, color=DARK)

# Footer note
tf = add_textbox(s, Inches(0.6), Inches(6.4), Inches(12.1), Inches(0.5))
set_text(tf, "Outcome: a leader with a clear capability map, a coach, and a development plan tied to your hospital strategy.",
         size=13, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

footer_band(s, 6)


# ============= SLIDE 7: COMMERCIALS - MAAROVA PRICING =============
s = prs.slides.add_slide(blank)
header_band(s, "Slide 07")
slide_title(s, "Maarova investment tiers", kicker="COMMERCIALS  /  MAAROVA")

# Pricing table
tiers = [
    ("Assessment Only", "N400K to N600K", "Per leader. Full report and 1:1 debrief.", False),
    ("Assessment + Group Coaching", "N800K to N1.2M", "6 monthly group sessions over 6 months.", False),
    ("Assessment + 1:1 Coaching", "N2M to N3.5M", "12 sessions over 6 months. Recommended starting point for senior leaders.", True),
    ("Leadership Intensive", "N4M to N6M", "1:1 plus group plus on-site workshop. For C-suite.", False),
    ("Executive Retreat", "N1.5M to N2.5M", "Per person. 2 to 3 day offsite. Premium venue, peer cohort.", False),
    ("Enterprise Cohort", "N15M to N25M", "10+ leaders. Stack of assessment, coaching, and retreat. Best per-head value.", True),
]
y = Inches(2.4)
row_h = Inches(0.65)
for i, (name, price, body, featured) in enumerate(tiers):
    row_y = y + row_h * i
    bg = GOLD if featured else (LIGHT if i % 2 == 0 else WHITE)
    text_color = NAVY if featured else DARK
    add_rect(s, Inches(0.6), row_y, Inches(12.1), row_h, fill=bg)
    tf = add_textbox(s, Inches(0.85), row_y + Inches(0.1), Inches(4.0), Inches(0.45))
    set_text(tf, name, size=14, bold=True, color=text_color)
    tf = add_textbox(s, Inches(4.9), row_y + Inches(0.1), Inches(2.8), Inches(0.45))
    set_text(tf, price, size=15, bold=True, color=NAVY if featured else GOLD)
    tf = add_textbox(s, Inches(7.8), row_y + Inches(0.13), Inches(4.8), Inches(0.5))
    set_text(tf, body, size=11, color=text_color)

# Anchor footer
tf = add_textbox(s, Inches(0.6), Inches(6.5), Inches(12.1), Inches(0.5))
set_text(tf, "ANCHOR: LBS Advanced Management Programme is N17 to 18M for ONE leader. Our enterprise cohort sends 10+ leaders for the same outlay.",
         size=11, bold=True, color=NAVY, align=PP_ALIGN.LEFT)

footer_band(s, 7)


# ============= SLIDE 8: COMMERCIALS - CADREHEALTH =============
s = prs.slides.add_slide(blank)
header_band(s, "Slide 08")
slide_title(s, "CadreHealth for hospitals", kicker="COMMERCIALS  /  CADREHEALTH")

# Two columns: Subscription tiers + Performance fees
left_x = Inches(0.6)
right_x = Inches(6.85)
col_w = Inches(6.0)
y = Inches(2.4)

# Left: Subscriptions
add_rect(s, left_x, y, col_w, Inches(0.6), fill=NAVY)
tf = add_textbox(s, left_x + Inches(0.25), y + Inches(0.12), col_w - Inches(0.5), Inches(0.4))
set_text(tf, "EMPLOYER SUBSCRIPTIONS", size=14, bold=True, color=WHITE)

sub_tiers = [
    ("Verified Profile", "N250K / year", "Verified employer badge, basic facility profile, response to reviews."),
    ("Talent Search", "N1.5M / year", "Profile plus search across 4,300+ verified professionals, 5 mandates per year."),
    ("Workforce Intelligence", "N4M / year", "Talent search plus salary benchmarks, attrition data, anonymous hospital review insights."),
]
for i, (name, price, body) in enumerate(sub_tiers):
    rh = Inches(1.0)
    ry = y + Inches(0.6) + rh * i + Inches(0.05) * i
    bg = LIGHT if i % 2 == 0 else WHITE
    add_rect(s, left_x, ry, col_w, rh, fill=bg)
    tf = add_textbox(s, left_x + Inches(0.25), ry + Inches(0.1), col_w - Inches(0.5), Inches(0.35))
    set_text(tf, name, size=13, bold=True, color=NAVY)
    tf = add_textbox(s, left_x + Inches(0.25), ry + Inches(0.4), col_w - Inches(0.5), Inches(0.3))
    set_text(tf, price, size=14, bold=True, color=GOLD)
    tf = add_textbox(s, left_x + Inches(0.25), ry + Inches(0.7), col_w - Inches(0.5), Inches(0.3))
    set_text(tf, body, size=10, color=DARK)

# Right: Placement + Locum
add_rect(s, right_x, y, col_w, Inches(0.6), fill=TEAL)
tf = add_textbox(s, right_x + Inches(0.25), y + Inches(0.12), col_w - Inches(0.5), Inches(0.4))
set_text(tf, "PERFORMANCE FEES", size=14, bold=True, color=WHITE)

perf = [
    ("Permanent Placement", "20 to 25%", "Of first-year compensation. Pay only on hire. 90-day replacement guarantee."),
    ("Locum Placement", "10 to 15%", "Margin on locum shifts. Instant payment to professionals via Paystack."),
    ("International Recruitment", "Custom", "Diaspora and returning professionals. Negotiated per mandate."),
]
for i, (name, price, body) in enumerate(perf):
    rh = Inches(1.0)
    ry = y + Inches(0.6) + rh * i + Inches(0.05) * i
    bg = LIGHT if i % 2 == 0 else WHITE
    add_rect(s, right_x, ry, col_w, rh, fill=bg)
    tf = add_textbox(s, right_x + Inches(0.25), ry + Inches(0.1), col_w - Inches(0.5), Inches(0.35))
    set_text(tf, name, size=13, bold=True, color=NAVY)
    tf = add_textbox(s, right_x + Inches(0.25), ry + Inches(0.4), col_w - Inches(0.5), Inches(0.3))
    set_text(tf, price, size=14, bold=True, color=GOLD)
    tf = add_textbox(s, right_x + Inches(0.25), ry + Inches(0.7), col_w - Inches(0.5), Inches(0.3))
    set_text(tf, body, size=10, color=DARK)

# Anchor
tf = add_textbox(s, Inches(0.6), Inches(6.5), Inches(12.1), Inches(0.5))
set_text(tf, "ANCHOR: Traditional headhunters charge 25 to 35% with no replacement guarantee. Our subscription plus placement model is structurally cheaper at any hire volume above 2 per year.",
         size=11, bold=True, color=NAVY)

footer_band(s, 8)


# ============= SLIDE 9: COMMERCIALS - BUNDLED STACK =============
s = prs.slides.add_slide(blank)
header_band(s, "Slide 09")
slide_title(s, "The integrated stack: best per-naira value", kicker="COMMERCIALS  /  BUNDLE")

# Recommendation cards: Small / Medium / Large hospital
profiles = [
    {
        "label": "UNDER 50 BEDS",
        "name": "Boutique Stack",
        "price": "N6M to N10M",
        "duration": "First year",
        "items": [
            "Maarova for 2 leaders (assessment + 1:1)",
            "CadreHealth Verified Profile",
            "2 placements credit",
            "Annual leadership check-in",
        ],
    },
    {
        "label": "50 TO 150 BEDS",
        "name": "Foundation Stack",
        "price": "N15M to N25M",
        "duration": "First year",
        "items": [
            "Maarova for 5 leaders (assessment + 1:1)",
            "CadreHealth Talent Search subscription",
            "5 placements credit",
            "Quarterly leadership review",
        ],
    },
    {
        "label": "150 TO 350 BEDS",
        "name": "Growth Stack",
        "price": "N35M to N55M",
        "duration": "First year",
        "items": [
            "Maarova Enterprise Cohort of 10 leaders",
            "Executive retreat included",
            "CadreHealth Workforce Intelligence",
            "12 placements credit",
            "Half-yearly board reporting",
        ],
        "featured": True,
    },
    {
        "label": "350+ BEDS / MULTI-SITE",
        "name": "Enterprise Stack",
        "price": "From N75M",
        "duration": "First year",
        "items": [
            "Maarova for 20+ leaders, multi-site",
            "Two retreats per year",
            "CadreHealth multi-facility deployment",
            "25+ placements credit",
            "Dedicated CFA engagement manager",
        ],
    },
]
y = Inches(2.4)
card_w = Inches(2.95)
gap = Inches(0.15)
start_x = Inches(0.6)

for i, p in enumerate(profiles):
    x = start_x + (card_w + gap) * i
    featured = p.get("featured", False)
    fill = NAVY if featured else LIGHT
    text_main = WHITE if featured else DARK
    label_color = GOLD if featured else NAVY
    add_rect(s, x, y, card_w, Inches(4.4), fill=fill)
    if featured:
        # Gold ribbon
        add_rect(s, x, y, card_w, Inches(0.3), fill=GOLD)
        tf = add_textbox(s, x + Inches(0.15), y + Inches(0.02), card_w - Inches(0.3), Inches(0.25))
        set_text(tf, "RECOMMENDED", size=9, bold=True, color=NAVY)
        ty = y + Inches(0.4)
    else:
        ty = y + Inches(0.2)

    tf = add_textbox(s, x + Inches(0.15), ty, card_w - Inches(0.3), Inches(0.3))
    set_text(tf, p["label"], size=9, bold=True, color=label_color)
    tf = add_textbox(s, x + Inches(0.15), ty + Inches(0.3), card_w - Inches(0.3), Inches(0.5))
    set_text(tf, p["name"], size=17, bold=True, color=text_main)
    tf = add_textbox(s, x + Inches(0.15), ty + Inches(0.95), card_w - Inches(0.3), Inches(0.5))
    set_text(tf, p["price"], size=18, bold=True, color=GOLD)
    tf = add_textbox(s, x + Inches(0.15), ty + Inches(1.45), card_w - Inches(0.3), Inches(0.3))
    set_text(tf, p["duration"], size=10, color=text_main)
    # Items
    item_tf = add_textbox(s, x + Inches(0.15), ty + Inches(1.85), card_w - Inches(0.3), Inches(2.3))
    set_text(item_tf, "• " + p["items"][0], size=10, color=text_main)
    for item in p["items"][1:]:
        add_paragraph(item_tf, "• " + item, size=10, color=text_main, space_before=Pt(4))

# Bottom note
tf = add_textbox(s, Inches(0.6), Inches(6.95), Inches(12.1), Inches(0.3))
set_text(tf, "Placements above the credit pool charged at standard 20 to 25% of first-year compensation. Bundle savings of 15 to 25% versus à la carte.",
         size=10, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

footer_band(s, 9)


# ============= SLIDE 10: ROI LOGIC =============
s = prs.slides.add_slide(blank)
header_band(s, "Slide 10")
slide_title(s, "The ROI math", kicker="COMMERCIALS  /  PAYBACK")

# Three ROI scenarios
scenarios = [
    {
        "title": "Avoid one bad senior hire",
        "save": "N20M+",
        "vs": "Foundation Stack: N12 to 18M",
        "verdict": "Pays for itself with one prevented mis-hire.",
    },
    {
        "title": "Reduce attrition by 5%",
        "save": "N40M+",
        "vs": "200-bed hospital, 600 staff at avg N1M/yr churn cost",
        "verdict": "Growth Stack pays back in under 12 months.",
    },
    {
        "title": "One leader who scales 3 more",
        "save": "Compounding",
        "vs": "Cost of NOT developing top leadership",
        "verdict": "The compounding return that no benchmark fully captures.",
    },
]
y = Inches(2.4)
card_w = Inches(3.95)
gap = Inches(0.2)
start_x = Inches(0.6)

for i, sc in enumerate(scenarios):
    x = start_x + (card_w + gap) * i
    add_rect(s, x, y, card_w, Inches(3.8), fill=LIGHT)
    add_rect(s, x, y, card_w, Inches(0.08), fill=GOLD)
    tf = add_textbox(s, x + Inches(0.25), y + Inches(0.3), card_w - Inches(0.5), Inches(0.6))
    set_text(tf, sc["title"], size=15, bold=True, color=NAVY)
    tf = add_textbox(s, x + Inches(0.25), y + Inches(1.1), card_w - Inches(0.5), Inches(0.7))
    set_text(tf, "You save", size=11, color=GREY)
    add_paragraph(tf, sc["save"], size=32, bold=True, color=GOLD, space_before=Pt(2))
    tf = add_textbox(s, x + Inches(0.25), y + Inches(2.4), card_w - Inches(0.5), Inches(0.6))
    set_text(tf, sc["vs"], size=10, color=DARK)
    tf = add_textbox(s, x + Inches(0.25), y + Inches(3.05), card_w - Inches(0.5), Inches(0.6))
    set_text(tf, sc["verdict"], size=11, bold=True, color=TEAL)

# Bottom punchline
tf = add_textbox(s, Inches(0.6), Inches(6.5), Inches(12.1), Inches(0.5))
set_text(tf, "Healthcare leadership and workforce decisions are not a cost line. They are the highest-leverage investment a hospital makes.",
         size=13, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

footer_band(s, 10)


# ============= SLIDE 11: HOW WE ENGAGE =============
s = prs.slides.add_slide(blank)
header_band(s, "Slide 11")
slide_title(s, "How we start", kicker="ENGAGEMENT PATH")

# 4-step horizontal timeline
steps = [
    ("Week 0", "Discovery Call", "30-minute call. We understand your context. No deck on our end, just questions."),
    ("Week 1 to 2", "Diagnostic", "Light-touch leadership and workforce diagnostic. Combination of interviews and platform data."),
    ("Week 3", "Tailored Proposal", "Stack recommendation with named coaches, scope, timeline, and fixed pricing in NGN."),
    ("Week 4+", "Phased Rollout", "Pilot with one cohort or facility. Scale based on signal. No multi-year lock-in."),
]
y = Inches(2.6)
card_w = Inches(2.95)
gap = Inches(0.15)
start_x = Inches(0.6)

for i, (when, title, body) in enumerate(steps):
    x = start_x + (card_w + gap) * i
    add_rect(s, x, y, card_w, Inches(3.6), fill=WHITE, line=NAVY)
    add_rect(s, x, y, card_w, Inches(0.6), fill=NAVY)
    tf = add_textbox(s, x + Inches(0.25), y + Inches(0.13), card_w - Inches(0.5), Inches(0.4))
    set_text(tf, when.upper(), size=11, bold=True, color=GOLD)
    tf = add_textbox(s, x + Inches(0.25), y + Inches(0.85), card_w - Inches(0.5), Inches(0.6))
    set_text(tf, title, size=18, bold=True, color=NAVY)
    tf = add_textbox(s, x + Inches(0.25), y + Inches(1.6), card_w - Inches(0.5), Inches(1.9))
    set_text(tf, body, size=11, color=DARK)
    # Step number badge
    badge = s.shapes.add_shape(MSO_SHAPE.OVAL, x + card_w - Inches(0.65), y - Inches(0.25), Inches(0.5), Inches(0.5))
    badge.fill.solid()
    badge.fill.fore_color.rgb = GOLD
    badge.line.fill.background()
    badge_tf = badge.text_frame
    badge_tf.margin_left = Inches(0)
    badge_tf.margin_right = Inches(0)
    set_text(badge_tf, str(i + 1), size=18, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

# Bottom emphasis
tf = add_textbox(s, Inches(0.6), Inches(6.5), Inches(12.1), Inches(0.5))
set_text(tf, "No retainer. No multi-year commitment. You buy outcomes, not seat licenses.",
         size=14, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

footer_band(s, 11)


# ============= SLIDE 12: NEXT STEP / CONTACT =============
s = prs.slides.add_slide(blank)
add_rect(s, 0, 0, SLIDE_W, SLIDE_H, fill=NAVY)

# Big CTA
tf = add_textbox(s, Inches(0.6), Inches(1.5), Inches(12), Inches(0.5))
set_text(tf, "NEXT STEP", size=14, bold=True, color=GOLD)
tf = add_textbox(s, Inches(0.6), Inches(2.0), Inches(12), Inches(2.0))
set_text(tf, "Book a 30-minute discovery call.", size=42, bold=True, color=WHITE)
add_paragraph(tf, "We will map your specific gaps and send a tailored proposal within 7 days.",
              size=22, color=GOLD, space_before=Pt(15))

# Contact block
add_rect(s, Inches(0.6), Inches(5.0), Inches(12.1), Inches(1.6), fill=WHITE)
tf = add_textbox(s, Inches(0.85), Inches(5.15), Inches(11.6), Inches(0.4))
set_text(tf, "CONTACT", size=11, bold=True, color=GOLD)
tf = add_textbox(s, Inches(0.85), Inches(5.5), Inches(5.5), Inches(1.0))
set_text(tf, "Consult For Africa", size=18, bold=True, color=NAVY)
add_paragraph(tf, "consultforafrica.com", size=13, color=DARK, space_before=Pt(5))
add_paragraph(tf, "hello@consultforafrica.com", size=13, color=DARK, space_before=Pt(3))

tf = add_textbox(s, Inches(7.0), Inches(5.5), Inches(5.5), Inches(1.0))
set_text(tf, "Maarova", size=14, bold=True, color=NAVY)
add_paragraph(tf, "consultforafrica.com/maarova", size=12, color=DARK, space_before=Pt(3))
add_paragraph(tf, "CadreHealth", size=14, bold=True, color=NAVY, space_before=Pt(8))
add_paragraph(tf, "consultforafrica.com/oncadre", size=12, color=DARK, space_before=Pt(3))

# Footer
tf = add_textbox(s, Inches(0.6), Inches(6.9), Inches(12), Inches(0.4))
set_text(tf, "African healthcare that runs better.", size=14, color=GOLD, align=PP_ALIGN.CENTER)


# Save
output_path = "/Users/debo/consult-for-africa/docs/cfa-pitch-deck.pptx"
prs.save(output_path)
print(f"Deck saved: {output_path}")
print(f"Slides: {len(prs.slides)}")
