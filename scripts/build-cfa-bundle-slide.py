"""Build standalone slide 9 (the integrated stack bundle pricing).

For drop-in merging into the main deck. 4 MECE tiers covering the
full Nigerian hospital bed-count spectrum.
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

NAVY = RGBColor(0x0B, 0x3C, 0x5D)
GOLD = RGBColor(0xD4, 0xAF, 0x37)
TEAL = RGBColor(0x1F, 0x7A, 0x8C)
LIGHT = RGBColor(0xF4, 0xF7, 0xFA)
DARK = RGBColor(0x11, 0x18, 0x27)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GREY = RGBColor(0x6B, 0x72, 0x80)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def set_text(frame, text, size=18, bold=False, color=DARK, align=PP_ALIGN.LEFT, font="Calibri"):
    frame.clear()
    p = frame.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font


def add_paragraph(frame, text, size=14, bold=False, color=DARK, align=PP_ALIGN.LEFT, space_before=Pt(6), font="Calibri"):
    p = frame.add_paragraph()
    p.alignment = align
    p.space_before = space_before
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font


def add_rect(slide, x, y, w, h, fill=NAVY, line=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
    shape.shadow.inherit = False
    return shape


def add_textbox(slide, x, y, w, h):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.05)
    tf.margin_bottom = Inches(0.05)
    return tf


prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
blank = prs.slide_layouts[6]
s = prs.slides.add_slide(blank)

# Header band
add_rect(s, 0, 0, SLIDE_W, Inches(0.6), fill=NAVY)
tf = add_textbox(s, Inches(0.6), Inches(0.1), Inches(8), Inches(0.4))
set_text(tf, "CONSULT FOR AFRICA", size=11, bold=True, color=WHITE)
tf = add_textbox(s, Inches(8.733), Inches(0.1), Inches(4.4), Inches(0.4))
set_text(tf, "INSERT SLIDE", size=11, bold=True, color=GOLD, align=PP_ALIGN.RIGHT)

# Title block
tf = add_textbox(s, Inches(0.6), Inches(0.9), Inches(12.1), Inches(0.4))
set_text(tf, "COMMERCIALS  /  BUNDLE", size=12, bold=True, color=GOLD)
tf = add_textbox(s, Inches(0.6), Inches(1.25), Inches(12.1), Inches(0.8))
set_text(tf, "The integrated stack: best per-naira value", size=32, bold=True, color=NAVY)
add_rect(s, Inches(0.6), Inches(2.05), Inches(0.8), Inches(0.05), fill=GOLD)

# 4 MECE pricing tiers
profiles = [
    {
        "label": "UNDER 50 BEDS",
        "name": "Boutique Stack",
        "price": "N6M to N10M",
        "duration": "First year",
        "items": [
            "Maarova for 2 leaders (assessment + 1:1)",
            "CadreHealth Verified Profile",
            "2 placements credit",
            "Annual leadership check-in",
        ],
    },
    {
        "label": "50 TO 150 BEDS",
        "name": "Foundation Stack",
        "price": "N15M to N25M",
        "duration": "First year",
        "items": [
            "Maarova for 5 leaders (assessment + 1:1)",
            "CadreHealth Talent Search subscription",
            "5 placements credit",
            "Quarterly leadership review",
        ],
    },
    {
        "label": "150 TO 350 BEDS",
        "name": "Growth Stack",
        "price": "N35M to N55M",
        "duration": "First year",
        "items": [
            "Maarova Enterprise Cohort of 10 leaders",
            "Executive retreat included",
            "CadreHealth Workforce Intelligence",
            "12 placements credit",
            "Half-yearly board reporting",
        ],
        "featured": True,
    },
    {
        "label": "350+ BEDS / MULTI-SITE",
        "name": "Enterprise Stack",
        "price": "From N75M",
        "duration": "First year",
        "items": [
            "Maarova for 20+ leaders, multi-site",
            "Two retreats per year",
            "CadreHealth multi-facility deployment",
            "25+ placements credit",
            "Dedicated CFA engagement manager",
        ],
    },
]
y = Inches(2.4)
card_w = Inches(2.95)
gap = Inches(0.15)
start_x = Inches(0.6)

for i, p in enumerate(profiles):
    x = start_x + (card_w + gap) * i
    featured = p.get("featured", False)
    fill = NAVY if featured else LIGHT
    text_main = WHITE if featured else DARK
    label_color = GOLD if featured else NAVY
    add_rect(s, x, y, card_w, Inches(4.4), fill=fill)
    if featured:
        add_rect(s, x, y, card_w, Inches(0.3), fill=GOLD)
        tf = add_textbox(s, x + Inches(0.15), y + Inches(0.02), card_w - Inches(0.3), Inches(0.25))
        set_text(tf, "RECOMMENDED", size=9, bold=True, color=NAVY)
        ty = y + Inches(0.4)
    else:
        ty = y + Inches(0.2)

    tf = add_textbox(s, x + Inches(0.15), ty, card_w - Inches(0.3), Inches(0.3))
    set_text(tf, p["label"], size=9, bold=True, color=label_color)
    tf = add_textbox(s, x + Inches(0.15), ty + Inches(0.3), card_w - Inches(0.3), Inches(0.5))
    set_text(tf, p["name"], size=17, bold=True, color=text_main)
    tf = add_textbox(s, x + Inches(0.15), ty + Inches(0.95), card_w - Inches(0.3), Inches(0.5))
    set_text(tf, p["price"], size=18, bold=True, color=GOLD)
    tf = add_textbox(s, x + Inches(0.15), ty + Inches(1.45), card_w - Inches(0.3), Inches(0.3))
    set_text(tf, p["duration"], size=10, color=text_main)
    item_tf = add_textbox(s, x + Inches(0.15), ty + Inches(1.85), card_w - Inches(0.3), Inches(2.3))
    set_text(item_tf, "• " + p["items"][0], size=10, color=text_main)
    for item in p["items"][1:]:
        add_paragraph(item_tf, "• " + item, size=10, color=text_main, space_before=Pt(4))

# Bottom note
tf = add_textbox(s, Inches(0.6), Inches(6.95), Inches(12.1), Inches(0.3))
set_text(tf, "Placements above the credit pool charged at standard 20 to 25% of first-year compensation. Bundle savings of 15 to 25% versus à la carte.",
         size=10, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

# Footer
add_rect(s, 0, Inches(7.2), SLIDE_W, Inches(0.3), fill=LIGHT)
tf = add_textbox(s, Inches(0.6), Inches(7.22), Inches(8), Inches(0.25))
set_text(tf, "consultforafrica.com  |  Maarova  |  CadreHealth", size=9, color=GREY)

output_path = "/Users/debo/consult-for-africa/docs/cfa-bundle-pricing-slide.pptx"
prs.save(output_path)
print(f"Standalone slide saved: {output_path}")
