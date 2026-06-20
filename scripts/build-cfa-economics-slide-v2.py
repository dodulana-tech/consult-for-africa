"""V2: economics slide with corrected MO locum rate (N35K not N100K).

Other figures held constant. The senior fractional and mid-level secondment
columns were market-realistic; only the locum example needed correction.
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

NAVY = RGBColor(0x0B, 0x3C, 0x5D)
GOLD = RGBColor(0xD4, 0xAF, 0x37)
TEAL = RGBColor(0x1F, 0x7A, 0x8C)
LIGHT = RGBColor(0xF4, 0xF7, 0xFA)
DARK = RGBColor(0x11, 0x18, 0x27)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GREY = RGBColor(0x6B, 0x72, 0x80)
SOFT = RGBColor(0xE5, 0xE7, 0xEB)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def set_text(frame, text, size=18, bold=False, color=DARK, align=PP_ALIGN.LEFT, font="Calibri"):
    frame.clear()
    p = frame.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font


def add_paragraph(frame, text, size=14, bold=False, color=DARK, align=PP_ALIGN.LEFT, space_before=Pt(6), font="Calibri"):
    p = frame.add_paragraph()
    p.alignment = align
    p.space_before = space_before
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font


def add_rect(slide, x, y, w, h, fill=NAVY, line=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
    shape.shadow.inherit = False
    return shape


def add_textbox(slide, x, y, w, h):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.05)
    tf.margin_bottom = Inches(0.05)
    return tf


def waterfall_row(slide, x, y, w, label, amount, kind="neutral"):
    color_map = {"gross": NAVY, "cfa": GOLD, "pro": TEAL, "net": NAVY, "deduction": GREY, "neutral": GREY}
    text_color_map = {"gross": WHITE, "cfa": NAVY, "pro": WHITE, "net": WHITE, "deduction": DARK, "neutral": DARK}
    bg = color_map.get(kind, GREY)
    fg = text_color_map.get(kind, DARK)
    h = Inches(0.34)
    add_rect(slide, x, y, w, h, fill=bg)
    tf = add_textbox(slide, x + Inches(0.12), y + Inches(0.05), w - Inches(2.2), Inches(0.25))
    set_text(tf, label, size=10, bold=True, color=fg)
    tf = add_textbox(slide, x + w - Inches(2.05), y + Inches(0.05), Inches(1.9), Inches(0.25))
    set_text(tf, amount, size=10, bold=True, color=fg, align=PP_ALIGN.RIGHT)


def light_row(slide, x, y, w, label, amount):
    h = Inches(0.3)
    add_rect(slide, x, y, w, h, fill=LIGHT)
    tf = add_textbox(slide, x + Inches(0.12), y + Inches(0.04), w - Inches(2.0), Inches(0.22))
    set_text(tf, label, size=9, color=DARK)
    tf = add_textbox(slide, x + w - Inches(1.85), y + Inches(0.04), Inches(1.7), Inches(0.22))
    set_text(tf, amount, size=9, bold=True, color=DARK, align=PP_ALIGN.RIGHT)


prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
blank = prs.slide_layouts[6]
s = prs.slides.add_slide(blank)

add_rect(s, 0, 0, SLIDE_W, Inches(0.6), fill=NAVY)
tf = add_textbox(s, Inches(0.6), Inches(0.1), Inches(8), Inches(0.4))
set_text(tf, "CONSULT FOR AFRICA", size=11, bold=True, color=WHITE)
tf = add_textbox(s, Inches(8.733), Inches(0.1), Inches(4.4), Inches(0.4))
set_text(tf, "INSERT SLIDE", size=11, bold=True, color=GOLD, align=PP_ALIGN.RIGHT)

tf = add_textbox(s, Inches(0.6), Inches(0.85), Inches(12.1), Inches(0.4))
set_text(tf, "COMMERCIALS  /  ECONOMICS AND TRANSPARENCY", size=11, bold=True, color=GOLD)
tf = add_textbox(s, Inches(0.6), Inches(1.2), Inches(12.1), Inches(0.7))
set_text(tf, "Where every naira goes: hospital, taxes, CFA, professional", size=24, bold=True, color=NAVY)
add_rect(s, Inches(0.6), Inches(1.95), Inches(0.8), Inches(0.05), fill=GOLD)

card_w = Inches(4.1)
gap = Inches(0.15)
start_x = Inches(0.6)
cards_y = Inches(2.2)
card_h = Inches(4.65)

# ---------- CARD 1: Senior Fractional (unchanged) ----------
x = start_x
add_rect(s, x, cards_y, card_w, card_h, fill=WHITE, line=SOFT)
add_rect(s, x, cards_y, card_w, Inches(0.55), fill=NAVY)
tf = add_textbox(s, x + Inches(0.15), cards_y + Inches(0.05), card_w - Inches(0.3), Inches(0.25))
set_text(tf, "SENIOR FRACTIONAL", size=9, bold=True, color=GOLD)
tf = add_textbox(s, x + Inches(0.15), cards_y + Inches(0.27), card_w - Inches(0.3), Inches(0.3))
set_text(tf, "Fractional CMO at N5M / month", size=12, bold=True, color=WHITE)

ry = cards_y + Inches(0.7)
waterfall_row(s, x + Inches(0.15), ry, card_w - Inches(0.3), "Hospital invoice (ex VAT)", "N5,000,000", "gross")
ry += Inches(0.38)
light_row(s, x + Inches(0.15), ry, card_w - Inches(0.3), "+ VAT 7.5% (input credit if registered)", "N375,000")
ry += Inches(0.32)
light_row(s, x + Inches(0.15), ry, card_w - Inches(0.3), "- WHT 10% (remitted to FIRS)", "N500,000")
ry += Inches(0.4)
waterfall_row(s, x + Inches(0.15), ry, card_w - Inches(0.3), "CFA platform fee (25%)", "N1,250,000", "cfa")
ry += Inches(0.42)
waterfall_row(s, x + Inches(0.15), ry, card_w - Inches(0.3), "Professional gross (contractor)", "N3,750,000", "pro")
ry += Inches(0.38)
light_row(s, x + Inches(0.15), ry, card_w - Inches(0.3), "- Professional own PIT (~25%)", "N940,000")
ry += Inches(0.4)
waterfall_row(s, x + Inches(0.15), ry, card_w - Inches(0.3), "Professional net cash", "N2,810,000", "net")
ry += Inches(0.5)
tf = add_textbox(s, x + Inches(0.15), ry, card_w - Inches(0.3), Inches(0.4))
set_text(tf, "Senior contractor model. Professional invoices CFA, handles own taxes.", size=8, color=GREY)

# ---------- CARD 2: Mid-Level Secondment (unchanged) ----------
x = start_x + (card_w + gap)
add_rect(s, x, cards_y, card_w, card_h, fill=WHITE, line=SOFT)
add_rect(s, x, cards_y, card_w, Inches(0.55), fill=TEAL)
tf = add_textbox(s, x + Inches(0.15), cards_y + Inches(0.05), card_w - Inches(0.3), Inches(0.25))
set_text(tf, "MID-LEVEL SECONDMENT", size=9, bold=True, color=GOLD)
tf = add_textbox(s, x + Inches(0.15), cards_y + Inches(0.27), card_w - Inches(0.3), Inches(0.3))
set_text(tf, "Finance Manager at N2M / month", size=12, bold=True, color=WHITE)

ry = cards_y + Inches(0.7)
waterfall_row(s, x + Inches(0.15), ry, card_w - Inches(0.3), "Hospital invoice (ex VAT)", "N2,000,000", "gross")
ry += Inches(0.38)
light_row(s, x + Inches(0.15), ry, card_w - Inches(0.3), "+ VAT 7.5%", "N150,000")
ry += Inches(0.32)
light_row(s, x + Inches(0.15), ry, card_w - Inches(0.3), "- WHT 10% (FIRS)", "N200,000")
ry += Inches(0.4)
waterfall_row(s, x + Inches(0.15), ry, card_w - Inches(0.3), "CFA platform fee (20%)", "N400,000", "cfa")
ry += Inches(0.38)
light_row(s, x + Inches(0.15), ry, card_w - Inches(0.3), "Employer pension + NSITF + ITF", "N190,000")
ry += Inches(0.4)
waterfall_row(s, x + Inches(0.15), ry, card_w - Inches(0.3), "Professional gross salary", "N1,410,000", "pro")
ry += Inches(0.38)
light_row(s, x + Inches(0.15), ry, card_w - Inches(0.3), "- PAYE + pension + NHIS", "N390,000")
ry += Inches(0.4)
waterfall_row(s, x + Inches(0.15), ry, card_w - Inches(0.3), "Professional take-home", "N1,020,000", "net")
ry += Inches(0.5)
tf = add_textbox(s, x + Inches(0.15), ry, card_w - Inches(0.3), Inches(0.4))
set_text(tf, "CFA is employer of record. Full statutory compliance handled.", size=8, color=GREY)

# ---------- CARD 3: Locum Shift (CORRECTED to N35K) ----------
x = start_x + (card_w + gap) * 2
add_rect(s, x, cards_y, card_w, card_h, fill=WHITE, line=SOFT)
add_rect(s, x, cards_y, card_w, Inches(0.55), fill=GOLD)
tf = add_textbox(s, x + Inches(0.15), cards_y + Inches(0.05), card_w - Inches(0.3), Inches(0.25))
set_text(tf, "LOCUM CLINICAL COVER", size=9, bold=True, color=NAVY)
tf = add_textbox(s, x + Inches(0.15), cards_y + Inches(0.27), card_w - Inches(0.3), Inches(0.3))
set_text(tf, "Medical Officer shift at N35K", size=12, bold=True, color=NAVY)

ry = cards_y + Inches(0.7)
waterfall_row(s, x + Inches(0.15), ry, card_w - Inches(0.3), "Hospital pays per shift", "N35,000", "gross")
ry += Inches(0.38)
light_row(s, x + Inches(0.15), ry, card_w - Inches(0.3), "- WHT 5% (individual, FIRS)", "N1,750")
ry += Inches(0.4)
waterfall_row(s, x + Inches(0.15), ry, card_w - Inches(0.3), "CFA platform margin (12%)", "N4,200", "cfa")
ry += Inches(0.42)
waterfall_row(s, x + Inches(0.15), ry, card_w - Inches(0.3), "Professional payout (instant)", "N29,050", "pro")
ry += Inches(0.38)
light_row(s, x + Inches(0.15), ry, card_w - Inches(0.3), "Indemnity (Turaco) bundled", "Included")
ry += Inches(0.32)
light_row(s, x + Inches(0.15), ry, card_w - Inches(0.3), "Professional own annual PIT", "Self-managed")
ry += Inches(0.4)
waterfall_row(s, x + Inches(0.15), ry, card_w - Inches(0.3), "Net cash to professional", "N29,050", "net")
ry += Inches(0.5)
tf = add_textbox(s, x + Inches(0.15), ry, card_w - Inches(0.3), Inches(0.4))
set_text(tf, "Per-shift contractor. Paid via Paystack within 24 hours of shift.", size=8, color=GREY)

# Bottom strip
strip_y = Inches(7.0)
add_rect(s, 0, strip_y, SLIDE_W, Inches(0.2), fill=NAVY)
tf = add_textbox(s, Inches(0.6), strip_y + Inches(0.0), Inches(12.1), Inches(0.2))
set_text(tf, "Statutory taxes (VAT, WHT, PAYE, pension) flow to government and pension funds. CFA fee covers sourcing, vetting, payroll, indemnity, and ongoing support.",
         size=9, color=WHITE, align=PP_ALIGN.CENTER)

add_rect(s, 0, Inches(7.2), SLIDE_W, Inches(0.3), fill=LIGHT)
tf = add_textbox(s, Inches(0.6), Inches(7.22), Inches(8), Inches(0.25))
set_text(tf, "consultforafrica.com  |  Maarova  |  CadreHealth  |  Indicative figures, subject to engagement specifics",
         size=8, color=GREY)

output_path = "/Users/debo/consult-for-africa/docs/cfa-economics-transparency-slide-v2.pptx"
prs.save(output_path)
print(f"Saved: {output_path}")
