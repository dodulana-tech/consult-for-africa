"""Build standalone slide explaining why CFA earns its platform fee.

Pairs with the economics transparency slide. Where that slide shows
WHERE the money goes, this slide shows WHY CFA is entitled to its cut.

Structure:
- 4 value pillars (what CFA layers on top)
- Comparison strip vs. doing it yourself or using a traditional agency
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

NAVY = RGBColor(0x0B, 0x3C, 0x5D)
GOLD = RGBColor(0xD4, 0xAF, 0x37)
TEAL = RGBColor(0x1F, 0x7A, 0x8C)
LIGHT = RGBColor(0xF4, 0xF7, 0xFA)
DARK = RGBColor(0x11, 0x18, 0x27)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GREY = RGBColor(0x6B, 0x72, 0x80)
SOFT = RGBColor(0xE5, 0xE7, 0xEB)
DEEP_TEAL = RGBColor(0x14, 0x55, 0x63)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def set_text(frame, text, size=18, bold=False, color=DARK, align=PP_ALIGN.LEFT, font="Calibri"):
    frame.clear()
    p = frame.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font


def add_paragraph(frame, text, size=14, bold=False, color=DARK, align=PP_ALIGN.LEFT, space_before=Pt(6), font="Calibri"):
    p = frame.add_paragraph()
    p.alignment = align
    p.space_before = space_before
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font


def add_rect(slide, x, y, w, h, fill=NAVY, line=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
    shape.shadow.inherit = False
    return shape


def add_textbox(slide, x, y, w, h):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.05)
    tf.margin_bottom = Inches(0.05)
    return tf


prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
blank = prs.slide_layouts[6]
s = prs.slides.add_slide(blank)

# Header
add_rect(s, 0, 0, SLIDE_W, Inches(0.6), fill=NAVY)
tf = add_textbox(s, Inches(0.6), Inches(0.1), Inches(8), Inches(0.4))
set_text(tf, "CONSULT FOR AFRICA", size=11, bold=True, color=WHITE)
tf = add_textbox(s, Inches(8.733), Inches(0.1), Inches(4.4), Inches(0.4))
set_text(tf, "INSERT SLIDE", size=11, bold=True, color=GOLD, align=PP_ALIGN.RIGHT)

# Title
tf = add_textbox(s, Inches(0.6), Inches(0.8), Inches(12.1), Inches(0.4))
set_text(tf, "COMMERCIALS  /  WHY THE CFA FEE", size=11, bold=True, color=GOLD)
tf = add_textbox(s, Inches(0.6), Inches(1.15), Inches(12.1), Inches(0.7))
set_text(tf, "What our fee actually pays for", size=28, bold=True, color=NAVY)
add_rect(s, Inches(0.6), Inches(1.85), Inches(0.8), Inches(0.05), fill=GOLD)

# Subtitle
tf = add_textbox(s, Inches(0.6), Inches(1.95), Inches(12.1), Inches(0.4))
set_text(tf, "Four layers of value the hospital does not build, hire, or carry alone.",
         size=14, color=DARK)

# ============= 4 VALUE PILLARS =============
pillars = [
    {
        "color": NAVY,
        "tag": "01",
        "title": "Vetted Network",
        "subtitle": "We have already done the sourcing",
        "items": [
            "4,300+ healthcare professionals, 16 cadres",
            "Senior CFA consultant bench across functions",
            "ICF-credentialed coaching panel",
            "License verified, reference checked, ID confirmed",
            "Skip the 6 to 12 month recruitment cycle",
        ],
    },
    {
        "color": TEAL,
        "tag": "02",
        "title": "Statutory Machinery",
        "subtitle": "Payroll, tax, pension, indemnity",
        "items": [
            "PAYE, pension, NHIS, NSITF, ITF, VAT, WHT",
            "FIRS remittance and audit-ready records",
            "Indemnity cover bundled (Turaco)",
            "NDPR-compliant data handling",
            "Contracts, IP, and exit management",
        ],
    },
    {
        "color": GOLD,
        "tag": "03",
        "title": "Risk Backstop",
        "subtitle": "We carry the wrong-hire risk",
        "items": [
            "Replacement guarantee on every engagement",
            "Performance accountability we own, not pass through",
            "Escalation line to CFA partners",
            "Substitute cover for cancelled locum shifts",
            "If it goes wrong, it is our problem to fix",
        ],
    },
    {
        "color": DEEP_TEAL,
        "tag": "04",
        "title": "Institutional Capability",
        "subtitle": "You hire the firm, not one person",
        "items": [
            "Playbooks and frameworks the consultant brings",
            "Cross-portfolio learning from other engagements",
            "Coaching and development during placement",
            "Advisory access from CFA partners",
            "Knowledge stays even when the person changes",
        ],
    },
]

cards_y = Inches(2.55)
card_w = Inches(2.95)
card_h = Inches(3.0)
gap = Inches(0.15)
start_x = Inches(0.6)

for i, p in enumerate(pillars):
    x = start_x + (card_w + gap) * i
    add_rect(s, x, cards_y, card_w, card_h, fill=WHITE, line=SOFT)
    # Top color band
    add_rect(s, x, cards_y, card_w, Inches(0.7), fill=p["color"])
    # Tag number (large, gold)
    tf = add_textbox(s, x + Inches(0.15), cards_y + Inches(0.05), Inches(0.6), Inches(0.4))
    set_text(tf, p["tag"], size=14, bold=True, color=GOLD)
    # Title
    tf = add_textbox(s, x + Inches(0.6), cards_y + Inches(0.05), card_w - Inches(0.7), Inches(0.4))
    set_text(tf, p["title"], size=15, bold=True, color=WHITE)
    # Subtitle
    tf = add_textbox(s, x + Inches(0.15), cards_y + Inches(0.42), card_w - Inches(0.3), Inches(0.25))
    set_text(tf, p["subtitle"], size=10, color=GOLD)
    # Items
    item_tf = add_textbox(s, x + Inches(0.15), cards_y + Inches(0.85), card_w - Inches(0.3), Inches(2.05))
    set_text(item_tf, "• " + p["items"][0], size=10, color=DARK)
    for item in p["items"][1:]:
        add_paragraph(item_tf, "• " + item, size=10, color=DARK, space_before=Pt(4))

# ============= COMPARISON STRIP =============
comp_y = Inches(5.7)
add_rect(s, Inches(0.6), comp_y, Inches(12.1), Inches(0.42), fill=NAVY)
tf = add_textbox(s, Inches(0.85), comp_y + Inches(0.07), Inches(11.6), Inches(0.3))
set_text(tf, "VS. THE ALTERNATIVES", size=11, bold=True, color=GOLD)

# 4 column comparison: dimension | DIY | Agency | CFA
table_y = comp_y + Inches(0.42)
col_widths = [Inches(3.0), Inches(3.0), Inches(3.05), Inches(3.05)]
col_x = [Inches(0.6)]
for w in col_widths[:-1]:
    col_x.append(col_x[-1] + w)

# Header row
hdrs = ["WHAT YOU NEED", "HIRE DIRECT", "TRADITIONAL AGENCY", "WITH CFA"]
hdr_colors = [LIGHT, LIGHT, LIGHT, GOLD]
hdr_text = [NAVY, GREY, GREY, NAVY]
for i, h in enumerate(hdrs):
    add_rect(s, col_x[i], table_y, col_widths[i], Inches(0.32), fill=hdr_colors[i])
    tf = add_textbox(s, col_x[i] + Inches(0.12), table_y + Inches(0.05), col_widths[i] - Inches(0.2), Inches(0.25))
    set_text(tf, h, size=10, bold=True, color=hdr_text[i])

rows = [
    ("Time to fill", "6 to 12 months", "2 to 3 months", "2 to 4 weeks"),
    ("Compliance burden", "All yours", "All yours", "All ours"),
    ("Replacement risk", "All yours", "Yours after 90 days", "Ours, ongoing"),
    ("Senior capability part-time", "Not possible", "Rare", "Standard offering"),
]
ry = table_y + Inches(0.32)
row_h = Inches(0.27)
for j, row in enumerate(rows):
    row_y = ry + row_h * j
    bg = WHITE if j % 2 == 0 else LIGHT
    for i, cell in enumerate(row):
        # CFA column always highlighted
        if i == 3:
            cell_bg = GOLD
            cell_color = NAVY
            bold = True
        else:
            cell_bg = bg
            cell_color = NAVY if i == 0 else DARK
            bold = (i == 0)
        add_rect(s, col_x[i], row_y, col_widths[i], row_h, fill=cell_bg)
        tf = add_textbox(s, col_x[i] + Inches(0.12), row_y + Inches(0.04), col_widths[i] - Inches(0.2), Inches(0.22))
        set_text(tf, cell, size=10, bold=bold, color=cell_color)

# Footer
add_rect(s, 0, Inches(7.2), SLIDE_W, Inches(0.3), fill=LIGHT)
tf = add_textbox(s, Inches(0.6), Inches(7.22), Inches(8), Inches(0.25))
set_text(tf, "consultforafrica.com  |  Maarova  |  CadreHealth", size=9, color=GREY)

output_path = "/Users/debo/consult-for-africa/docs/cfa-value-add-slide.pptx"
prs.save(output_path)
print(f"Standalone slide saved: {output_path}")
